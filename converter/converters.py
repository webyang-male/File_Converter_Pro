"""
Advanced Converter Engine — File Converter Pro
converter/converters.py  v3.0 — PEAK

HTML → PDF  : Playwright (headless Chrome) › pdfkit › fitz › reportlab+images
PDF → HTML  : fitz dict-mode (flow layout, base64 images, adaptive sizing)
XLSX → PDF  : smart portrait/landscape, auto col-widths, multi-sheet
PPTX → PDF  : LibreOffice › python-pptx+reportlab (images embedded)
TXT/RTF     : reportlab (styled) › pypandoc › striprtf
EPUB → PDF  : pypandoc › spine-order native (images embedded)
Images      : Pillow max-quality, EXIF preserved, HEIC via pillow-heif
Audio/Video : ffmpeg binary auto-located, quality presets per format

Author: Hyacinthe
Version: 1.0
"""

from __future__ import annotations

import base64
import csv
import io
import json
import os
import re
import subprocess
import tempfile
import time
import zipfile
from pathlib import Path

#  Result dataclass

class ConversionResult:
    __slots__ = ("success", "source", "target", "elapsed", "error", "file_size")

    def __init__(self, success, source, target, elapsed=0.0, error="", file_size=0):
        self.success   = success
        self.source    = source
        self.target    = target
        self.elapsed   = elapsed
        self.error     = error
        self.file_size = file_size

    def __repr__(self):
        s = "OK" if self.success else f"ERR({self.error})"
        return f"<ConversionResult {s} {self.source!r}→{self.target!r}>"

#  Helpers

def _timed(fn):
    t0 = time.perf_counter()
    fn()
    return time.perf_counter() - t0

def _build_dst(src, dst_dir, new_ext):
    return str(Path(dst_dir) / f"{Path(src).stem}.{new_ext.lstrip('.')}")

def _img_to_b64(data, mime="image/png"):
    return f"data:{mime};base64,{base64.b64encode(data).decode()}"

def _mime_for_ext(ext):
    return {
        "png": "image/png", "jpg": "image/jpeg", "jpeg": "image/jpeg",
        "gif": "image/gif", "webp": "image/webp", "svg": "image/svg+xml",
        "bmp": "image/bmp", "tiff": "image/tiff",
    }.get(ext.lower(), "image/png")

def _safe_html(text):
    return text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

def _read_file_b64(path):
    with open(path, "rb") as f:
        return base64.b64encode(f.read()).decode()

#  Category map

CATEGORY_MAP = {
    "txt_to_pdf":  "document", "rtf_to_pdf":   "document",
    "txt_to_docx": "document", "rtf_to_docx":  "document",
    "csv_to_json": "document", "json_to_csv":  "document",
    "xlsx_to_pdf": "document", "xlsx_to_json": "document",
    "xlsx_to_csv": "document", "pptx_to_pdf":  "document",
    "html_to_pdf": "document", "pdf_to_html":  "document",
    "epub_to_pdf": "document",
    "jpeg_to_png": "image",    "png_to_jpg":   "image",
    "jpg_to_png":  "image",    "webp_to_png":  "image",
    "bmp_to_png":  "image",    "tiff_to_png":  "image",
    "heic_to_png": "image",    "gif_to_png":   "image",
    "image_to_ico": "image",
    "wav_to_mp3":  "audio",    "mp3_to_wav":   "audio",
    "acc_to_mp3":  "audio",    "mp3_to_acc":   "audio",
    "flac_to_mp3": "audio",    "ogg_to_mp3":   "audio",
    "avi_to_mp4":  "video",    "webm_to_mp4":  "video",
    "mkv_to_mp4":  "video",    "mov_to_mp4":   "video",
    "mp4_to_mp3":  "audio",    "avi_to_mp3":   "audio",
    "webm_to_mp3": "audio",    "mkv_to_mp3":   "audio",
}

#  Engine

class AdvancedConverterEngine:

    _DISPATCH = {
        "txt_to_pdf":   ("_txt_to_pdf",     "pdf"),
        "rtf_to_pdf":   ("_rtf_to_pdf",     "pdf"),
        "txt_to_docx":  ("_txt_to_docx",    "docx"),
        "rtf_to_docx":  ("_rtf_to_docx",    "docx"),
        "csv_to_json":  ("_csv_to_json",    "json"),
        "json_to_csv":  ("_json_to_csv",    "csv"),
        "xlsx_to_pdf":  ("_xlsx_to_pdf",    "pdf"),
        "xlsx_to_json": ("_xlsx_to_json",   "json"),
        "xlsx_to_csv":  ("_xlsx_to_csv",    "csv"),
        "pptx_to_pdf":  ("_pptx_to_pdf",    "pdf"),
        "html_to_pdf":  ("_html_to_pdf",    "pdf"),
        "pdf_to_html":  ("_pdf_to_html",    "html"),
        "epub_to_pdf":  ("_epub_to_pdf",    "pdf"),
        "jpeg_to_png":  ("_image_convert",  "png"),
        "png_to_jpg":   ("_image_convert",  "jpg"),
        "jpg_to_png":   ("_image_convert",  "png"),
        "webp_to_png":  ("_image_convert",  "png"),
        "bmp_to_png":   ("_image_convert",  "png"),
        "tiff_to_png":  ("_image_convert",  "png"),
        "heic_to_png":  ("_heic_to_png",    "png"),
        "gif_to_png":   ("_image_convert",  "png"),
        "image_to_ico": ("_image_to_ico",   "ico"),
        "wav_to_mp3":   ("_ffmpeg_convert", "mp3"),
        "mp3_to_wav":   ("_ffmpeg_convert", "wav"),
        "acc_to_mp3":   ("_ffmpeg_convert", "mp3"),
        "mp3_to_acc":   ("_ffmpeg_convert", "aac"),
        "flac_to_mp3":  ("_ffmpeg_convert", "mp3"),
        "ogg_to_mp3":   ("_ffmpeg_convert", "mp3"),
        "avi_to_mp4":   ("_ffmpeg_convert", "mp4"),
        "webm_to_mp4":  ("_ffmpeg_convert", "mp4"),
        "mkv_to_mp4":   ("_ffmpeg_convert", "mp4"),
        "mov_to_mp4":   ("_ffmpeg_convert", "mp4"),
        "mp4_to_mp3":   ("_ffmpeg_convert", "mp3"),
        "avi_to_mp3":   ("_ffmpeg_convert", "mp3"),
        "webm_to_mp3":  ("_ffmpeg_convert", "mp3"),
        "mkv_to_mp3":   ("_ffmpeg_convert", "mp3"),
    }

    # COM / Office automation helper
    @staticmethod
    def _office_to_pdf_com(src: str, dst: str, app_name: str) -> bool:
        """
        Convert src → dst PDF via Microsoft Office COM automation.
        app_name: "PowerPoint.Application" | "Excel.Application" | "Word.Application"
        Returns True on success, False if Office/comtypes not available.
        Works only on Windows with Office installed.
        """
        try:
            import comtypes.client
            import comtypes
            src_abs = str(Path(src).resolve())
            dst_abs = str(Path(dst).resolve())

            # Excel constants
            XL_PORTRAIT  = 1
            XL_LANDSCAPE = 2
            XL_PDF       = 0    # xlTypePDF

            # PowerPoint constants
            PP_PDF           = 32  # ppSaveAsPDF
            PP_WIN_MINIMIZED = 2

            # Word constants
            WD_PDF = 17   # wdFormatPDF

            if "Excel" in app_name:
                app = comtypes.client.CreateObject(app_name)
                try:
                    app.Visible = False
                except Exception:
                    pass
                app.DisplayAlerts = False
                try:
                    wb = app.Workbooks.Open(src_abs)
                    try:
                        for sheet in wb.Worksheets:
                            try:
                                used   = sheet.UsedRange
                                n_cols = used.Columns.Count
                                n_rows = used.Rows.Count
                                if n_cols > 6 or (n_cols > 0 and n_cols > n_rows * 1.5):
                                    sheet.PageSetup.Orientation = XL_LANDSCAPE
                                else:
                                    sheet.PageSetup.Orientation = XL_PORTRAIT
                                sheet.PageSetup.Zoom           = False
                                sheet.PageSetup.FitToPagesWide = 1
                                sheet.PageSetup.FitToPagesTall = 0
                            except Exception:
                                pass
                        wb.ExportAsFixedFormat(XL_PDF, dst_abs)
                    finally:
                        wb.Close(False)
                finally:
                    app.Quit()

            elif "PowerPoint" in app_name:
                import shutil

                tmp_dst = str(Path(tempfile.gettempdir()) / "pptx_com_out.pdf")
                src_esc = src_abs.replace("\\", "\\\\").replace("'", "\\'")
                dst_esc = tmp_dst.replace("\\", "\\\\")

                # MsoTriState values: msoTrue=-1, msoFalse=0
                # Strategy A: SaveAs format 32 (ppSaveAsPDF)
                ps_a = f"""
$app = New-Object -ComObject PowerPoint.Application
$app.Visible = -1
try {{
    $prs = $app.Presentations.Open('{src_esc}', 0, 0, -1)
    $prs.SaveAs('{dst_esc}', 32)
    $prs.Close()
}} finally {{ $app.Quit() }}
"""
                r = subprocess.run(
                    ["powershell", "-NoProfile", "-NonInteractive",
                     "-ExecutionPolicy", "Bypass", "-Command", ps_a],
                    capture_output=True, timeout=120
                )
                if Path(tmp_dst).exists():
                    shutil.move(tmp_dst, dst_abs)
                    return Path(dst_abs).exists()

                # Strategy B: PrintOut to file (like "Print to PDF")
                ps_b = f"""
$app = New-Object -ComObject PowerPoint.Application
$app.Visible = -1
try {{
    $prs = $app.Presentations.Open('{src_esc}', 0, 0, -1)
    $prs.PrintOut(1, $prs.Slides.Count, '{dst_esc}', 0, 2)
    $prs.Close()
}} finally {{ $app.Quit() }}
"""
                r2 = subprocess.run(
                    ["powershell", "-NoProfile", "-NonInteractive",
                     "-ExecutionPolicy", "Bypass", "-Command", ps_b],
                    capture_output=True, timeout=120
                )
                if Path(tmp_dst).exists():
                    shutil.move(tmp_dst, dst_abs)
                    return Path(dst_abs).exists()

                print(f"[COM-PS] stderr: {r2.stderr.decode('utf-8','replace')[:400]}")

            elif "Word" in app_name:
                app = comtypes.client.CreateObject(app_name)
                try:
                    app.Visible = False
                except Exception:
                    pass
                try:
                    doc = app.Documents.Open(src_abs)
                    try:
                        doc.SaveAs2(dst_abs, WD_PDF)
                    finally:
                        doc.Close(False)
                finally:
                    app.Quit()

            return Path(dst).exists()

        except Exception as _com_exc:
            print(f"[COM] {app_name} failed: {_com_exc}")
            return False

    # public
    def convert(self, conversion_type, src, dst_dir):
        if conversion_type not in self._DISPATCH:
            return ConversionResult(False, src, "", error=f"Unknown type: {conversion_type}")
        method_name, ext = self._DISPATCH[conversion_type]
        dst = _build_dst(src, dst_dir, ext)
        file_size = os.path.getsize(src) if os.path.exists(src) else 0
        try:
            method = getattr(self, method_name)
            if method_name in ("_image_convert", "_ffmpeg_convert", "_heic_to_png"):
                elapsed = _timed(lambda: method(src, dst, conversion_type))
            else:
                elapsed = _timed(lambda: method(src, dst))
            return ConversionResult(True, src, dst, elapsed=elapsed, file_size=file_size)
        except Exception as exc:
            return ConversionResult(False, src, dst, error=str(exc), file_size=file_size)

    def convert_batch(self, conversion_type, sources, dst_dir, progress_cb=None):
        results = []
        for i, src in enumerate(sources, 1):
            result = self.convert(conversion_type, src, dst_dir)
            results.append(result)
            if progress_cb:
                progress_cb(i, len(sources), src)
        return results

    #  TXT / RTF

    def _txt_to_pdf(self, src, dst):
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.units import mm
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer

        with open(src, "r", encoding="utf-8", errors="replace") as f:
            lines = f.readlines()

        styles  = getSampleStyleSheet()
        normal  = ParagraphStyle("N", parent=styles["Normal"],
                                 fontSize=10, leading=14, spaceAfter=2)
        heading = ParagraphStyle("H", parent=styles["Heading2"],
                                 fontSize=12, leading=16, spaceBefore=8, spaceAfter=4)
        story = []
        for line in lines:
            line = line.rstrip()
            safe = _safe_html(line)
            if not safe:
                story.append(Spacer(1, 6))
            elif safe.isupper() and 0 < len(safe) < 80:
                story.append(Paragraph(safe, heading))
            else:
                story.append(Paragraph(safe or "&nbsp;", normal))

        SimpleDocTemplate(dst, pagesize=A4,
                          leftMargin=20*mm, rightMargin=20*mm,
                          topMargin=20*mm, bottomMargin=20*mm).build(story)

    def _rtf_to_pdf(self, src, dst):
        """
        RTF → PDF with images and tables preserved.
        1. Word COM (Windows)  — pixel-perfect, full fidelity
        2. pypandoc            — preserves images+tables if pandoc installed
        3. Native RTF parser   — extracts embedded images (pict/objdata),
                                 tables (trowd), text with formatting
        """
        src_abs = str(Path(src).resolve())

        # 1. Word COM — best quality, handles everything
        if self._office_to_pdf_com(src_abs, dst, "Word.Application"):
            return

        # 2. pypandoc
        try:
            import pypandoc
            pypandoc.convert_file(src, "pdf", outputfile=dst)
            return
        except Exception:
            pass

        # 3. Native — parse RTF, extract images + tables + text → PDF
        self._rtf_to_pdf_native(src, dst)

    #  RTF Native Parser — shared engine used by both _rtf_to_pdf_native
    #  and the native fallback in _rtf_to_docx.
    #
    #  Preserves: bold · italic · underline · font sizes · colors
    #             tables (trowd/cell/row) · embedded PNG/JPEG images (pict)
    #  No dependency on striprtf, pandoc, or Office.

    @staticmethod
    def _rtf_read_raw(src: str) -> str:
        """Read RTF bytes and decode robustly (utf-8 → cp1252 → latin-1)."""
        with open(src, "rb") as fh:
            raw_bytes = fh.read()
        for enc in ("utf-8", "cp1252", "latin-1"):
            try:
                return raw_bytes.decode(enc)
            except Exception:
                pass
        return raw_bytes.decode("utf-8", errors="replace")

    @staticmethod
    def _rtf_parse_colortbl(raw: str) -> list:
        """
        Extract \\colortbl color entries.
        Returns list of (r,g,b) tuples; index 0 = default (black).
        """
        result = [(0, 0, 0)]
        m = re.search(r"\\colortbl\s*;(.*?)}", raw, re.S)
        if not m:
            return result
        for entry in m.group(1).split(";"):
            if not entry.strip():
                continue
            def _ci(tag):
                hit = re.search(rf"\\{tag}(\d+)", entry)
                return int(hit.group(1)) if hit else 0
            result.append((_ci("red"), _ci("green"), _ci("blue")))
        return result

    @staticmethod
    def _rtf_extract_images(raw: str) -> list:
        """
        Extract embedded PNG/JPEG images from \\pict blocks (hex-encoded).
        Returns list of (image_bytes, 'png'|'jpg').
        WMF/EMF bitmaps are skipped (not portably renderable).
        """
        images = []
        for m in re.finditer(r"\{\\pict((?:[^{}]|\{[^{}]*\})*)\}", raw, re.S):
            block = m.group(1)
            if re.search(r"\\pngblip\b", block):
                fmt = "png"
            elif re.search(r"\\jpegblip\b", block):
                fmt = "jpg"
            else:
                continue  # wmf/emf — skip
            hex_data = re.sub(r"\\[a-zA-Z]+[-0-9]*\s?", "", block)
            hex_data = re.sub(r"[^0-9a-fA-F]", "", hex_data)
            if len(hex_data) < 8:
                continue
            try:
                images.append((bytes.fromhex(hex_data), fmt))
            except Exception:
                pass
        return images

    @staticmethod
    def _rtf_tokenize(raw: str) -> list:
        """
        Tokenize RTF source into:
          ('open',)               — {
          ('close',)              — }
          ('ctrl', name, param)   — \\cmd, \\cmd-3, \\cmd42
          ('text', chars)         — literal characters
        Handles \\' hex escapes, \\\\, \\{, \\}, \\~ \\- and CRLF.
        """
        tokens = []
        i, n = 0, len(raw)
        ctrl_re = re.compile(r"\\([a-zA-Z]+)(-?\d+)?[ ]?")
        while i < n:
            c = raw[i]
            if c == "{":
                tokens.append(("open",)); i += 1
            elif c == "}":
                tokens.append(("close",)); i += 1
            elif c == "\\":
                if i + 1 >= n:
                    i += 1; continue
                nc = raw[i + 1]
                if nc in ("\r", "\n"):
                    tokens.append(("ctrl", "par", None)); i += 2
                elif nc == "\\":
                    tokens.append(("text", "\\")); i += 2
                elif nc == "{":
                    tokens.append(("text", "{")); i += 2
                elif nc == "}":
                    tokens.append(("text", "}")); i += 2
                elif nc == "~":
                    tokens.append(("text", "\u00a0")); i += 2
                elif nc == "-":
                    tokens.append(("text", "\u00ad")); i += 2
                elif nc == "*":
                    tokens.append(("ctrl", "*", None)); i += 2
                elif nc == "'":
                    hx = raw[i+2:i+4]
                    try:
                        ch = bytes.fromhex(hx).decode("cp1252", errors="replace")
                    except Exception:
                        ch = "?"
                    tokens.append(("text", ch)); i += 4
                else:
                    m2 = ctrl_re.match(raw, i)
                    if m2:
                        name = m2.group(1)
                        param = int(m2.group(2)) if m2.group(2) else None
                        tokens.append(("ctrl", name, param))
                        i = m2.end()
                    else:
                        i += 1
            elif c in ("\r", "\n"):
                i += 1
            else:
                j = i
                while j < n and raw[j] not in ("{", "}", "\\", "\r", "\n"):
                    j += 1
                tokens.append(("text", raw[i:j]))
                i = j
        return tokens

    @staticmethod
    def _rtf_parse_spans(tokens: list, color_table: list) -> list:
        """
        Walk the token stream with a formatting state-stack and emit span dicts:
          { text, bold, italic, underline, fontsize (half-pts),
            color (r,g,b)|None, par, in_table, cell_end, row_end }

        Groups flagged by {\\* ...} and structural groups (fonttbl, colortbl,
        stylesheet, info, pict, object, header, footer, footnote…) are skipped.
        """
        SKIP_GROUPS = {
            # Structural / metadata
            "pict", "object", "objdata", "objclass", "objname",
            "info", "fonttbl", "colortbl", "stylesheet",
            "header", "footer", "headerf", "footerf", "headerl", "headerr",
            "footerl", "footerr", "footnote", "annotation", "comment",
            # Fields — fldinst only, fldrslt holds visible text so keep it
            "fldinst",
            # Binary blobs / Office theme/chart data
            "themedata", "colorschememapping", "datastore", "datastoreprop",
            "wgrffmtfilter", "blipuid", "mmathPr",
            # List/revision/doc metadata
            "listtable", "listoverridetable", "rsidtbl",
            "generator", "pgdsctbl", "docvar", "xmlnstbl",
            # Drawing/shapes — too complex, skip entirely
            "shp", "shpinst", "shprslt", "sp", "sn", "sv",
            "dptxbxtext", "dptxbx",
        }

        class St:
            __slots__ = ("bold","italic","underline","fontsize","color_idx",
                         "in_table","skip_depth")
            def __init__(self):
                self.bold=False; self.italic=False; self.underline=False
                self.fontsize=24; self.color_idx=0
                self.in_table=False; self.skip_depth=0
            def copy(self):
                s = St()
                for a in self.__slots__: setattr(s, a, getattr(self, a))
                return s

        stack, state, spans = [], St(), []

        def _span(text="", par=False, cell_end=False, row_end=False):
            if state.skip_depth > 0:
                return
            col = (color_table[state.color_idx]
                   if 0 < state.color_idx < len(color_table) else None)
            spans.append(dict(
                text=text, bold=state.bold, italic=state.italic,
                underline=state.underline, fontsize=state.fontsize,
                color=col, par=par, in_table=state.in_table,
                cell_end=cell_end, row_end=row_end,
            ))

        for i, tok in enumerate(tokens):
            kind = tok[0]
            if kind == "open":
                stack.append(state.copy())
                # Lookahead: if next token is \* this is a destination group → skip it
                if i + 1 < len(tokens) and tokens[i + 1] == ("ctrl", "*", None):
                    state.skip_depth += 1
            elif kind == "close":
                if state.skip_depth > 0:
                    state.skip_depth -= 1
                if stack:
                    state = stack.pop()
            elif kind == "ctrl":
                name, param = tok[1], tok[2]
                if name == "*":
                    continue  # handled via lookahead on open
                if name in SKIP_GROUPS:
                    state.skip_depth += 1; continue
                if state.skip_depth > 0:
                    continue
                if name == "b":
                    state.bold = (param != 0) if param is not None else True
                elif name == "i":
                    state.italic = (param != 0) if param is not None else True
                elif name == "ul":
                    state.underline = True
                elif name == "ulnone":
                    state.underline = False
                elif name == "fs":
                    state.fontsize = param if param else 24
                elif name == "cf":
                    state.color_idx = param if param is not None else 0
                elif name in ("par", "line"):
                    _span(par=True)
                elif name == "pard":
                    state.bold=False; state.italic=False
                    state.underline=False; state.fontsize=24; state.color_idx=0
                    state.in_table=False  # \pard exits table mode
                elif name == "intbl":
                    state.in_table = True
                elif name == "cell":
                    spans.append(dict(text="", bold=False, italic=False,
                                      underline=False, fontsize=24, color=None,
                                      par=False, in_table=True,
                                      cell_end=True, row_end=False))
                elif name == "row":
                    spans.append(dict(text="", bold=False, italic=False,
                                      underline=False, fontsize=24, color=None,
                                      par=False, in_table=True,
                                      cell_end=False, row_end=True))
                    state.in_table = False
                elif name == "tab":
                    _span(text="\t")
            elif kind == "text":
                if state.skip_depth == 0 and tok[1]:
                    _span(text=tok[1])

        return spans

    @staticmethod
    def _rtf_spans_to_paragraphs(spans: list) -> list:
        """
        Assemble spans into a document structure:
          { 'type': 'para',      'para': { runs, text_content } }
          { 'type': 'table_row', 'cells': [ [para, ...], ... ] }

        Adjacent runs with identical formatting are merged.
        """
        items = []
        cur_runs = []
        cur_cell_paras = []
        cur_row_cells = []

        def _flush(is_cell=False):
            merged = []
            for run in cur_runs:
                if (merged
                        and merged[-1]["bold"]      == run["bold"]
                        and merged[-1]["italic"]    == run["italic"]
                        and merged[-1]["underline"] == run["underline"]
                        and merged[-1]["fontsize"]  == run["fontsize"]
                        and merged[-1]["color"]     == run["color"]):
                    merged[-1]["text"] += run["text"]
                else:
                    merged.append(dict(run))
            cur_runs.clear()
            return {"runs": merged,
                    "text_content": "".join(r["text"] for r in merged).strip(),
                    "is_table_cell": is_cell}

        for sp in spans:
            if sp["cell_end"]:
                # Each \cell ends one column — flush runs into current cell paras,
                # then push that cell as a new column in cur_row_cells
                cur_cell_paras.append(_flush(is_cell=True))
                cur_row_cells.append(list(cur_cell_paras))
                cur_cell_paras.clear()
                continue
            if sp["row_end"]:
                # \row ends the table row — flush any trailing runs
                leftover = _flush(is_cell=True)
                if leftover["text_content"] or leftover["runs"]:
                    cur_cell_paras.append(leftover)
                    cur_row_cells.append(list(cur_cell_paras))
                    cur_cell_paras.clear()
                items.append({"type": "table_row", "cells": list(cur_row_cells)})
                cur_row_cells.clear()
                continue
            if sp["par"]:
                para = _flush()
                if sp["in_table"]:
                    cur_cell_paras.append(para)
                else:
                    items.append({"type": "para", "para": para})
                continue
            if sp["text"]:
                cur_runs.append({k: sp[k] for k in
                                 ("text","bold","italic","underline","fontsize","color")})

        if cur_runs:
            items.append({"type": "para", "para": _flush()})
        return items

    def _rtf_to_pdf_native(self, src, dst):
        """
        Full-fidelity native RTF → PDF (no Office, no pandoc needed).
        Preserves bold · italic · underline · font sizes · RGB colors ·
        tables · embedded PNG/JPEG images.
        """
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.units import mm
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib import colors as rl_colors
        from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer,
                                        Table, TableStyle, Image as RLImage)

        raw         = self._rtf_read_raw(src)
        color_table = self._rtf_parse_colortbl(raw)
        img_data    = self._rtf_extract_images(raw)
        tokens      = self._rtf_tokenize(raw)
        spans       = self._rtf_parse_spans(tokens, color_table)
        paras       = self._rtf_spans_to_paragraphs(spans)

        sty  = getSampleStyleSheet()
        base = ParagraphStyle("RTFBase", parent=sty["Normal"],
                              fontSize=11, leading=15, spaceAfter=3,
                              fontName="Helvetica")
        td_s = ParagraphStyle("RTFTd", parent=sty["Normal"],
                              fontSize=9, leading=12, fontName="Helvetica")

        def _runs_to_xml(runs, base_fs=11):
            parts = []
            for r in runs:
                txt = _safe_html(r["text"])
                if not txt:
                    continue
                fs = max(6, r["fontsize"] // 2)
                o, c = [], []
                if r["bold"]:      o.append("<b>");  c.insert(0,"</b>")
                if r["italic"]:    o.append("<i>");  c.insert(0,"</i>")
                if r["underline"]: o.append("<u>");  c.insert(0,"</u>")
                if r["color"]:
                    rv,gv,bv = r["color"]
                    o.append(f'<font color="#{rv:02x}{gv:02x}{bv:02x}">'); c.insert(0,"</font>")
                if fs != base_fs:
                    o.append(f'<font size="{fs}">'); c.insert(0,"</font>")
                parts.append("".join(o) + txt + "".join(c))
            return "".join(parts)

        story = []

        def _build_rl_table(rows_data):
            """Convert list of table_row items into a single ReportLab Table."""
            max_cols = max(len(r["cells"]) for r in rows_data)
            tbl_rows = []
            for row_item in rows_data:
                tbl_row = []
                for ci in range(max_cols):
                    cell_paras = row_item["cells"][ci] if ci < len(row_item["cells"]) else []
                    cell_content = []
                    for cp in cell_paras:
                        x = _runs_to_xml(cp["runs"], base_fs=9)
                        if x.strip():
                            cell_content.append(Paragraph(x, td_s))
                    tbl_row.append(cell_content or [Paragraph("", td_s)])
                tbl_rows.append(tbl_row)
            tbl = Table(tbl_rows, hAlign="LEFT", colWidths=[None]*max_cols)
            tbl.setStyle(TableStyle([
                ("GRID",          (0,0),(-1,-1), 0.4, rl_colors.HexColor("#bbbbbb")),
                ("BACKGROUND",    (0,0),(-1, 0), rl_colors.HexColor("#f5f5f5")),
                ("VALIGN",        (0,0),(-1,-1), "TOP"),
                ("LEFTPADDING",   (0,0),(-1,-1), 4),
                ("RIGHTPADDING",  (0,0),(-1,-1), 4),
                ("TOPPADDING",    (0,0),(-1,-1), 3),
                ("BOTTOMPADDING", (0,0),(-1,-1), 3),
            ]))
            return tbl

        idx = 0
        while idx < len(paras):
            item = paras[idx]

            if item["type"] == "table_row":
                # Collect all consecutive table_row items → one table
                rows_data = []
                while idx < len(paras) and paras[idx]["type"] == "table_row":
                    rows_data.append(paras[idx]); idx += 1
                story.append(_build_rl_table(rows_data))
                continue

            # Regular paragraph
            para = item["para"]
            xml  = _runs_to_xml(para["runs"])
            idx += 1
            if not xml.strip():
                story.append(Spacer(1, 4)); continue
            plain = para["text_content"]
            if (plain.isupper() and 0 < len(plain) < 80
                    and all(r["bold"] or r["fontsize"] >= 28
                            for r in para["runs"] if r["text"].strip())):
                hs = ParagraphStyle("RTFHead", parent=base, fontSize=13,
                                    leading=17, spaceBefore=6, spaceAfter=4,
                                    fontName="Helvetica-Bold")
                story.append(Paragraph(xml, hs))
            else:
                sizes = [r["fontsize"] for r in para["runs"] if r["text"].strip()]
                dom   = max(6, (max(sizes) // 2)) if sizes else 11
                ps    = ParagraphStyle("RTFPar", parent=base,
                                       fontSize=dom, leading=max(dom+3,13))
                story.append(Paragraph(xml, ps))

        # Append embedded images (RTF embeds them in-stream but we append at end)
        for img_bytes, fmt in img_data:
            try:
                from PIL import Image as PILImg
                buf   = io.BytesIO(img_bytes)
                pil   = PILImg.open(buf)
                w, h  = pil.size
                max_w = 120 * mm          # cap at 120 mm wide
                scale = min(1.0, max_w / max(w, 1))
                buf.seek(0)
                story.append(Spacer(1, 6))
                story.append(RLImage(buf, width=w*scale, height=h*scale))
                story.append(Spacer(1, 6))
            except Exception:
                pass

        if not story:
            story.append(Paragraph("(empty document)", base))

        SimpleDocTemplate(
            dst, pagesize=A4,
            leftMargin=20*mm, rightMargin=20*mm,
            topMargin=20*mm, bottomMargin=20*mm,
        ).build(story)

    def _txt_to_docx(self, src, dst):
        from docx import Document
        from docx.shared import Pt
        with open(src, "r", encoding="utf-8", errors="replace") as f:
            lines = f.readlines()
        doc = Document()
        for p in list(doc.paragraphs):
            p._element.getparent().remove(p._element)
        for line in lines:
            line = line.rstrip()
            if line.isupper() and 0 < len(line) < 80:
                doc.add_heading(line, level=2)
            else:
                p = doc.add_paragraph(line)
                p.style.font.size = Pt(11)
        doc.save(dst)

    def _rtf_to_docx(self, src, dst):
        """
        RTF → DOCX with formatting, tables and images preserved.
        1. Word COM (Windows + Office) — save directly as DOCX (wdFormatXMLDocument)
        2. pypandoc                    — full fidelity if pandoc installed
        3. Native parser               — bold/italic/underline/colors/fontsize/
                                         tables/embedded PNG+JPEG images
        """
        src_abs = str(Path(src).resolve())
        dst_abs = str(Path(dst).resolve())

        # 1. Word COM — save directly as DOCX (wdFormatXMLDocument = 12)
        try:
            import comtypes.client
            app = comtypes.client.CreateObject("Word.Application")
            try:
                app.Visible = False
            except Exception:
                pass
            try:
                doc = app.Documents.Open(src_abs)
                try:
                    doc.SaveAs2(dst_abs, 12)
                finally:
                    doc.Close(False)
            finally:
                app.Quit()
            if Path(dst).exists():
                return
        except Exception:
            pass

        # 2. pypandoc
        try:
            import pypandoc
            pypandoc.convert_file(src, "docx", outputfile=dst)
            if Path(dst).exists():
                return
        except Exception:
            pass

        # 3. Native — full formatting via shared RTF parser
        from docx import Document
        from docx.shared import Pt, RGBColor

        raw         = self._rtf_read_raw(src)
        color_table = self._rtf_parse_colortbl(raw)
        img_data    = self._rtf_extract_images(raw)
        tokens      = self._rtf_tokenize(raw)
        spans       = self._rtf_parse_spans(tokens, color_table)
        paras       = self._rtf_spans_to_paragraphs(spans)

        doc = Document()
        for p in list(doc.paragraphs):
            p._element.getparent().remove(p._element)

        def _add_runs(docx_para, runs):
            for run in runs:
                if not run["text"]:
                    continue
                r = docx_para.add_run(run["text"])
                r.bold      = run["bold"]
                r.italic    = run["italic"]
                r.underline = run["underline"]
                r.font.size = Pt(max(6, run["fontsize"] // 2))
                if run["color"]:
                    r.font.color.rgb = RGBColor(*run["color"])

        i = 0
        while i < len(paras):
            item = paras[i]

            if item["type"] == "table_row":
                # Collect all consecutive table rows into one DOCX table
                rows = []
                while i < len(paras) and paras[i]["type"] == "table_row":
                    rows.append(paras[i]["cells"]); i += 1
                max_cols = max(len(r) for r in rows)
                tbl = doc.add_table(rows=len(rows), cols=max_cols)
                tbl.style = "Table Grid"
                for ri, row_cells in enumerate(rows):
                    for ci, cell_paras in enumerate(row_cells):
                        if ci >= max_cols:
                            break
                        cell = tbl.cell(ri, ci)
                        for cp in list(cell.paragraphs):
                            cp._element.getparent().remove(cp._element)
                        for cp in cell_paras:
                            if not cp["text_content"] and not cp["runs"]:
                                cell.add_paragraph("")
                                continue
                            np_ = cell.add_paragraph()
                            np_.style.font.size = Pt(9)
                            _add_runs(np_, cp["runs"])
                continue

            # Regular paragraph
            para  = item["para"]
            plain = para["text_content"]
            if not plain:
                doc.add_paragraph(""); i += 1; continue

            if (plain.isupper() and 0 < len(plain) < 80
                    and all(r["bold"] or r["fontsize"] >= 28
                            for r in para["runs"] if r["text"].strip())):
                doc.add_heading(plain, level=2)
            else:
                p = doc.add_paragraph()
                _add_runs(p, para["runs"])
            i += 1

        # Append embedded images
        for img_bytes, _ in img_data:
            try:
                doc.add_picture(io.BytesIO(img_bytes))
            except Exception:
                pass

        doc.save(dst)

    #  CSV / JSON

    def _csv_to_json(self, src, dst):
        def _cast(val):
            """Try to restore original JSON types from CSV string."""
            if val == "":
                return None
            if val.lower() == "true":
                return True
            if val.lower() == "false":
                return False
            # Integer?
            try:
                i = int(val)
                # Avoid casting phone numbers / long digit strings as int
                if len(val) < 16:
                    return i
            except ValueError:
                pass
            # Float?
            try:
                return float(val)
            except ValueError:
                pass
            return val

        with open(src, newline="", encoding="utf-8-sig") as f:
            sample = f.read(4096)
            f.seek(0)
            try:
                dialect = csv.Sniffer().sniff(sample, delimiters=",;\t|")
            except Exception:
                dialect = csv.excel
            rows = [{k: _cast(v) for k, v in row.items()}
                    for row in csv.DictReader(f, dialect=dialect)]

        with open(dst, "w", encoding="utf-8") as f:
            json.dump(rows, f, indent=2, ensure_ascii=False)

    def _json_to_csv(self, src, dst):
        with open(src, encoding="utf-8") as f:
            data = json.load(f)
        if isinstance(data, dict):
            for v in data.values():
                if isinstance(v, list) and v:
                    data = v
                    break
            else:
                data = [data]
        if not data:
            raise ValueError("JSON is empty.")
        flat = []
        for item in data:
            if isinstance(item, dict):
                flat.append({k: (json.dumps(v) if isinstance(v, (dict, list)) else v)
                             for k, v in item.items()})
            else:
                flat.append({"value": item})
        keys = list(flat[0].keys())
        with open(dst, "w", newline="", encoding="utf-8-sig") as f:
            w = csv.DictWriter(f, fieldnames=keys, extrasaction="ignore")
            w.writeheader()
            w.writerows(flat)

    #  XLSX

    def _xlsx_to_pdf(self, src, dst):
        """
        XLSX/XLS → PDF.
        1. Microsoft Office COM (Windows + Excel) — perfect, smart orientation per sheet
        2. LibreOffice CLI      — if Excel not available
        3. reportlab native     — fallback, always works
        """
        src_abs = str(Path(src).resolve())

        # 1. COM / Microsoft Excel (handles XLS + XLSX, smart orientation built-in)
        if self._office_to_pdf_com(src_abs, dst, "Excel.Application"):
            return

        # 2. LibreOffice (if Excel not available)
        try:
            dst_dir = str(Path(dst).parent)
            lo_candidates = [
                "libreoffice", "soffice",
                r"C:\Program Files\LibreOffice\program\soffice.exe",
                r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
                "/Applications/LibreOffice.app/Contents/MacOS/soffice",
            ]
            import shutil as _shutil
            lo_bin = next((c for c in lo_candidates
                           if _shutil.which(c) or __import__("os").path.isfile(c)), None)
            if lo_bin:
                subprocess.run(
                    [lo_bin, "--headless", "--convert-to", "pdf",
                     "--outdir", dst_dir, src_abs],
                    check=True, capture_output=True, timeout=120)
                lo_out = Path(dst_dir) / (Path(src).stem + ".pdf")
                if lo_out.exists() and str(lo_out) != dst:
                    lo_out.rename(dst)
                if Path(dst).exists():
                    return
            else:
                print("[INFO] LibreOffice not found, falling back to reportlab.")
        except Exception as _lo_err:
            print(f"[WARNING] LibreOffice XLSX→PDF failed: {_lo_err}")

        # 3. reportlab native fallback
        import openpyxl
        from reportlab.lib.pagesizes import A4, landscape
        from reportlab.lib import colors
        from reportlab.lib.units import mm
        from reportlab.platypus import (SimpleDocTemplate, Table, TableStyle,
                                        Spacer, Paragraph, PageBreak)
        from reportlab.lib.styles import getSampleStyleSheet

        wb     = openpyxl.load_workbook(src, data_only=True)
        styles = getSampleStyleSheet()
        story  = []
        HEADER_BG  = colors.HexColor("#4f46e5")
        ROW_ALT    = colors.HexColor("#f3f4f6")
        GRID_COLOR = colors.HexColor("#d1d5db")

        # Smart orientation: analyse all sheets
        max_cols = 0
        max_rows = 0
        for sheet in wb.worksheets:
            raw_check = list(sheet.iter_rows())
            if raw_check:
                max_cols = max(max_cols, len(raw_check[0]))
                max_rows = max(max_rows, len(raw_check))
        use_landscape = (max_cols > 6) or (max_cols > 0 and max_cols > max_rows * 1.5)
        page_size = landscape(A4) if use_landscape else A4
        page_w    = page_size[0] - 40 * mm

        for si, sheet in enumerate(wb.worksheets):
            if si > 0:
                story.append(PageBreak())
            story.append(Paragraph(f"<b>{_safe_html(sheet.title)}</b>", styles["Heading2"]))
            story.append(Spacer(1, 6))
            raw = [[str(c.value if c.value is not None else "") for c in row]
                   for row in sheet.iter_rows()]
            if not raw:
                continue
            n = len(raw[0])
            col_max = [0] * n
            for row in raw:
                for ci, cell in enumerate(row):
                    col_max[ci] = max(col_max[ci], len(str(cell)))
            raw_pts = [min(col_max[ci] * 5.5 + 8, 80 * mm) for ci in range(n)]
            total   = sum(raw_pts) or 1
            col_pts = [p / total * page_w for p in raw_pts]
            t = Table(raw, colWidths=col_pts, repeatRows=1)
            t.setStyle(TableStyle([
                ("BACKGROUND",    (0, 0), (-1, 0),  HEADER_BG),
                ("TEXTCOLOR",     (0, 0), (-1, 0),  colors.white),
                ("FONTNAME",      (0, 0), (-1, 0),  "Helvetica-Bold"),
                ("FONTSIZE",      (0, 0), (-1, 0),  9),
                ("FONTSIZE",      (0, 1), (-1, -1), 8),
                ("ROWBACKGROUNDS",(0, 1), (-1, -1), [colors.white, ROW_ALT]),
                ("GRID",          (0, 0), (-1, -1), 0.25, GRID_COLOR),
                ("VALIGN",        (0, 0), (-1, -1), "MIDDLE"),
                ("LEFTPADDING",   (0, 0), (-1, -1), 4),
                ("RIGHTPADDING",  (0, 0), (-1, -1), 4),
                ("TOPPADDING",    (0, 0), (-1, -1), 3),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
                ("WORDWRAP",      (0, 0), (-1, -1), True),
            ]))
            story.append(t)

        SimpleDocTemplate(dst, pagesize=page_size,
                          leftMargin=20 * mm, rightMargin=20 * mm,
                          topMargin=15 * mm, bottomMargin=15 * mm).build(story)

    def _xlsx_to_json(self, src, dst):
        import openpyxl
        wb  = openpyxl.load_workbook(src, data_only=True)
        out = {}
        for sheet in wb.worksheets:
            rows = list(sheet.iter_rows(values_only=True))
            if not rows:
                out[sheet.title] = []
                continue
            headers = [str(h) if h is not None else f"col_{i}"
                       for i, h in enumerate(rows[0])]
            out[sheet.title] = [
                {headers[i]: (v.isoformat() if hasattr(v, "isoformat") else v)
                 for i, v in enumerate(row)}
                for row in rows[1:]
            ]
        with open(dst, "w", encoding="utf-8") as f:
            json.dump(out, f, indent=2, ensure_ascii=False, default=str)

    def _xlsx_to_csv(self, src, dst):
        import openpyxl
        wb    = openpyxl.load_workbook(src, data_only=True)
        sheet = wb.active
        with open(dst, "w", newline="", encoding="utf-8-sig") as f:
            w = csv.writer(f)
            for row in sheet.iter_rows(values_only=True):
                w.writerow([v.isoformat() if hasattr(v, "isoformat") else
                            (v if v is not None else "") for v in row])

    #  PPTX → PDF

    def _pptx_to_pdf(self, src, dst):
        """
        PPTX/PPT → PDF.
        1. Microsoft Office COM (Windows + Office installed) — pixel-perfect, handles PPT
        2. LibreOffice CLI      (if installed)
        3. Native engine        (python-pptx + reportlab + matplotlib)
           Note: python-pptx cannot open .ppt (old binary format) — if COM and
           LibreOffice both fail on a .ppt file, we raise a clear error.
        """
        src_abs = str(Path(src).resolve())

        # 1. COM / Microsoft PowerPoint
        if self._office_to_pdf_com(src_abs, dst, "PowerPoint.Application"):
            return

        # 2. LibreOffice
        try:
            dst_dir = str(Path(dst).parent)
            lo_candidates = [
                "libreoffice", "soffice",
                r"C:\Program Files\LibreOffice\program\soffice.exe",
                r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
                "/Applications/LibreOffice.app/Contents/MacOS/soffice",
            ]
            import shutil as _shutil
            lo_bin = next((c for c in lo_candidates
                           if _shutil.which(c) or __import__("os").path.isfile(c)), None)
            if lo_bin:
                subprocess.run(
                    [lo_bin, "--headless", "--convert-to", "pdf",
                     "--outdir", dst_dir, src_abs],
                    check=True, capture_output=True, timeout=120)
                lo_out = Path(dst_dir) / (Path(src).stem + ".pdf")
                if lo_out.exists() and str(lo_out) != dst:
                    lo_out.rename(dst)
                if Path(dst).exists():
                    return
            else:
                print("[INFO] LibreOffice not found, falling back to native engine.")
        except Exception as _lo_err:
            print(f"[WARNING] LibreOffice PPTX→PDF failed: {_lo_err}")

        # 3. Native (python-pptx — PPTX only, not PPT)
        if Path(src).suffix.lower() == ".ppt":
            raise RuntimeError(
                "Le format .ppt nécessite Microsoft Office ou LibreOffice.\n"
                "COM PowerPoint a été tenté mais a échoué (voir [COM] dans la console).\n"
                "Vérifiez que PowerPoint est bien fermé avant la conversion, "
                "ou convertissez d'abord le fichier en .pptx."
            )
        self._pptx_to_pdf_native(src, dst)

    def _pptx_to_pdf_native(self, src, dst):
        """
        Native PPTX → PDF using python-pptx + reportlab + matplotlib.
        Handles: text, bullets, images, tables, charts (bar/line/pie/donut).
        All python-pptx API calls are wrapped in try/except.
        """
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from reportlab.lib.pagesizes import landscape, A4
        from reportlab.lib.units import mm
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib import colors
        from reportlab.lib.enums import TA_CENTER, TA_RIGHT
        from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer,
                                        Image as RLImage, PageBreak,
                                        Table, TableStyle, HRFlowable)

        prs = Presentation(src)
        try:    sw = int(prs.slide_width)  or 9144000
        except: sw = 9144000
        try:    sh = int(prs.slide_height) or 6858000
        except: sh = 6858000

        ratio   = float(sh) / float(sw)
        pw      = landscape(A4)[0]
        ph      = pw * ratio
        margin  = 14 * mm
        inner_w = pw - 2 * margin

        S = getSampleStyleSheet()
        def sty(name, **kw):
            p = kw.pop("parent", S["Normal"])
            return ParagraphStyle(name, parent=p, **kw)

        ST = {
            "title"  : sty("PT", parent=S["Heading1"], fontSize=20, leading=24,
                           spaceAfter=8, textColor=colors.HexColor("#1a1a2e"),
                           alignment=TA_CENTER),
            "h2"     : sty("PH2", parent=S["Heading2"], fontSize=15, leading=19,
                           spaceAfter=5, textColor=colors.HexColor("#1e3a5f")),
            "body"   : sty("PB", fontSize=11, leading=16, spaceAfter=4),
            "bul0"   : sty("PB0", fontSize=11, leading=16, spaceAfter=3, leftIndent=12),
            "bul1"   : sty("PB1", fontSize=10, leading=15, spaceAfter=2, leftIndent=24,
                           textColor=colors.HexColor("#333")),
            "bul2"   : sty("PB2", fontSize=9,  leading=14, spaceAfter=2, leftIndent=36,
                           textColor=colors.HexColor("#555")),
            "slide_n": sty("PSN", fontSize=8, leading=10,
                           textColor=colors.HexColor("#aaa"), alignment=TA_RIGHT),
            "td"     : sty("PTd", fontSize=9, leading=13),
            "th"     : sty("PTh", fontSize=9, leading=13, fontName="Helvetica-Bold"),
            "caption": sty("PC",  fontSize=8, leading=11,
                           textColor=colors.HexColor("#777"), alignment=TA_CENTER),
        }

        BULLETS = set("•–-*◦▪▸")
        story   = []
        tmp_imgs = []

        def _safe(txt):
            return _safe_html(str(txt or ""))

        def _add_img_blob(blob, ext):
            try:
                tf = tempfile.NamedTemporaryFile(suffix=f".{ext or 'png'}", delete=False)
                tf.write(blob); tf.close()
                tmp_imgs.append(tf.name)
                rl = RLImage(tf.name)
                scale = min(inner_w / rl.imageWidth,
                            (ph - 2*margin)*0.6 / rl.imageHeight, 1.0)
                rl.drawWidth  = rl.imageWidth  * scale
                rl.drawHeight = rl.imageHeight * scale
                rl.hAlign = "CENTER"
                story.append(Spacer(1, 4))
                story.append(rl)
                story.append(Spacer(1, 4))
            except Exception:
                pass

        def _render_chart(shape):
            """Extract chart data and render as matplotlib image."""
            try:
                import matplotlib
                matplotlib.use("Agg")
                import matplotlib.pyplot as plt
                import matplotlib.patches as mpatches
                import numpy as np
                from pptx.enum.chart import XL_CHART_TYPE
                from lxml import etree

                chart = shape.chart
                ctype = chart.chart_type

                # Extract series data
                NS_C = "http://schemas.openxmlformats.org/drawingml/2006/chart"
                NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
                cxml = chart._element

                series_list = []
                for ser in cxml.findall(f".//{{{NS_C}}}ser"):
                    # name
                    name = ""
                    v_el = ser.find(f".//{{{NS_C}}}tx//{{{NS_C}}}v")
                    if v_el is not None: name = v_el.text or ""

                    # categories
                    cats = [el.find(f"{{{NS_C}}}v").text
                            for el in ser.findall(f".//{{{NS_C}}}cat//{{{NS_C}}}pt")
                            if el.find(f"{{{NS_C}}}v") is not None]

                    # values
                    vals = []
                    for pt in ser.findall(f".//{{{NS_C}}}val//{{{NS_C}}}pt"):
                        v = pt.find(f"{{{NS_C}}}v")
                        if v is not None:
                            try: vals.append(float(v.text))
                            except: vals.append(0.0)

                    # per-point colours from dPt
                    pt_colors = {}
                    for dpt in ser.findall(f"{{{NS_C}}}dPt"):
                        idx_el = dpt.find(f"{{{NS_C}}}idx")
                        clr_el = dpt.find(f".//{{{NS_A}}}srgbClr")
                        if idx_el is not None and clr_el is not None:
                            idx = int(idx_el.get("val", 0))
                            pt_colors[idx] = "#" + clr_el.get("val", "4472C4")

                    # series-level colour
                    ser_clr = None
                    clr_el  = ser.find(f".//{{{NS_A}}}srgbClr")
                    if clr_el is not None:
                        ser_clr = "#" + clr_el.get("val", "4472C4")

                    series_list.append({
                        "name": name, "cats": cats, "vals": vals,
                        "pt_colors": pt_colors, "ser_color": ser_clr,
                    })

                if not series_list:
                    return False

                # Draw with matplotlib
                fig_w = float(inner_w) / 72 / 1.333 * 1.5
                fig_h = fig_w * 0.6
                fig, ax = plt.subplots(figsize=(fig_w, fig_h))
                fig.patch.set_facecolor("white")
                ax.set_facecolor("#f8f9fa")

                # Determine chart family
                is_pie  = ctype in (
                    XL_CHART_TYPE.PIE, XL_CHART_TYPE.PIE_EXPLODED,
                    XL_CHART_TYPE.DOUGHNUT, XL_CHART_TYPE.DOUGHNUT_EXPLODED,
                )
                is_bar  = ctype in (
                    XL_CHART_TYPE.BAR_CLUSTERED, XL_CHART_TYPE.BAR_STACKED,
                    XL_CHART_TYPE.BAR_STACKED_100,
                    XL_CHART_TYPE.COLUMN_CLUSTERED, XL_CHART_TYPE.COLUMN_STACKED,
                )
                is_line = ctype in (
                    XL_CHART_TYPE.LINE, XL_CHART_TYPE.LINE_MARKERS,
                )

                ser0 = series_list[0]

                if is_pie:
                    vals_  = ser0["vals"]
                    labels = ser0["cats"] or [f"Cat {i+1}" for i in range(len(vals_))]
                    clrs   = [ser0["pt_colors"].get(i, None) for i in range(len(vals_))]
                    # Fill None colors with defaults
                    default_clrs = plt.rcParams["axes.prop_cycle"].by_key()["color"]
                    clrs = [c if c else default_clrs[i % len(default_clrs)]
                            for i, c in enumerate(clrs)]

                    wedge_kw = {"width": 0.55} if "DOUGHNUT" in str(ctype) else {}
                    wedges, texts, autotexts = ax.pie(
                        vals_, labels=None, colors=clrs,
                        autopct="%1.1f%%", startangle=90,
                        wedgeprops=wedge_kw,
                        pctdistance=0.75,
                    )
                    for at in autotexts:
                        at.set_fontsize(8)
                    ax.legend(wedges, labels, loc="center left",
                              bbox_to_anchor=(1, 0.5), fontsize=8)
                    if "DOUGHNUT" in str(ctype):
                        total = sum(vals_)
                        ax.text(0, 0, f"{total:.0f}", ha="center", va="center",
                                fontsize=12, fontweight="bold")
                    ax.set_title(ser0["name"] or "Chart", fontsize=11, pad=10)

                elif is_bar or (not is_line):
                    # Default: grouped bar chart
                    cats    = ser0["cats"] or [str(i+1) for i in range(len(ser0["vals"]))]
                    x       = np.arange(len(cats))
                    n_ser   = len(series_list)
                    width   = 0.8 / max(n_ser, 1)
                    offsets = np.linspace(-(n_ser-1)*width/2, (n_ser-1)*width/2, n_ser)
                    clr_cycle = plt.rcParams["axes.prop_cycle"].by_key()["color"]

                    for si2, ser2 in enumerate(series_list):
                        clr = ser2["ser_color"] or clr_cycle[si2 % len(clr_cycle)]
                        ax.bar(x + offsets[si2], ser2["vals"], width,
                               label=ser2["name"], color=clr, alpha=0.88)

                    ax.set_xticks(x)
                    ax.set_xticklabels(cats, rotation=30, ha="right", fontsize=8)
                    ax.set_ylabel("Value", fontsize=9)
                    ax.tick_params(axis="y", labelsize=8)
                    ax.grid(axis="y", alpha=0.3, linestyle="--")
                    if any(s["name"] for s in series_list):
                        ax.legend(fontsize=8)

                else:
                    # Line
                    cats = ser0["cats"] or [str(i+1) for i in range(len(ser0["vals"]))]
                    for si2, ser2 in enumerate(series_list):
                        ax.plot(cats, ser2["vals"], marker="o", label=ser2["name"],
                                linewidth=1.8, markersize=4)
                    ax.set_xticks(range(len(cats)))
                    ax.set_xticklabels(cats, rotation=30, ha="right", fontsize=8)
                    ax.grid(alpha=0.3, linestyle="--")
                    ax.legend(fontsize=8)

                plt.tight_layout(pad=1.0)

                tf = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
                fig.savefig(tf.name, dpi=150, bbox_inches="tight",
                            facecolor="white")
                plt.close(fig)
                tf.close()
                tmp_imgs.append(tf.name)

                rl = RLImage(tf.name)
                scale = min(inner_w / rl.imageWidth,
                            (ph - 2*margin)*0.55 / rl.imageHeight, 1.0)
                rl.drawWidth  = rl.imageWidth  * scale
                rl.drawHeight = rl.imageHeight * scale
                rl.hAlign = "CENTER"
                story.append(Spacer(1, 4))
                story.append(rl)
                story.append(Spacer(1, 6))
                return True

            except Exception as e:
                # Fallback: show chart as text summary
                try:
                    chart = shape.chart
                    story.append(Paragraph(f"[Chart — {chart.chart_type}]", ST["h2"]))
                    for ser in chart.series:
                        story.append(Paragraph(f"• {ser.name}", ST["body"]))
                except Exception:
                    story.append(Paragraph("[Chart]", ST["body"]))
                return False

        def _render_table(tbl_obj):
            try:
                rows_data = []
                for ri, row in enumerate(tbl_obj.rows):
                    cells = []
                    for ci, cell in enumerate(row.cells):
                        try:   txt = cell.text.strip()
                        except: txt = ""
                        ps = ST["th"] if ri == 0 else ST["td"]
                        cells.append(Paragraph(_safe(txt) or " ", ps))
                    rows_data.append(cells)
                if not rows_data:
                    return
                t = Table(rows_data, repeatRows=1)
                t.setStyle(TableStyle([
                    ("GRID",         (0,0),(-1,-1), 0.5, colors.HexColor("#ccc")),
                    ("BACKGROUND",   (0,0),(-1,0),  colors.HexColor("#e8edf5")),
                    ("VALIGN",       (0,0),(-1,-1), "TOP"),
                    ("LEFTPADDING",  (0,0),(-1,-1), 4),
                    ("RIGHTPADDING", (0,0),(-1,-1), 4),
                    ("TOPPADDING",   (0,0),(-1,-1), 3),
                    ("BOTTOMPADDING",(0,0),(-1,-1), 3),
                ]))
                story.append(Spacer(1, 6))
                story.append(t)
                story.append(Spacer(1, 6))
            except Exception:
                pass

        def _render_text_frame(tf_obj, default_sty):
            for para in tf_obj.paragraphs:
                try:
                    text = para.text.strip()
                    if not text:
                        story.append(Spacer(1, 3))
                        continue
                    try:    level = int(para.level or 0)
                    except: level = 0
                    if level == 0:
                        is_bul = text[0] in BULLETS
                        ps     = ST["bul0"] if is_bul else default_sty
                        prefix = "• " if is_bul and text[0] not in BULLETS else ""
                    elif level == 1:
                        ps, prefix = ST["bul1"], "  ◦ "
                    else:
                        ps, prefix = ST["bul2"], "    ▪ "
                    parts = []
                    for run in para.runs:
                        try:
                            rt = run.text
                            if not rt: continue
                            s = _safe(rt)
                            try:
                                b, i = run.font.bold, run.font.italic
                                if b and i: s = f"<b><i>{s}</i></b>"
                                elif b:     s = f"<b>{s}</b>"
                                elif i:     s = f"<i>{s}</i>"
                            except Exception:
                                pass
                            parts.append(s)
                        except Exception:
                            pass
                    inner = prefix + ("".join(parts) if parts else _safe(prefix + text))
                    story.append(Paragraph(inner, ps))
                except Exception:
                    try:
                        raw = para.text.strip()
                        if raw:
                            story.append(Paragraph(_safe(raw), default_sty))
                    except Exception:
                        pass

        # Process each slide
        for sn, slide in enumerate(prs.slides, 1):
            shapes_sorted = sorted(
                slide.shapes,
                key=lambda sh: (
                    (sh.top  or 0),
                    (sh.left or 0),
                )
            )

            for shape in shapes_sorted:
                try:
                    # Picture
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try: _add_img_blob(shape.image.blob, shape.image.ext)
                        except Exception: pass
                        continue

                    # Chart
                    if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                        _render_chart(shape)
                        continue

                    # Table
                    if shape.has_table:
                        _render_table(shape.table)
                        continue

                    # Text frame
                    if not shape.has_text_frame:
                        continue
                    if not shape.text_frame.text.strip():
                        continue

                    is_title = shape.name.lower().startswith("title")

                    if is_title:
                        story.append(Paragraph(_safe(shape.text_frame.text.strip()),
                                               ST["title"]))
                        story.append(HRFlowable(width="100%", thickness=0.5,
                                               color=colors.HexColor("#1e3a8a"),
                                               spaceAfter=6))
                    else:
                        _render_text_frame(shape.text_frame, ST["body"])

                except Exception:
                    pass

            # Slide number footer
            story.append(Spacer(1, 6))
            story.append(Paragraph(f"— {sn} / {len(prs.slides)} —", ST["slide_n"]))
            if sn < len(prs.slides):
                story.append(PageBreak())

        if not story:
            story.append(Paragraph("(empty presentation)", ST["body"]))

        def _safe_title(p):
            try: return p.slides[0].shapes.title.text or ""
            except: return ""

        SimpleDocTemplate(
            dst, pagesize=(pw, ph),
            leftMargin=margin, rightMargin=margin,
            topMargin=margin, bottomMargin=margin,
            title=_safe_title(prs),
        ).build(story)

        for f in tmp_imgs:
            try: os.remove(f)
            except Exception: pass

    #  HTML → PDF  — lean stack (exe-friendly)

    def _html_to_pdf(self, src, dst):
        """
        HTML → PDF — lightweight strategy stack, PyInstaller-compatible.

        1. pdfkit    (wkhtmltopdf separate binary, optional)
        2. weasyprint (pip install weasyprint)
        3. reportlab  — parses HTML manually, never fails, no duplicate content
           (fitz insert_htmlbox removed — duplicates content on some HTML inputs)
        """
        src_path = Path(src)
        base_dir = src_path.parent

        with open(src, "r", encoding="utf-8", errors="replace") as f:
            html_raw = f.read()

        # Inline all local resources → self-contained HTML
        html = self._inline_all_resources(html_raw, base_dir)

        # Strategy 1: pdfkit (wkhtmltopdf)
        try:
            import pdfkit
            tmp = tempfile.NamedTemporaryFile(
                suffix=".html", delete=False, mode="w", encoding="utf-8")
            tmp.write(html); tmp.close()
            try:
                pdfkit.from_file(tmp.name, dst, options={
                    "enable-local-file-access": "",
                    "load-error-handling":       "ignore",
                    "load-media-error-handling": "ignore",
                    "print-media-type":          "",
                    "quiet":                     "",
                    "margin-top":    "15mm",
                    "margin-bottom": "15mm",
                    "margin-left":   "15mm",
                    "margin-right":  "15mm",
                })
                return
            finally:
                try: os.remove(tmp.name)
                except Exception: pass
        except Exception:
            pass

        # Strategy 2: weasyprint
        try:
            import warnings
            import weasyprint
            with warnings.catch_warnings():
                warnings.simplefilter("ignore")
                weasyprint.HTML(
                    string=html,
                    base_url=base_dir.as_uri()
                ).write_pdf(dst)
            return
        except Exception:
            pass

        # Strategy 3: reportlab — parse HTML, no duplicates
        self._reportlab_html_to_pdf(html, dst)

    def _inline_all_resources(self, html, base_dir):
        """Inline all local img/css/url resources as base64 data-URIs."""
        SRC_RE  = re.compile(r"""src=(['"])([^'"]+)\1""",  re.I)
        HREF_RE = re.compile(r"""href=(['"])([^'"]+)\1""", re.I)

        # 1. <link rel="stylesheet"> -> <style>
        def _link_to_style(m):
            tag = m.group(0)
            if "stylesheet" not in tag.lower():
                return tag
            hm = HREF_RE.search(tag)
            if not hm:
                return tag
            href = hm.group(2)
            if href.startswith(("http://", "https://", "data:", "//")):
                return tag
            css_path = base_dir / href
            if not css_path.exists():
                return tag
            try:
                css = css_path.read_text(encoding="utf-8", errors="replace")
                css = self._inline_css_urls(css, css_path.parent)
                return "<style>" + css + "</style>"
            except Exception:
                return tag
        html = re.sub(r"<link[^>]+>", _link_to_style, html, flags=re.I | re.S)

        # 2. url() inside <style> blocks
        def _style_block(m):
            return "<style>" + self._inline_css_urls(m.group(1), base_dir) + "</style>"
        html = re.sub(r"<style[^>]*>(.*?)</style>",
                      _style_block, html, flags=re.I | re.S)

        # 3. <img src="...">
        def _img_src(m):
            tag = m.group(0)
            sm  = SRC_RE.search(tag)
            if not sm:
                return tag
            val = sm.group(2)
            if val.startswith(("http://", "https://", "data:")):
                return tag
            p = base_dir / val
            if not p.exists():
                return tag
            try:
                b64 = _img_to_b64(p.read_bytes(), _mime_for_ext(p.suffix.lstrip(".")))
                return SRC_RE.sub('src="' + b64 + '"', tag, count=1)
            except Exception:
                return tag
        html = re.sub(r"<img[^>]+>", _img_src, html, flags=re.I | re.S)
        return html

    def _inline_css_urls(self, css, base_dir):
        """Replace url(path) in CSS with base64 data URIs."""
        URL_RE = re.compile(r"url\\(\\s*([\"']?)([^)'\"]+)\\1\\s*\\)", re.I)
        def repl(m):
            raw = m.group(2).strip()
            if raw.startswith(("http://", "https://", "data:")):
                return m.group(0)
            p = base_dir / raw
            if not p.exists():
                return m.group(0)
            try:
                b64 = _img_to_b64(p.read_bytes(), _mime_for_ext(p.suffix.lstrip(".")))
                return "url('" + b64 + "')"
            except Exception:
                return m.group(0)
        return URL_RE.sub(repl, css)

    def _reportlab_html_to_pdf(self, html, dst):
        """
        Reportlab HTML fallback — faithful CSS class + inline style support.
        Parses <style> blocks to extract .class rules (text-align, margin-left).
        Applies class + inline style to every <p>/<div>.
        <br> inside <p> becomes a real line break inside the paragraph.
        .pn (page number) divs are skipped.
        """
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.units import mm
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib import colors
        from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT, TA_JUSTIFY
        from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer,
                                        Image as RLImage, HRFlowable,
                                        Table, TableStyle)

        MAX_W = A4[0] - 40 * mm
        base  = getSampleStyleSheet()

        # 1. Parse <style> blocks → class rules dict
        class_rules = {}
        for sb in re.finditer(r'<style[^>]*>(.*?)</style>', html, re.I | re.S):
            css = sb.group(1)
            for rule in re.finditer(r'\.([\w-]+)\s*\{([^}]*)\}', css):
                cls   = rule.group(1)
                decls = {}
                for d in rule.group(2).split(';'):
                    d = d.strip()
                    if ':' not in d:
                        continue
                    p, _, v = d.partition(':')
                    decls[p.strip().lower()] = v.strip()
                class_rules[cls] = decls

        ALIGN_MAP = {
            'right':   TA_RIGHT,
            'center':  TA_CENTER,
            'left':    TA_LEFT,
            'justify': TA_JUSTIFY,
        }

        def _decls_to_align(decls):
            return ALIGN_MAP.get(decls.get('text-align', '').lower())

        def _decls_to_indent(decls):
            ml = decls.get('margin-left', '')
            m  = re.search(r'([\d.]+)px', ml)
            return float(m.group(1)) * 0.75 if m else 0.0

        # 2. Style cache
        _cache = {}

        def _sty(align=TA_LEFT, indent=0.0, size=10.5, bold=False):
            key = (align, round(indent), size, bold)
            if key not in _cache:
                _cache[key] = ParagraphStyle(
                    'RS%d' % len(_cache),
                    parent=base['Normal'],
                    fontName='Helvetica-Bold' if bold else 'Helvetica',
                    fontSize=size, leading=size * 1.55,
                    spaceAfter=5, alignment=align, leftIndent=indent,
                )
            return _cache[key]

        H_STYS = {
            'h1': _sty(size=16, bold=True),
            'h2': _sty(size=14, bold=True),
            'h3': _sty(size=12, bold=True),
            'h4': _sty(size=11, bold=True),
            'h5': _sty(size=10, bold=True),
            'h6': _sty(size=10, bold=True),
        }
        TD_S  = _sty(size=9)
        TH_S  = _sty(size=9, bold=True)
        PRE_S = ParagraphStyle('RPRE', parent=base['Code'],
                               fontName='Courier', fontSize=9, leading=13,
                               spaceAfter=6, backColor=colors.HexColor('#f5f5f5'))

        # 3. Resolve align+indent from class= and style=
        _CLS_RE  = re.compile(r'''class=['"]([^'"]+)['"]''',  re.I)
        _STY_RE  = re.compile(r'''style=['"]([^'"]*?)['"]''', re.I)
        _SRC_RE  = re.compile(r'''src=['"]([^'"]+)['"]''',    re.I)

        def _resolve(attrs):
            align  = None
            indent = 0.0
            cm = _CLS_RE.search(attrs)
            if cm:
                for c in cm.group(1).split():
                    rd = class_rules.get(c, {})
                    a  = _decls_to_align(rd)
                    if a is not None:
                        align = a
                    i = _decls_to_indent(rd)
                    if i:
                        indent = i
            sm = _STY_RE.search(attrs)
            if sm:
                decl_str = sm.group(1)
                rd2 = {}
                for d in decl_str.split(';'):
                    d = d.strip()
                    if ':' not in d:
                        continue
                    p, _, v = d.partition(':')
                    rd2[p.strip().lower()] = v.strip()
                a2 = _decls_to_align(rd2)
                if a2 is not None:
                    align = a2
                i2 = _decls_to_indent(rd2)
                if i2:
                    indent = i2
            return (align or TA_LEFT), indent

        # 4. Inline markup
        def _decode(t):
            return (t.replace('&nbsp;', '\xa0')
                     .replace('&amp;',  '&')
                     .replace('&lt;',   '<')
                     .replace('&gt;',   '>')
                     .replace('&quot;', '"')
                     .replace('&#39;',  "'"))

        def _inline(frag):
            h = frag
            # <span style="...">
            def _span(m):
                inner = m.group(2)
                sm2   = _STY_RE.search(m.group(1) or '')
                if sm2:
                    s  = sm2.group(1)
                    fw = re.search(r'font-weight\s*:\s*(\w+)',  s, re.I)
                    fi = re.search(r'font-style\s*:\s*(\w+)',   s, re.I)
                    fc = re.search(r'color\s*:\s*(#[\da-fA-F]{3,8}|\w+)', s, re.I)
                    fs = re.search(r'font-size\s*:\s*([\d.]+)px', s, re.I)
                    if fw and fw.group(1).lower() in ('bold','700','800','900'):
                        inner = '<b>' + inner + '</b>'
                    if fi and fi.group(1).lower() == 'italic':
                        inner = '<i>' + inner + '</i>'
                    if fc:
                        inner = '<font color="' + fc.group(1) + '">' + inner + '</font>'
                    if fs:
                        inner = '<font size="' + fs.group(1) + '">' + inner + '</font>'
                return inner

            h = re.sub(r'<span([^>]*)>(.*?)</span>',     _span,            h, flags=re.I|re.S)
            h = re.sub(r'<strong[^>]*>(.*?)</strong>',   r'<b>\1</b>',     h, flags=re.I|re.S)
            h = re.sub(r'<b[^>]*>(.*?)</b>',             r'<b>\1</b>',     h, flags=re.I|re.S)
            h = re.sub(r'<em[^>]*>(.*?)</em>',           r'<i>\1</i>',     h, flags=re.I|re.S)
            h = re.sub(r'<i[^>]*>(.*?)</i>',             r'<i>\1</i>',     h, flags=re.I|re.S)
            h = re.sub(r'<a[^>]*>(.*?)</a>',             r'\1',            h, flags=re.I|re.S)
            h = re.sub(r'<code[^>]*>(.*?)</code>',       r'\1',            h, flags=re.I|re.S)

            # CRITICAL: save <br/> BEFORE stripping all tags
            # otherwise <br> gets stripped and address lines collapse
            BR = '\x00BR\x00'
            h = re.sub(r'<br\s*/?>', BR, h, flags=re.I)

            # Strip remaining HTML tags
            h = re.sub(r'<[^>]+>', '', h)
            h = _decode(h)

            # Save b/i/font tags before HTML-escaping
            SAVED = {}
            for tag in ('<b>', '</b>', '<i>', '</i>',
                        '<font color=', '</font>'):
                key = '\x00T%d\x00' % len(SAVED)
                SAVED[key] = tag
                h = h.replace(tag, key)

            h = h.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')

            # Restore saved tags
            for key, tag in SAVED.items():
                h = h.replace(key, tag)

            # Restore <br/>
            h = h.replace(BR, '<br/>')

            return re.sub(r' {2,}', ' ', h).strip()

        # 5. Image helper
        tmp_imgs = []

        def _add_img(src_val):
            if not src_val.startswith('data:'):
                return
            try:
                b64  = src_val.split(',', 1)[1]
                raw  = base64.b64decode(b64)
                ext  = src_val[5:src_val.index(';')].split('/')[-1] or 'png'
                tf   = tempfile.NamedTemporaryFile(suffix='.' + ext, delete=False)
                tf.write(raw); tf.close()
                tmp_imgs.append(tf.name)
                rl   = RLImage(tf.name)
                if rl.imageWidth > MAX_W:
                    s = MAX_W / rl.imageWidth
                    rl.drawWidth  = rl.imageWidth  * s
                    rl.drawHeight = rl.imageHeight * s
                rl.hAlign = 'CENTER'
                story.append(Spacer(1, 4))
                story.append(rl)
                story.append(Spacer(1, 4))
            except Exception:
                pass

        # 6. Table helper
        def _parse_tbl(tbl_html):
            rows = []
            for rm in re.finditer(r'<tr[^>]*>(.*?)</tr>', tbl_html, re.I | re.S):
                cells = []; hdr = False
                for cm in re.finditer(r'<(td|th)[^>]*>(.*?)</(td|th)>',
                                      rm.group(1), re.I | re.S):
                    t = cm.group(1).lower()
                    cells.append(Paragraph(_inline(cm.group(2)) or ' ',
                                           TH_S if t == 'th' else TD_S))
                    if t == 'th': hdr = True
                if cells: rows.append((cells, hdr))
            if not rows: return None
            n    = max(len(r) for r, _ in rows)
            data = []
            cmds = [('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#ccc')),
                    ('VALIGN', (0,0), (-1,-1), 'TOP'),
                    ('LEFTPADDING',   (0,0), (-1,-1), 4),
                    ('RIGHTPADDING',  (0,0), (-1,-1), 4),
                    ('TOPPADDING',    (0,0), (-1,-1), 3),
                    ('BOTTOMPADDING', (0,0), (-1,-1), 3)]
            for ri, (row, hdr) in enumerate(rows):
                while len(row) < n: row.append(Paragraph(' ', TD_S))
                data.append(row)
                if hdr: cmds.append(('BACKGROUND', (0,ri), (-1,ri),
                                     colors.HexColor('#f0f0f0')))
            t = Table(data, repeatRows=1)
            t.setStyle(TableStyle(cmds))
            return t

        # 7. Strip head / get body
        h2 = re.sub(r'<head[^>]*>.*?</head>',     '', html, flags=re.I | re.S)
        h2 = re.sub(r'<script[^>]*>.*?</script>', '', h2,   flags=re.I | re.S)
        bm = re.search(r'<body[^>]*>(.*?)</body>', h2,       re.I | re.S)
        body = bm.group(1) if bm else h2

        story = []

        # 8. Tokenise
        TAG_RE = re.compile(r'<(/?)(\w+)((?:\s[^>]*)?)/?>',  re.I)
        pos    = 0

        for m in TAG_RE.finditer(body):
            pos     = m.end()
            closing = m.group(1)
            tag     = m.group(2).lower()
            attrs   = m.group(3) or ''

            if tag == 'hr' and not closing:
                story.append(HRFlowable(width='100%', thickness=0.5,
                                        color=colors.HexColor('#ccc'),
                                        spaceAfter=5, spaceBefore=5))

            elif tag in ('h1','h2','h3','h4','h5','h6') and not closing:
                em = re.search(r'</' + tag + r'\s*>', body[pos:], re.I)
                if em:
                    txt = _inline(body[pos:pos + em.start()])
                    pos += em.end()
                    if txt:
                        story.append(Paragraph(txt, H_STYS.get(tag, H_STYS['h4'])))

            elif tag == 'p' and not closing:
                align, indent = _resolve(attrs)
                em = re.search(r'</p\s*>', body[pos:], re.I)
                if em:
                    inner = body[pos:pos + em.start()]
                    pos  += em.end()
                    # Images inside <p>
                    for im in re.finditer(r'<img[^>]+>', inner, re.I):
                        sm = _SRC_RE.search(im.group(0))
                        if sm: _add_img(sm.group(1))
                    txt = _inline(inner)
                    if txt.strip():
                        story.append(Paragraph(txt, _sty(align, indent)))

            elif tag == 'div' and not closing:
                # Skip .pn (page number) divs entirely
                cm2 = _CLS_RE.search(attrs)
                if cm2 and 'pn' in cm2.group(1).split():
                    em = re.search(r'</div\s*>', body[pos:], re.I)
                    if em: pos += em.end()

            elif tag == 'li' and not closing:
                em = re.search(r'</li\s*>', body[pos:], re.I)
                if em:
                    txt = _inline(body[pos:pos + em.start()])
                    pos += em.end()
                    if txt.strip():
                        story.append(Paragraph('• ' + txt, _sty()))

            elif tag == 'blockquote' and not closing:
                em = re.search(r'</blockquote\s*>', body[pos:], re.I)
                if em:
                    txt = _inline(body[pos:pos + em.start()])
                    pos += em.end()
                    if txt:
                        story.append(Paragraph(txt, _sty(indent=18)))

            elif tag == 'pre' and not closing:
                em = re.search(r'</pre\s*>', body[pos:], re.I)
                if em:
                    raw = re.sub(r'<[^>]+>', '', body[pos:pos + em.start()])
                    pos += em.end()
                    safe = (_decode(raw)
                            .replace('&', '&amp;')
                            .replace('<', '&lt;')
                            .replace('>', '&gt;')
                            .replace('\n', '<br/>'))
                    story.append(Paragraph(safe, PRE_S))

            elif tag == 'img' and not closing:
                sm = _SRC_RE.search(attrs)
                if sm: _add_img(sm.group(1))

            elif tag == 'table' and not closing:
                em = re.search(r'</table\s*>', body[pos:], re.I)
                if em:
                    tbl = _parse_tbl(body[pos:pos + em.start()])
                    pos += em.end()
                    if tbl:
                        story.append(Spacer(1, 6))
                        story.append(tbl)
                        story.append(Spacer(1, 6))

        if not story:
            story.append(Paragraph('(empty)', _sty()))

        SimpleDocTemplate(dst, pagesize=A4,
                          leftMargin=20*mm, rightMargin=20*mm,
                          topMargin=20*mm, bottomMargin=20*mm).build(story)
        for f in tmp_imgs:
            try: os.remove(f)
            except Exception: pass

    #  PDF → HTML — peak quality

    def _pdf_to_html(self, src, dst):
        """
        PDF → HTML — self-contained, faithful layout.

        Key improvements:
        - Margins reconstructed from block X positions relative to page width
          → text appears centred with proper left/right margins, not stuck left
        - Multi-column detection: blocks with x0 > 50% of page width → right col
        - Justified text for body paragraphs
        - Line-height and letter-spacing tuned per block font size
        - Consecutive lines of the same block merged into one <p> (no <br> soup)
        - Bold/italic/colour/size spans preserved
        - Images at their exact position in the flow
        - Superscript/subscript detected via vertical origin offset
        """
        import fitz

        doc  = fitz.open(src)
        name = Path(dst).stem

        # Page geometry
        _widths  = [p.rect.width  for p in doc]
        _heights = [p.rect.height for p in doc]
        MED_W  = sorted(_widths)[len(_widths)  // 2] if _widths  else 595.0
        MED_H  = sorted(_heights)[len(_heights) // 2] if _heights else 842.0

        # CSS page div width: map PDF points → screen pixels (96dpi / 72pt = 1.333)
        CSS_W  = max(620, min(1080, int(MED_W * 1.333)))
        # Padding mirrors PDF margins (typically ~56pt on A4 → ~75px)
        PAD    = max(40, int(CSS_W * 0.075))
        INNER  = CSS_W - PAD * 2   # usable text width in px

        # Base font: A4 body text is usually 10–12pt
        BASE_F = max(11, min(15, round(MED_W / 50)))

        CSS = f"""<style>
  *, *::before, *::after {{ box-sizing: border-box; margin: 0; padding: 0; }}
  html {{ font-size: {BASE_F}px; }}
  body {{
    font-family: 'Georgia', 'Times New Roman', serif;
    background: #d8d8d8;
    color: #111;
    line-height: 1.65;
  }}
  /* ── page card ── */
  .page {{
    background: #fff;
    width: {CSS_W}px;
    margin: 28px auto;
    padding: {PAD}px;
    box-shadow: 0 3px 18px rgba(0,0,0,.18);
    border-radius: 3px;
    position: relative;
  }}
  /* ── page number ── */
  .pn {{
    font-size: 9px;
    color: #bbb;
    text-align: right;
    margin-bottom: 14px;
    font-family: 'Segoe UI', Arial, sans-serif;
  }}
  /* ── body paragraph ── */
  p {{
    margin: 0 0 6px 0;
    text-align: justify;
    hyphens: auto;
    word-spacing: 0.02em;
  }}
  p.center {{ text-align: center; }}
  p.right  {{ text-align: right;  }}
  /* ── headings ── */
  h1 {{ font-size: {BASE_F + 7}px; margin: 16px 0 8px; line-height: 1.25; }}
  h2 {{ font-size: {BASE_F + 5}px; margin: 14px 0 6px; line-height: 1.28; }}
  h3 {{ font-size: {BASE_F + 3}px; margin: 10px 0 5px; line-height: 1.3;  }}
  h4 {{ font-size: {BASE_F + 1}px; margin: 8px  0 4px; }}
  /* ── inline ── */
  img {{
    max-width: 100%;
    height: auto;
    display: block;
    margin: 14px auto;
  }}
  sup {{ font-size: 0.72em; vertical-align: super; }}
  sub {{ font-size: 0.72em; vertical-align: sub;   }}
  hr  {{ border: none; border-top: 1px solid #e0e0e0; margin: 18px 0; }}
  /* ── two-column blocks ── */
  .col-r {{ margin-left: 50%; }}
  /* ── indent levels ── */
  .ind1 {{ margin-left: {int(INNER * 0.05)}px; }}
  .ind2 {{ margin-left: {int(INNER * 0.10)}px; }}
  .ind3 {{ margin-left: {int(INNER * 0.15)}px; }}
</style>"""

        # Span → HTML
        def _span_html(span, page_origin_y):
            txt   = span.get("text", "")
            if not txt.strip():
                return txt  # preserve spaces
            flags = span.get("flags",  0)
            size  = span.get("size",   BASE_F)
            color = span.get("color",  0)
            orig  = span.get("origin", (0, 0))   # (x, baseline_y)

            styles = []
            if flags & 16:  styles.append("font-weight:700")
            if flags & 2:   styles.append("font-style:italic")

            # Font size: only annotate when significantly different from base
            ratio = size / BASE_F if BASE_F else 1
            if ratio >= 1.35:
                styles.append(f"font-size:{min(int(size * 1.333), 38)}px")
            elif ratio <= 0.78:
                styles.append(f"font-size:{max(int(size * 1.333), 8)}px")

            # Colour — skip near-black
            if color and color != 0:
                r2 = (color >> 16) & 0xFF
                g2 = (color >>  8) & 0xFF
                b2 =  color        & 0xFF
                if not (r2 < 30 and g2 < 30 and b2 < 30):
                    styles.append(f"color:#{r2:02x}{g2:02x}{b2:02x}")

            safe = (_safe_html(txt)
                    .replace("	", "    ")
                    .replace("  ", " &nbsp;"))

            # Superscript / subscript via baseline offset
            if len(orig) >= 2 and page_origin_y:
                # fitz: origin[1] is baseline y; compare to block bbox top
                pass   # handled at block level below

            if not styles:
                return safe
            style_str = ";".join(styles)
            return f'<span style="{style_str}">{safe}</span>'

        # Block → HTML
        def _block_html(block, page_w):
            """
            Convert a fitz text block to an HTML element.
            Returns (html_string, alignment_class).
            """
            bbox      = block.get("bbox", [0, 0, page_w, 0])
            x0, x1    = bbox[0], bbox[2]
            block_w   = x1 - x0

            # Collect all spans text to determine block characteristics
            all_spans = [sp for ln in block.get("lines", [])
                         for sp in ln.get("spans", [])]
            if not all_spans:
                return "", ""

            first_size = all_spans[0].get("size", BASE_F)
            avg_size   = sum(sp.get("size", BASE_F) for sp in all_spans) / len(all_spans)
            all_bold   = all(sp.get("flags", 0) & 16 for sp in all_spans)

            # Heading detection
            ratio = first_size / BASE_F if BASE_F else 1
            if ratio >= 1.6 or (ratio >= 1.3 and all_bold):
                tag = "h1"
            elif ratio >= 1.35 or (ratio >= 1.15 and all_bold):
                tag = "h2"
            elif ratio >= 1.15:
                tag = "h3"
            elif ratio >= 1.05 and all_bold:
                tag = "h4"
            else:
                tag = "p"

            # Alignment from X position
            # Margin from left edge of page as fraction of page width
            left_margin_frac  = x0 / page_w if page_w else 0
            right_margin_frac = (page_w - x1) / page_w if page_w else 0

            align_cls = ""
            indent_cls = ""

            if tag == "p":
                # Centre: both margins roughly equal and both > 15%
                if (abs(left_margin_frac - right_margin_frac) < 0.08
                        and left_margin_frac > 0.15):
                    align_cls = "center"
                # Right-aligned: big left margin, small right margin
                elif left_margin_frac > 0.45 and right_margin_frac < 0.12:
                    align_cls = "right"
                # Indented body
                elif 0.08 < left_margin_frac < 0.20:
                    indent_cls = "ind1"
                elif 0.20 <= left_margin_frac < 0.30:
                    indent_cls = "ind2"
                elif left_margin_frac >= 0.30:
                    indent_cls = "ind3"

            # Build inner HTML
            # Strategy: collect each line's text and width, then decide
            # whether to join with space (natural wrap) or <br> (intentional).
            line_texts  = []
            line_widths = []   # actual text width in points

            for line in block.get("lines", []):
                spans_html = "".join(
                    _span_html(sp, None) for sp in line.get("spans", [])
                )
                if not spans_html.strip():
                    continue
                line_texts.append(spans_html)
                # bbox width of this line
                lbbox = line.get("bbox", [0, 0, 0, 0])
                line_widths.append(lbbox[2] - lbbox[0])

            if not line_texts:
                return "", ""

            # Block text width (points)
            block_text_w = bbox[2] - bbox[0]

            # Decide join mode:
            # "wrap"  → lines are just natural word-wrap artefacts → join with space
            # "break" → each line is intentionally short (address, sig, etc.) → <br>
            #
            # Heuristic: if MOST lines are short (< 75% of block width),
            # they are intentional line breaks.
            # Exception: single-line blocks always join (nothing to join anyway).
            if len(line_texts) == 1:
                inner = line_texts[0]
            else:
                short_lines = sum(
                    1 for w in line_widths if block_text_w > 0 and w / block_text_w < 0.75
                )
                short_ratio = short_lines / len(line_widths)

                # Also: if the block itself is narrow (< 40% page width),
                # it's likely a sidebar or address block → always use <br>
                block_frac = block_text_w / page_w if page_w else 1

                if short_ratio >= 0.5 or block_frac < 0.40:
                    # Intentional line breaks → preserve them
                    inner = "<br>".join(line_texts)
                else:
                    # Natural word-wrap → join with space
                    inner = " ".join(line_texts)

            # Build class list
            classes = " ".join(filter(None, [align_cls, indent_cls]))
            cls_attr = f' class="{classes}"' if classes else ""

            return f"<{tag}{cls_attr}>{inner}</{tag}>", align_cls

        # Per-page rendering
        pages_html = []

        for page in doc:
            pn     = page.number + 1
            page_w = page.rect.width or MED_W

            # Collect images
            img_b64: dict[int, str] = {}
            for info in page.get_images(full=True):
                xref = info[0]
                if xref in img_b64:
                    continue
                try:
                    bi = doc.extract_image(xref)
                    img_b64[xref] = _img_to_b64(bi["image"],
                                                  _mime_for_ext(bi["ext"]))
                except Exception:
                    pass

            # Build items sorted top→bottom, left→right
            items = []
            for block in page.get_text("dict", sort=True).get("blocks", []):
                y0, x0 = block["bbox"][1], block["bbox"][0]
                if block.get("type") == 0:
                    html, _ = _block_html(block, page_w)
                    if html:
                        items.append((y0, x0, "txt", html))
                elif block.get("type") == 1:
                    xref = block.get("xref", 0)
                    if xref and xref in img_b64:
                        items.append((y0, x0, "img", img_b64[xref]))

            # Images referenced but not in dict blocks
            placed = {d for _, _, k, d in items if k == "img"}
            for xref, uri in img_b64.items():
                if uri not in placed:
                    items.append((9999, 0, "img", uri))

            items.sort(key=lambda x: (x[0], x[1]))

            body_parts = []
            for _, _, kind, data in items:
                if kind == "txt":
                    body_parts.append(data)
                elif kind == "tbl":
                    body_parts.append(data)
                else:
                    body_parts.append(f'<img src="{data}" alt="img p{pn}">')

            pages_html.append(
                f'<div class="page">'
                f'<div class="pn">Page {pn} / {len(doc)}</div>'
                + "\n".join(body_parts)
                + "</div>"
            )

        with open(dst, "w", encoding="utf-8") as f:
            f.write(
                "<!DOCTYPE html>\n"
                "<html lang=\"fr\">\n<head>\n"
                "  <meta charset=\"utf-8\">\n"
                f"  <title>{_safe_html(name)}</title>\n"
                f"  {CSS}\n"
                "</head>\n<body>\n"
                + "\n<hr>\n".join(pages_html)
                + "\n</body>\n</html>"
            )

    #  EPUB → PDF  — peak quality

    def _epub_to_pdf(self, src, dst):
        """
        EPUB → PDF strategies:
        1. pypandoc (pandoc binary) — best typographic output
        2. Native engine             — fully self-contained, no binary needed
           Handles: spine order, cover image, metadata title, h1-h6,
           p, li (ul/ol), blockquote, br, strong/em/b/i, inline images
           at correct position in text flow.
        """
        try:
            import pypandoc
            pypandoc.convert_file(src, "pdf", outputfile=dst)
            return
        except Exception:
            pass
        self._epub_to_pdf_native(src, dst)

    def _epub_to_pdf_native(self, src, dst):
        """
        Native EPUB → PDF — peak quality.
        Improvements over previous version:
        - Tables rendered as reportlab Table objects
        - <ol> with real numbering (1. 2. 3.)
        - <pre>/<code> in monospace box
        - <figure>/<figcaption> with caption
        - <span style="..."> inline CSS (font-size, font-weight, color)
        - Image path resolution handles ../../../ deep relative paths
        - CSS font-size extracted from embedded <style> blocks
        """
        from xml.etree import ElementTree as ET
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.units import mm
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib import colors
        from reportlab.platypus import (
            SimpleDocTemplate, Paragraph, Spacer,
            Image as RLImage, PageBreak, HRFlowable,
            KeepTogether, Table, TableStyle,
        )
        from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY, TA_LEFT

        base = getSampleStyleSheet()
        def _sty(name, **kw):
            parent = kw.pop("parent", base["Normal"])
            return ParagraphStyle(name, parent=parent, **kw)

        sty = {
            "title"  : _sty("ET",  parent=base["Title"],
                            fontSize=22, leading=28, spaceAfter=16,
                            textColor=colors.HexColor("#1a1a2e"), alignment=TA_CENTER),
            "author" : _sty("EA",  fontSize=13, leading=18, spaceAfter=6,
                            textColor=colors.HexColor("#555"), alignment=TA_CENTER),
            "h1"     : _sty("EH1", parent=base["Heading1"],
                            fontSize=18, leading=22, spaceBefore=14, spaceAfter=8,
                            textColor=colors.HexColor("#1a1a2e")),
            "h2"     : _sty("EH2", parent=base["Heading2"],
                            fontSize=15, leading=19, spaceBefore=10, spaceAfter=6,
                            textColor=colors.HexColor("#1e3a5f")),
            "h3"     : _sty("EH3", parent=base["Heading3"],
                            fontSize=13, leading=17, spaceBefore=8, spaceAfter=4,
                            textColor=colors.HexColor("#1e3a5f")),
            "h4"     : _sty("EH4", fontSize=11, leading=15, spaceBefore=6, spaceAfter=3,
                            fontName="Helvetica-Bold", textColor=colors.HexColor("#333")),
            "h5"     : _sty("EH5", fontSize=10.5, leading=14, spaceBefore=4, spaceAfter=2,
                            fontName="Helvetica-Bold"),
            "h6"     : _sty("EH6", fontSize=10, leading=13, spaceBefore=4, spaceAfter=2,
                            fontName="Helvetica-BoldOblique"),
            "body"   : _sty("EB",  fontSize=10.5, leading=16, spaceAfter=5,
                            alignment=TA_JUSTIFY),
            "bq"     : _sty("EBQ", fontSize=10, leading=15, spaceAfter=5,
                            leftIndent=24, rightIndent=12,
                            textColor=colors.HexColor("#555")),
            "li_ul"  : _sty("ELU", fontSize=10.5, leading=15, spaceAfter=3, leftIndent=16),
            "li_ol"  : _sty("ELO", fontSize=10.5, leading=15, spaceAfter=3, leftIndent=20),
            "pre"    : _sty("EPR", fontName="Courier", fontSize=9, leading=13,
                            spaceAfter=8, spaceBefore=4,
                            leftIndent=12, rightIndent=12,
                            backColor=colors.HexColor("#f5f5f5"),
                            borderColor=colors.HexColor("#ddd"),
                            borderWidth=0.5, borderPad=6),
            "caption": _sty("EC",  fontSize=9, leading=12, spaceAfter=8,
                            textColor=colors.HexColor("#777"), alignment=TA_CENTER),
            "td"     : _sty("ETd", fontSize=9, leading=12, spaceAfter=0),
            "th"     : _sty("ETh", fontSize=9, leading=12, spaceAfter=0,
                            fontName="Helvetica-Bold"),
        }

        H_STYS = {"h1":"h1","h2":"h2","h3":"h3","h4":"h4","h5":"h5","h6":"h6"}

        story    = []
        tmp_imgs = []
        PAGE_W   = A4[0] - 40 * mm

        def _decode(text):
            return (text
                    .replace("&amp;",  "&").replace("&lt;",   "<")
                    .replace("&gt;",   ">").replace("&quot;", '"')
                    .replace("&apos;", "'").replace("&nbsp;", "\xa0"))

        def _rl_safe(text):
            d = _decode(text)
            return (d.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;"))

        def _css_to_inline(style_str):
            """Parse a CSS style attribute into reportlab-compatible markup hints."""
            bold = False; italic = False; color = None; size = None
            for decl in style_str.split(";"):
                decl = decl.strip()
                if not decl or ":" not in decl:
                    continue
                prop, _, val = decl.partition(":")
                prop = prop.strip().lower()
                val  = val.strip()
                if prop == "font-weight" and val in ("bold","700","800","900"):
                    bold = True
                elif prop == "font-style" and val == "italic":
                    italic = True
                elif prop == "color":
                    color = val
                elif prop == "font-size":
                    try:
                        size = float(re.sub(r"[^0-9.]","",val))
                    except Exception:
                        pass
            return bold, italic, color, size

        def _inline_markup(frag):
            """Convert inline HTML to reportlab XML, preserving b/i/span styles."""
            h = frag
            # Handle <span style="...">
            def _span_style(m):
                style_attr = re.search(r'style=["\']([^"\']*)["\']', m.group(1) or "", re.I)
                bold, italic, color, size = (False, False, None, None)
                if style_attr:
                    bold, italic, color, size = _css_to_inline(style_attr.group(1))
                inner = m.group(2)
                if bold:   inner = f"<b>{inner}</b>"
                if italic: inner = f"<i>{inner}</i>"
                return inner

            h = re.sub(r"<span([^>]*)>(.*?)</span>", _span_style, h, flags=re.I|re.S)
            h = re.sub(r"<strong[^>]*>(.*?)</strong>", r"<b>\1</b>", h, flags=re.I|re.S)
            h = re.sub(r"<b[^>]*>(.*?)</b>",           r"<b>\1</b>", h, flags=re.I|re.S)
            h = re.sub(r"<em[^>]*>(.*?)</em>",         r"<i>\1</i>", h, flags=re.I|re.S)
            h = re.sub(r"<i[^>]*>(.*?)</i>",           r"<i>\1</i>", h, flags=re.I|re.S)
            h = re.sub(r"<br\s*/?>",                   " ",           h, flags=re.I)
            h = re.sub(r"<a[^>]*>(.*?)</a>",           r"\1",         h, flags=re.I|re.S)
            h = re.sub(r"<code[^>]*>(.*?)</code>",     r"<font name='Courier'>\1</font>",
                        h, flags=re.I|re.S)
            h = re.sub(r"<[^>]+>", "", h)
            h = _decode(h)
            # Escape & restore b/i/font tags
            h = h.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")
            for tag in ("b","i"):
                h = (h.replace(f"&lt;{tag}&gt;", f"<{tag}>")
                      .replace(f"&lt;/{tag}&gt;", f"</{tag}>"))
            h = (h.replace("&lt;font name='Courier'&gt;", "<font name='Courier'>")
                  .replace("&lt;/font&gt;", "</font>"))
            return re.sub(r"\s{2,}", " ", h).strip()

        def _add_img(raw_bytes, ext, alt=""):
            try:
                tf = tempfile.NamedTemporaryFile(suffix=f".{ext or 'png'}", delete=False)
                tf.write(raw_bytes); tf.close()
                tmp_imgs.append(tf.name)
                rl    = RLImage(tf.name)
                scale = min(PAGE_W / rl.imageWidth,
                            (A4[1] * 0.5) / rl.imageHeight, 1.0)
                rl.drawWidth  = rl.imageWidth  * scale
                rl.drawHeight = rl.imageHeight * scale
                rl.hAlign     = "CENTER"
                story.append(Spacer(1, 6))
                story.append(rl)
                if alt and alt.lower() not in ("", "image", "cover"):
                    story.append(Paragraph(_rl_safe(alt), sty["caption"]))
                story.append(Spacer(1, 6))
                return True
            except Exception:
                return False

        def _resolve_img(src_attr, chap_dir, img_data):
            """Resolve image src with robust path normalisation."""
            if not src_attr:
                return None
            if src_attr.startswith("data:"):
                try:
                    b64  = src_attr.split(",", 1)[1]
                    raw  = base64.b64decode(b64)
                    ext  = re.search(r"data:image/(\w+)", src_attr)
                    return raw, (ext.group(1) if ext else "png")
                except Exception:
                    return None

            # Normalise path
            raw_p  = src_attr.split("?")[0].split("#")[0]
            # Resolve relative path from chapter directory
            try:
                from pathlib import PurePosixPath
                resolved = str(PurePosixPath(chap_dir) / raw_p)
            except Exception:
                resolved = f"{chap_dir}/{raw_p}"

            candidates = set()
            for c in [raw_p, resolved, raw_p.lstrip("/"),
                       resolved.lstrip("/")]:
                c_norm = c.replace("\\", "/")
                # Remove leading ./
                while c_norm.startswith("./"):
                    c_norm = c_norm[2:]
                candidates.add(c_norm)
                # Also try removing leading path segments to find by filename
                parts = c_norm.split("/")
                if len(parts) > 1:
                    candidates.add(parts[-1])

            for c in candidates:
                b = img_data.get(c)
                if b:
                    return b, Path(c).suffix.lstrip(".")
                # Fuzzy match by filename
                fname = c.split("/")[-1]
                for k in img_data:
                    if k.split("/")[-1] == fname:
                        return img_data[k], Path(fname).suffix.lstrip(".")
            return None

        def _parse_table(table_html, chap_dir, img_data):
            """Render <table> as reportlab Table."""
            rows_data = []
            for row_m in re.finditer(r"<tr[^>]*>(.*?)</tr>", table_html, re.I|re.S):
                cells = []
                is_hdr = False
                for cell_m in re.finditer(r"<(td|th)[^>]*>(.*?)</(td|th)>",
                                           row_m.group(1), re.I|re.S):
                    tag  = cell_m.group(1).lower()
                    text = _inline_markup(cell_m.group(2))
                    ps   = sty["th"] if tag == "th" else sty["td"]
                    cells.append(Paragraph(text or " ", ps))
                    if tag == "th":
                        is_hdr = True
                if cells:
                    rows_data.append((cells, is_hdr))
            if not rows_data:
                return None
            n_cols = max(len(r) for r, _ in rows_data)
            tbl_data = []
            style_cmds = [
                ("GRID",         (0,0),(-1,-1), 0.5, colors.HexColor("#ccc")),
                ("VALIGN",       (0,0),(-1,-1), "TOP"),
                ("LEFTPADDING",  (0,0),(-1,-1), 4),
                ("RIGHTPADDING", (0,0),(-1,-1), 4),
                ("TOPPADDING",   (0,0),(-1,-1), 3),
                ("BOTTOMPADDING",(0,0),(-1,-1), 3),
            ]
            for ri, (row, is_hdr) in enumerate(rows_data):
                while len(row) < n_cols:
                    row.append(Paragraph(" ", sty["td"]))
                tbl_data.append(row)
                if is_hdr:
                    style_cmds.append(("BACKGROUND",(0,ri),(-1,ri),
                                       colors.HexColor("#e8edf5")))
            t = Table(tbl_data, repeatRows=1)
            t.setStyle(TableStyle(style_cmds))
            return t

        def _parse_chapter(html_raw, chap_dir, img_data):
            html = re.sub(r'\s+xmlns(?::\w+)?=["\'][^"\']*["\']', "", html_raw)
            bm   = re.search(r"<body[^>]*>(.*?)</body>", html, re.I|re.S)
            body = bm.group(1) if bm else html
            body = re.sub(r"<head[^>]*>.*?</head>",     "", body, flags=re.I|re.S)
            body = re.sub(r"<script[^>]*>.*?</script>", "", body, flags=re.I|re.S)
            body = re.sub(r"<style[^>]*>.*?</style>",   "", body, flags=re.I|re.S)

            TAG_RE = re.compile(
                r"<(/?)(\w+)((?:\s[^>]*)?)/?>", re.I
            )
            pos      = 0
            buf      = ""
            ol_count = [0]   # stack for ordered list numbering

            def flush_buf():
                nonlocal buf
                text = _inline_markup(buf)
                buf  = ""
                if text.strip():
                    story.append(Paragraph(text, sty["body"]))

            for m in TAG_RE.finditer(body):
                buf += body[pos:m.start()]
                pos  = m.end()
                closing, tag, attrs_str = m.group(1), m.group(2).lower(), m.group(3) or ""

                block_tags = {"p","div","h1","h2","h3","h4","h5","h6",
                              "br","hr","li","img","figure","figcaption",
                              "table","ul","ol","pre","blockquote","section",
                              "article","header","footer","aside","nav"}
                if tag in block_tags:
                    flush_buf()

                if tag == "br" and not closing:
                    story.append(Spacer(1, 4))

                elif tag == "hr" and not closing:
                    story.append(HRFlowable(width="100%", thickness=0.5,
                                            color=colors.HexColor("#ccc"),
                                            spaceAfter=6, spaceBefore=6))

                elif tag in H_STYS and not closing:
                    end_m = re.search(rf"</{tag}\s*>", body[pos:], re.I)
                    if end_m:
                        inner = _inline_markup(body[pos:pos+end_m.start()])
                        pos  += end_m.end()
                        if inner:
                            story.append(Spacer(1, 4))
                            story.append(Paragraph(inner, sty[H_STYS[tag]]))

                elif tag == "p" and not closing:
                    end_m = re.search(r"</p\s*>", body[pos:], re.I)
                    if end_m:
                        inner_html = body[pos:pos+end_m.start()]
                        pos       += end_m.end()
                        # Inline images inside <p>
                        for img_m in re.finditer(r"<img[^>]+>", inner_html, re.I):
                            sm = re.search(r"""src=(['"])([^'"]+)\1""", img_m.group(0), re.I)
                            if sm:
                                res = _resolve_img(sm.group(2), chap_dir, img_data)
                                if res:
                                    alt_m = re.search(r"""alt=(['"])([^'"]*)\1""",
                                                      img_m.group(0), re.I)
                                    _add_img(res[0], res[1],
                                             alt_m.group(2) if alt_m else "")
                        text = _inline_markup(inner_html)
                        if text.strip():
                            story.append(Paragraph(text, sty["body"]))

                elif tag in ("div","section","article","header",
                             "footer","aside","nav") and not closing:
                    pass  # just let content flow through

                elif tag == "blockquote" and not closing:
                    end_m = re.search(r"</blockquote\s*>", body[pos:], re.I)
                    if end_m:
                        inner = _inline_markup(body[pos:pos+end_m.start()])
                        pos  += end_m.end()
                        if inner:
                            story.append(Paragraph(inner, sty["bq"]))

                elif tag == "pre" and not closing:
                    end_m = re.search(r"</pre\s*>", body[pos:], re.I)
                    if end_m:
                        raw_pre = body[pos:pos+end_m.start()]
                        pos    += end_m.end()
                        # Strip only tags, preserve whitespace
                        text = re.sub(r"<[^>]+>", "", raw_pre)
                        text = _decode(text)
                        safe = text.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")
                        story.append(Paragraph(safe.replace("\n","<br/>"), sty["pre"]))

                elif tag == "ol" and not closing:
                    ol_count[0] = 0
                elif tag == "ul" and not closing:
                    ol_count[0] = -1   # sentinel: unordered

                elif tag == "li" and not closing:
                    end_m = re.search(r"</li\s*>", body[pos:], re.I)
                    if end_m:
                        inner = _inline_markup(body[pos:pos+end_m.start()])
                        pos  += end_m.end()
                        if inner.strip():
                            if ol_count[0] >= 0:
                                ol_count[0] += 1
                                prefix = f"{ol_count[0]}. "
                                ps = sty["li_ol"]
                            else:
                                prefix = "• "
                                ps = sty["li_ul"]
                            story.append(Paragraph(prefix + inner, ps))

                elif tag == "img" and not closing:
                    sm = re.search(r"""src=(['"])([^'"]+)\1""", attrs_str, re.I)
                    if sm:
                        res = _resolve_img(sm.group(2), chap_dir, img_data)
                        if res:
                            alt_m = re.search(r"""alt=(['"])([^'"]*)\1""", attrs_str, re.I)
                            _add_img(res[0], res[1], alt_m.group(2) if alt_m else "")

                elif tag == "figure" and not closing:
                    pass  # content flows through

                elif tag == "figcaption" and not closing:
                    end_m = re.search(r"</figcaption\s*>", body[pos:], re.I)
                    if end_m:
                        inner = _inline_markup(body[pos:pos+end_m.start()])
                        pos  += end_m.end()
                        if inner:
                            story.append(Paragraph(inner, sty["caption"]))

                elif tag == "table" and not closing:
                    end_m = re.search(r"</table\s*>", body[pos:], re.I)
                    if end_m:
                        tbl = _parse_table(body[pos:pos+end_m.start()],
                                           chap_dir, img_data)
                        pos += end_m.end()
                        if tbl:
                            story.append(Spacer(1, 8))
                            story.append(tbl)
                            story.append(Spacer(1, 8))

            buf += body[pos:]
            flush_buf()

        # Open EPUB
        with zipfile.ZipFile(src) as zf:
            names = zf.namelist()
            spine_order  = []
            book_title   = Path(src).stem
            book_authors = []
            cover_data   = None
            opf_dir      = ""

            try:
                container = zf.read("META-INF/container.xml").decode("utf-8","replace")
                opf_m = re.search(r"full-path=[\"']([^\"']+\.opf)[\"']", container, re.I)
                if opf_m:
                    opf_path = opf_m.group(1)
                    opf_dir  = str(Path(opf_path).parent)
                    opf_xml  = zf.read(opf_path).decode("utf-8","replace")
                    root     = ET.fromstring(opf_xml)
                    ns_opf   = {"o":"http://www.idpf.org/2007/opf",
                                "dc":"http://purl.org/dc/elements/1.1/"}
                    t_el = root.find(".//dc:title", ns_opf)
                    if t_el is not None and t_el.text:
                        book_title = t_el.text.strip()
                    for cr in root.findall(".//dc:creator", ns_opf):
                        if cr.text:
                            book_authors.append(cr.text.strip())
                    manifest = {}
                    for item in root.findall(".//o:item", ns_opf):
                        iid  = item.get("id","")
                        href = item.get("href","")
                        mtype= item.get("media-type","")
                        full = f"{opf_dir}/{href}".lstrip("/")
                        manifest[iid] = {"href":href,"full":full,"type":mtype}
                    cover_id = None
                    mc = root.find(".//o:meta[@name='cover']", ns_opf)
                    if mc is not None:
                        cover_id = mc.get("content","")
                    if not cover_id:
                        for iid,v in manifest.items():
                            if "cover" in iid.lower() and "image" in v["type"]:
                                cover_id = iid; break
                    if cover_id and cover_id in manifest:
                        cpath = manifest[cover_id]["full"]
                        try:
                            cover_data = (zf.read(cpath),
                                          Path(cpath).suffix.lstrip("."))
                        except Exception:
                            pass
                    for ref in root.findall(".//o:itemref", ns_opf):
                        iid  = ref.get("idref","")
                        if iid in manifest:
                            full = manifest[iid]["full"]
                            if full in names:
                                spine_order.append(full)
            except Exception:
                pass

            if not spine_order:
                spine_order = sorted(
                    n for n in names
                    if n.endswith((".xhtml",".html",".htm"))
                )

            # Pre-load all images
            img_data = {}
            for n in names:
                ext = Path(n).suffix.lower().lstrip(".")
                if ext in ("png","jpg","jpeg","gif","webp","bmp","svg"):
                    try:
                        key = n.replace("\\","/").lstrip("/")
                        img_data[key] = zf.read(n)
                    except Exception:
                        pass

            # Cover page
            if cover_data:
                _add_img(cover_data[0], cover_data[1], "cover")
                story.append(PageBreak())

            # Title page
            story.append(Spacer(1, 30*mm))
            story.append(Paragraph(_rl_safe(book_title), sty["title"]))
            if book_authors:
                story.append(Spacer(1, 6))
                for auth in book_authors:
                    story.append(Paragraph(_rl_safe(auth), sty["author"]))
            story.append(PageBreak())

            # Chapters
            for chap in spine_order:
                try:
                    raw      = zf.read(chap).decode("utf-8","replace")
                    chap_dir = str(Path(chap).parent).replace("\\","/")
                    _parse_chapter(raw, chap_dir, img_data)
                    story.append(PageBreak())
                except Exception:
                    continue

        if len(story) <= 2:
            raise RuntimeError("No content could be extracted from this EPUB.")

        def _safe_str(s):
            return s if isinstance(s, str) else ""

        SimpleDocTemplate(
            dst, pagesize=A4,
            leftMargin=22*mm, rightMargin=22*mm,
            topMargin=22*mm, bottomMargin=22*mm,
            title=_safe_str(book_title),
            author=", ".join(book_authors) if book_authors else "",
        ).build(story)

        for f in tmp_imgs:
            try: os.remove(f)
            except Exception: pass

    #  Images — max quality, EXIF preserved

    def _image_convert(self, src, dst, conversion_type):
        from PIL import Image, ImageOps
        img = Image.open(src)
        try: img = ImageOps.exif_transpose(img)
        except Exception: pass

        ext = Path(dst).suffix.lower().lstrip(".")
        if ext in ("jpg", "jpeg"):
            if img.mode in ("RGBA", "P", "LA"):
                bg   = Image.new("RGB", img.size, (255, 255, 255))
                if img.mode == "P": img = img.convert("RGBA")
                mask = img.split()[-1] if img.mode == "RGBA" else None
                bg.paste(img, mask=mask); img = bg
            elif img.mode != "RGB":
                img = img.convert("RGB")
            img.save(dst, "JPEG", quality=95, subsampling=0, optimize=True)
        elif ext == "png":
            if img.mode not in ("RGB", "RGBA", "L", "LA", "P"):
                img = img.convert("RGBA")
            img.save(dst, "PNG", optimize=True, compress_level=6)
        elif ext == "webp":
            img.save(dst, "WEBP", quality=92, method=6)
        elif ext == "bmp":
            if img.mode not in ("RGB", "RGBA"):
                img = img.convert("RGB")
            img.save(dst, "BMP")
        else:
            img.save(dst)

    # All standard sizes used by Windows + macOS + Linux icons.
    _ICO_SIZES = [(16,16),(24,24),(32,32),(48,48),(64,64),(128,128),(256,256)]

    def _image_to_ico(self, src, dst, conversion_type=None):
        """
        Convert any Pillow-readable image to a multi-resolution .ico file.

        Key fix: Pillow's `append_images=` parameter is silently ignored for
        ICO — only the first frame gets written.  The correct approach is to
        pass `sizes=` directly on the *original* full-resolution image and let
        Pillow handle the downsampling internally using its own LANCZOS pass.

        Only sizes ≤ the source dimensions are included to avoid upscaling.
        """
        from PIL import Image, ImageOps

        img = Image.open(src)

        try:
            img = ImageOps.exif_transpose(img)
        except Exception:
            pass

        img = img.convert("RGBA")

        src_min = min(img.size)
        sizes = [(w, h) for (w, h) in self._ICO_SIZES if w <= src_min]
        if not sizes:
            sizes = [(16, 16)]

        img.save(dst, format="ICO", sizes=sizes)

    def _heic_to_png(self, src, dst, conversion_type=None):
        # 1. pillow-heif
        try:
            from pillow_heif import register_heif_opener
            from PIL import Image, ImageOps
            register_heif_opener()
            img = Image.open(src)
            try: img = ImageOps.exif_transpose(img)
            except Exception: pass
            if img.mode not in ("RGB", "RGBA"):
                img = img.convert("RGB")
            img.save(dst, "PNG", optimize=True, compress_level=6)
            return
        except Exception:
            pass
        # 2. pyheif
        try:
            import pyheif
            from PIL import Image
            heif_file = pyheif.read(src)
            img = Image.frombytes(heif_file.mode, heif_file.size,
                                  heif_file.data, "raw",
                                  heif_file.mode, heif_file.stride)
            img.save(dst, "PNG", optimize=True, compress_level=6)
            return
        except Exception:
            pass
        # 3. ImageMagick CLI
        subprocess.run(["magick", "convert", src, dst],
                       check=True, capture_output=True, timeout=120)

    #  Audio / Video — ffmpeg, auto-located, quality presets

    _FFMPEG_PRESETS = {
        "mp3":  ["-codec:a","libmp3lame","-q:a","2","-ar","44100"],
        "wav":  ["-codec:a","pcm_s16le","-ar","44100"],
        "aac":  ["-codec:a","aac","-b:a","192k","-ar","44100"],
        "ogg":  ["-codec:a","libvorbis","-q:a","5"],
        "flac": ["-codec:a","flac","-compression_level","8"],
        "mp4":  ["-codec:v","libx264","-crf","20","-preset","slow",
                 "-codec:a","aac","-b:a","192k","-movflags","+faststart"],
        "avi":  ["-codec:v","libxvid","-qscale:v","3",
                 "-codec:a","libmp3lame","-q:a","4"],
        "mkv":  ["-codec:v","libx264","-crf","20",
                 "-codec:a","aac","-b:a","192k"],
        "webm": ["-codec:v","libvpx-vp9","-crf","30","-b:v","0",
                 "-codec:a","libopus","-b:a","128k"],
    }

    @staticmethod
    def _find_ffmpeg():
        import shutil, sys
        found = shutil.which("ffmpeg")
        if found: return found
        candidates = [
            r"C:\ffmpeg\bin\ffmpeg.exe",
            r"C:\Program Files\ffmpeg\bin\ffmpeg.exe",
            r"C:\Program Files (x86)\ffmpeg\bin\ffmpeg.exe",
            os.path.join(os.environ.get("LOCALAPPDATA",""), "ffmpeg","bin","ffmpeg.exe"),
            os.path.join(os.environ.get("APPDATA",""),      "ffmpeg","bin","ffmpeg.exe"),
            os.path.join(getattr(sys,"_MEIPASS",""),         "ffmpeg.exe"),
            os.path.join(os.path.dirname(os.path.abspath(__file__)), "ffmpeg.exe"),
            "/usr/bin/ffmpeg", "/usr/local/bin/ffmpeg",
            "/opt/homebrew/bin/ffmpeg", "/snap/bin/ffmpeg",
        ]
        for c in candidates:
            if c and os.path.isfile(c): return c
        raise FileNotFoundError(
            "ffmpeg introuvable.\n"
            "Windows : winget install Gyan.FFmpeg puis redémarrer le terminal.\n"
            "Linux   : sudo apt install ffmpeg\n"
            "macOS   : brew install ffmpeg"
        )

    def _has_audio_stream(self, src: str, ffmpeg_bin: str) -> bool:
        """Return True if *src* contains at least one audio stream."""
        try:
            result = subprocess.run(
                [ffmpeg_bin, "-i", src],
                capture_output=True, timeout=30,
            )
            output = result.stderr.decode("utf-8", errors="replace")
            return "Audio:" in output
        except Exception:
            return True  # on error, let ffmpeg attempt the conversion anyway

    def _ffmpeg_convert(self, src, dst, conversion_type):
        target_ext = Path(dst).suffix.lower().lstrip(".")
        presets    = self._FFMPEG_PRESETS.get(target_ext, [])
        ffmpeg_bin = self._find_ffmpeg()

        # Check for audio track before video-to-audio extractions
        _AUDIO_EXTS = {"mp3", "wav", "aac", "ogg", "flac"}
        _VIDEO_EXTS = {".mp4", ".avi", ".mkv", ".webm", ".mov",
                       ".m4v", ".flv", ".wmv", ".ts", ".3gp"}
        if target_ext in _AUDIO_EXTS and Path(src).suffix.lower() in _VIDEO_EXTS:
            if not self._has_audio_stream(src, ffmpeg_bin):
                raise RuntimeError(
                    f"This file does not contain an audio track : {Path(src).name}"
                )

        try:
            import ffmpeg as fl
            kwargs: dict = {}
            if target_ext == "mp3":
                kwargs = {"acodec":"libmp3lame","audio_bitrate":"320k","ar":"44100"}
            elif target_ext == "mp4":
                kwargs = {"vcodec":"libx264","crf":"20","preset":"slow",
                          "acodec":"aac","audio_bitrate":"192k","movflags":"+faststart"}
            elif target_ext == "wav":
                kwargs = {"acodec":"pcm_s16le","ar":"44100"}
            elif target_ext == "aac":
                kwargs = {"acodec":"aac","audio_bitrate":"192k"}
            elif target_ext == "webm":
                kwargs = {"vcodec":"libvpx-vp9","crf":"30","b:v":"0",
                          "acodec":"libopus","audio_bitrate":"128k"}
            (fl.input(src).output(dst, **kwargs)
               .overwrite_output().run(cmd=ffmpeg_bin, quiet=True))
            return
        except Exception:
            pass

        subprocess.run(
            [ffmpeg_bin, "-y", "-i", src] + presets + [dst],
            check=True, capture_output=True, timeout=1800,
        )