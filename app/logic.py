"""
app/logic.py — AppLogicMixin
File Converter Pro

Contains all conversion, processing, and data-management methods of
FileConverterApp, extracted as a mixin for clarity.  Import order:

    logic.py  (AppLogicMixin)
        ↑
    ui.py     (AppUIMixin)
        ↑
    __init__.py (FileConverterApp)

Author: Hyacinthe
Version: 1.0
"""

"""
Main Application Logic - File Converter Pro

Core application window and conversion engine implementation.

Classes:
    FileConverterApp(QMainWindow):
        - Complete UI layout (File list, Conversion panel, Toolbar)
        - Multi-format conversion engine (PDF, Word, Images, Archives)
        - Integration with external systems (Achievements, Notifications)
        - Project management (Save/Load .fcproj)
    
    FadingMainWindow(FileConverterApp):
        - Subclass with fade-in animation support on launch

Key Functionalities:
    - Drag & Drop file handling
    - Batch conversions and operations
    - Encrypted PDF handling (Decrypt → Convert)
    - Office file optimization (Metadata removal, Compression)
    - Achievement system tracking (Unlocks, Stats, Ranks)

Design:
    - Modular imports (config, database, widgets, dialogs, etc.)
    - Signal/Slot architecture for decoupled communication
    - Robust resource path management (Dev + PyInstaller)

Author: Hyacinthe
Version: 1.0
"""

import sys
import os
import tempfile
import shutil
import io
import zipfile
from pathlib import Path
from PySide6.QtWidgets import (QWidget, QVBoxLayout, QLabel, QListWidgetItem,
                               QFileDialog, QMessageBox, QGroupBox, QScrollArea,
                               QLineEdit, QDialog, QDialogButtonBox, QTextEdit,
                               QTabWidget, QRadioButton)
from PySide6.QtCore import Qt, QTimer, QRect
from PySide6.QtGui import QIcon, QGuiApplication
from datetime import datetime
import time

# NOTE: fitz (PyMuPDF) is intentionally NOT imported here — it is heavy (~70 MB
# in RAM) and only needed inside specific PDF-processing methods. Each method
# does `import fitz` locally so the library loads on first use, not at startup.

try:
    from pdf2docx import Converter as _Converter
    Converter = _Converter
except ImportError as _e:
    Converter = None
    print(f"[IMPORT] pdf2docx not available: {_e}")

try:
    from PyPDF2 import PdfReader, PdfWriter
except ImportError as _e:
    PdfMerger = PdfReader = PdfWriter = None
    print(f"[IMPORT] PyPDF2 not available: {_e}")

try:
    from docx import Document
    from docx.shared import Inches
except ImportError as _e:
    Document = Inches = None
    print(f"[IMPORT] python-docx not available: {_e}")

try:
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.pdfgen import canvas
    from reportlab.lib.utils import simpleSplit
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
except ImportError as _e:
    letter = A4 = canvas = simpleSplit = None
    getSampleStyleSheet = ParagraphStyle = None
    SimpleDocTemplate = Paragraph = Spacer = None
    print(f"[IMPORT] reportlab not available: {_e}")

try:
    from PIL import Image
    OCR_AVAILABLE = True
except ImportError as _e:
    Image = None
    OCR_AVAILABLE = False
    print(f"[IMPORT] Pillow not available: {_e}")

from config import is_dark_mode_qt
from database import DatabaseManager
from translations import TranslationManager
from widgets import AnimatedCheckBox
from dialogs import (PasswordDialog, SplitDialog, CompressionDialog,
                     BatchConvertDialog, WordToPdfOptionsDialog,
                     PdfToWordDialog, TermsAndPrivacyDialog)
from conversion_worker import ConversionWorker
from achievements import AchievementSystem
from special_events_manager import SpecialEventsManager
from system_notifier import SystemNotifier
import re as _re

def _sanitize_xml(t: str) -> str:
    """Remove control characters that break ReportLab's XML parser."""
    return _re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]', '', t or '')


class AppLogicMixin:
    """Mixin: conversion engine and data logic for FileConverterApp."""

    def __init__(self, config_manager):
        super().__init__()
        self.config_manager = config_manager
        self.config = config_manager.load_config()
        self.current_language = self.config.get("language", "fr")
        if self.config.get("use_system_theme", True):
            self.dark_mode = is_dark_mode_qt()
        else:
            self.dark_mode = self.config.get("dark_mode", False)
        
        self.translation_manager = TranslationManager()
        self.translation_manager.set_language(self.current_language)
        
        self.files_list = []
        self.current_project = None
        self._project_data = {}
        self.preview_dialog = None
        self.temp_files = []
        
        self.terms_accepted = self.config.get("accepted_terms", False) and self.config.get("accepted_privacy", False)
        if not self.terms_accepted:
            self.show_terms_dialog()
            if not (self.config.get("accepted_terms") and self.config.get("accepted_privacy")):
                sys.exit(0)
        
        self.db_manager = DatabaseManager()
        self.template_settings = {}
        self.template_manager = None

        self.adv_db_manager = None
        try:
            from converter import AdvancedDatabaseManager
            self.adv_db_manager = AdvancedDatabaseManager()
        except Exception:
            pass
        
        # Achievements system
        self.achievement_system = AchievementSystem(self.config_manager)
        self.achievement_system.achievement_unlocked.connect(self.queue_achievement)
        self.achievement_system.rank_unlocked.connect(self.show_rank_popup)
        self.achievement_queue = []
        self.is_showing_achievement = False
        self.rank_queue = []

        # Record app launch
        self.achievement_system.record_app_launch()
        self.dashboard_dialog = None
        self.history_dialog = None
        self.templates_dialog = None
        self.achievements_dialog = None
        self.dark_mode_timer_start = None

        # If app starts in dark mode, initialize the timer
        if self.dark_mode:
            self.dark_mode_timer_start = time.time()
        self.dark_mode_timer = QTimer()
        self.dark_mode_timer.timeout.connect(self.update_dark_mode_time)
        self.dark_mode_timer.start(60000)  # Each minute

        self.system_notifier = SystemNotifier(self, self.current_language)
        self.system_notifier.set_translator(self.translation_manager)

        self.setup_ui()

        # Restore the window size/position
        _geom = self.config.get("window_geometry")
        _maximized = self.config.get("window_maximized", False)
        if _maximized:
            self.showMaximized()
        elif _geom:
            restored = QRect(
                _geom.get("x", 100),
                _geom.get("y", 100),
                _geom.get("width", 1400),
                _geom.get("height", 900),
            )
            self.setGeometry(restored)
            if not QGuiApplication.screenAt(self.geometry().center()):
                self.setGeometry(100, 100, 1400, 900)

        if self.dark_mode:
            self.apply_modern_dark_theme()
        else:
            self.apply_modern_light_theme()
        
        self.update_texts()
        self.special_events = SpecialEventsManager(self)
        _argv_project = None
        for _arg in sys.argv[1:]:
            if _arg.lower().endswith('.fcproj') and os.path.exists(_arg):
                _argv_project = _arg
                break

        if _argv_project:
            QTimer.singleShot(500, lambda p=_argv_project: self.open_project_file(p))
        elif self.config.get("auto_open_last_project") and self.current_project and os.path.exists(self.current_project):
            QTimer.singleShot(1000, self.open_last_project)
        if self.config.get("show_dashboard_on_startup", False):
            QTimer.singleShot(1500, self.show_dashboard)

        QTimer.singleShot(2000, self._check_donor_return)

    def _ensure_template_manager(self):
        """Initialize template_manager lazily and return it (or None on error)."""
        if self.template_manager is None:
            try:
                from templates import TemplateManager
                self.template_manager = TemplateManager(self.db_manager)
            except Exception as e:
                print(f"[WARN] Could not load TemplateManager: {e}")
        return self.template_manager

    def show_terms_dialog(self):
        dialog = TermsAndPrivacyDialog(self, self.current_language, dark_mode=self.dark_mode)
        result = dialog.exec()
        if result == QDialog.Accepted:
            from datetime import datetime
            now = datetime.now().isoformat()
            
            self.config["accepted_terms"] = True
            self.config["accepted_privacy"] = True
            
            if self.config.get("terms_acceptance_timestamp") is not None:
                self.config["terms_reacceptance_timestamp"] = now
                print(f"[TERMS DEBUG] ✅ Re-acceptance detected - terms_reacceptance_timestamp added: {now}")
            else:
                self.config["terms_acceptance_timestamp"] = now
                print(f"[TERMS DEBUG] ℹ️ First acceptance - terms_acceptance_timestamp set: {now}")
            
            self.config_manager.save_config(self.config)
            self.terms_accepted = True
        else:
            QMessageBox.information(
                self,
                self.translate_text("Conditions requises"),
                self.translate_text("Vous devez accepter les conditions d'utilisation et la politique de confidentialité pour utiliser cette application.")
            )
            sys.exit(0)

    def closeEvent(self, event):
        """Flush and close the persistent SQLite connection on exit."""
        try:
            self.db_manager.close()
        except Exception:
            pass
        super().closeEvent(event)

    def cleanup_temp_files(self):
        """Clean up temporary files"""
        for temp_file in self.temp_files:
            try:
                if os.path.exists(temp_file):
                    os.remove(temp_file)
                    print(f"[DEBUG] Temporary file deleted: {temp_file}")
            except Exception as e:
                print(f"[ERROR] Could not delete temporary file {temp_file}: {e}")
        self.temp_files.clear()

    def create_temp_file(self, suffix=".tmp"):
        import tempfile
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
        temp_path = temp_file.name
        temp_file.close()
        self.temp_files.append(temp_path)
        return temp_path

    def optimize_office_files(self, office_files, optimization_type, quality_level, remove_metadata, compress_images, keep_backup):
        if hasattr(self, 'active_templates') and 'office_optimization' in self.active_templates:
            del self.active_templates['office_optimization']
        """Optimize office and image files"""
        output_dir = self.get_output_directory()
        if not output_dir:
            return
        
        IMAGE_EXTS  = {'.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.webp', '.gif'}
        EXCEL_EXTS  = {'.xlsx', '.xls'}
        AUDIO_EXTS  = {'.mp3', '.wav', '.aac', '.flac', '.ogg'}
        VIDEO_EXTS  = {'.mp4', '.avi', '.mkv', '.webm'}
        WEB_EXTS    = {'.json', '.html', '.htm'}
        EPUB_EXTS   = {'.epub'}
        
        self.show_progress(True, self.translate_text("len_off").format(len(office_files)))
        
        success_count = 0
        total_original_size = 0
        total_compressed_size = 0
        start_time = datetime.now()
        
        for i, file_path in enumerate(office_files):
            try:
                file_ext = Path(file_path).suffix.lower()
                original_size = os.path.getsize(file_path)
                total_original_size += original_size
                
                # Determine output file path
                if keep_backup:
                    output_file = os.path.join(output_dir, f"optimized_{Path(file_path).name}")
                else:
                    output_file = file_path
                
                operation_start = datetime.now()
                
                # Determine compression level based on quality_level
                if quality_level == 0:
                    compression_level = "high"
                elif quality_level == 1:
                    compression_level = "normal"
                elif quality_level == 2:
                    compression_level = "very_reduced"
                
                if file_ext == '.pdf':
                    success = self.optimize_pdf_file(
                        file_path,
                        output_file,
                        compression_level,
                        remove_metadata,
                        compress_images
                    )
                elif file_ext in ['.docx', '.doc']:
                    success = self.optimize_word_file(
                        file_path,
                        output_file,
                        compression_level,
                        remove_metadata,
                        compress_images
                    )
                elif file_ext in ['.pptx', '.ppt']:
                    success = self.optimize_powerpoint_file(
                        file_path,
                        output_file,
                        compression_level,
                        remove_metadata,
                        compress_images
                    )
                elif file_ext in EXCEL_EXTS:
                    success = self.optimize_excel_file(
                        file_path,
                        output_file,
                        compression_level,
                        remove_metadata
                    )
                elif file_ext in IMAGE_EXTS:
                    success = self.optimize_image_file(
                        file_path,
                        output_file,
                        quality_level
                    )
                elif file_ext in AUDIO_EXTS:
                    success = self.optimize_av_file(
                        file_path,
                        output_file,
                        quality_level,
                        'audio'
                    )
                elif file_ext in VIDEO_EXTS:
                    success = self.optimize_av_file(
                        file_path,
                        output_file,
                        quality_level,
                        'video'
                    )
                elif file_ext in WEB_EXTS:
                    success = self.optimize_web_file(
                        file_path,
                        output_file,
                        file_ext
                    )
                elif file_ext in EPUB_EXTS:
                    success = self.optimize_epub_file(
                        file_path,
                        output_file,
                        compress_images,
                        quality_level
                    )
                else:
                    success = False
                
                if success:
                    compressed_size = os.path.getsize(output_file) if os.path.exists(output_file) else original_size
                    total_compressed_size += compressed_size
                    
                    operation_time = (datetime.now() - operation_start).total_seconds()
                    
                    self.db_manager.add_conversion_record(
                        source_file=file_path,
                        source_format=file_ext.upper().replace('.', ''),
                        target_file=output_file,
                        target_format=file_ext.upper().replace('.', ''),
                        operation_type="office_optimization",
                        file_size=original_size,
                        conversion_time=operation_time,
                        success=True,
                        notes=f"Type: {optimization_type}, Quality: {quality_level}"
                    )
                    
                    success_count += 1
                
                self.progress_bar.setValue(int((i + 1) / len(office_files) * 100))
                
            except Exception as e:
                self.db_manager.add_conversion_record(
                    source_file=file_path,
                    source_format=Path(file_path).suffix.upper().replace('.', ''),
                    target_file="",
                    target_format="",
                    operation_type="office_optimization",
                    file_size=0,
                    conversion_time=0,
                    success=False,
                    notes=f"Error: {str(e)}"
                )
                print(f"Optimization error {file_path}: {e}")
        
        total_time = (datetime.now() - start_time).total_seconds()
        self.show_progress(False)
        
        # Calculate stats
        if total_original_size > 0:
            compression_rate = ((total_original_size - total_compressed_size) / total_original_size * 100)
            savings_mb = (total_original_size - total_compressed_size) / (1024 * 1024)
            
            message = self.translate_text("msg_1").format(success_count, len(office_files), f"{total_time:.1f}")
            message += self.translate_text("msg_2").format(f"{savings_mb:.2f}", f"{compression_rate:.1f}")
            message += self.translate_text("msg_3").format(output_dir)
        else:
            message = self.translate_text("msg_4").format(success_count, len(office_files))
        
        QMessageBox.information(self, self.translate_text("Succès"), self.translate_text(message))

    def optimize_pdf_file(self, pdf_path, output_path, compression_level, remove_metadata, compress_images):
        """Optimize a PDF file"""
        try:
            import fitz
            pdf_document = fitz.open(pdf_path)
            
            # Compression level settings
            if compression_level == "very_reduced":
                save_options = {
                    'garbage': 4,
                    'deflate': True,
                    'clean': True,
                    'deflate_images': compress_images,
                    'deflate_fonts': True,
                    'optimize': True
                }
            elif compression_level == "high":
                save_options = {
                    'garbage': 3,
                    'deflate': True,
                    'clean': True,
                    'deflate_images': compress_images,
                    'deflate_fonts': True
                }
            else:
                save_options = {
                    'garbage': 2,
                    'deflate': True,
                    'clean': True
                }
            
            if remove_metadata:
                pdf_document.set_metadata({})
            
            pdf_document.save(output_path, **save_options)
            pdf_document.close()
            
            return True
        
        except Exception as e:
            print(f"PDF optimization error {pdf_path}: {e}")
            return False

    def optimize_word_file(self, word_path, output_path, compression_level, remove_metadata, compress_images):
        """Optimize a Word file"""
        try:
            from docx import Document
            from docx.oxml.ns import qn
            import io as _io

            doc = Document(word_path)

            # Recompress embedded images
            if compress_images:
                quality_map = {"high": 85, "normal": 75, "very_reduced": 55}
                jpeg_quality = quality_map.get(compression_level, 75)
                try:
                    from PIL import Image as PILImage
                    for rel in doc.part.rels.values():
                        if "image" in rel.reltype:
                            try:
                                blob = rel.target_part.blob
                                ext  = Path(rel.target_part.partname).suffix.lower()
                                if ext in ('.jpg', '.jpeg', '.png', '.bmp', '.tiff'):
                                    buf_in  = _io.BytesIO(blob)
                                    buf_out = _io.BytesIO()
                                    with PILImage.open(buf_in) as img:
                                        img.convert('RGB').save(
                                            buf_out, format='JPEG',
                                            quality=jpeg_quality, optimize=True)
                                    new_blob = buf_out.getvalue()
                                    if len(new_blob) < len(blob):
                                        rel.target_part._blob = new_blob
                            except Exception:
                                pass
                except ImportError:
                    pass

            # Remove metadata
            if remove_metadata:
                doc.core_properties.title    = ""
                doc.core_properties.author   = ""
                doc.core_properties.subject  = ""
                doc.core_properties.keywords = ""
                doc.core_properties.comments = ""
                doc.core_properties.last_modified_by = ""

            doc.save(output_path)
            return True

        except Exception as e:
            print(f"Word optimization error {word_path}: {e}")
            return False

    def optimize_powerpoint_file(self, ppt_path, output_path, compression_level, remove_metadata, compress_images):
        """Optimize a PowerPoint file"""
        try:
            from pptx import Presentation
            from pptx.oxml.ns import qn
            import io

            prs = Presentation(ppt_path)

            slides_to_remove = []
            for i, slide in enumerate(prs.slides):
                if not slide.shapes:
                    slides_to_remove.append(i)

            xml_slides = prs.slides._sldIdLst
            for i in reversed(slides_to_remove):
                sld_id_elem = xml_slides[i]
                xml_slides.remove(sld_id_elem)

            if compress_images:
                quality_map = {"high": 85, "normal": 75, "very_reduced": 55}
                jpeg_quality = quality_map.get(compression_level, 75)
                try:
                    from PIL import Image as PILImage
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if shape.shape_type == 13:
                                try:
                                    img_part = shape.image
                                    blob = img_part.blob
                                    ext  = img_part.ext.lower()
                                    if ext in ("jpg", "jpeg", "png", "bmp", "tiff"):
                                        buf_in  = io.BytesIO(blob)
                                        buf_out = io.BytesIO()
                                        with PILImage.open(buf_in) as img:
                                            rgb = img.convert("RGB")
                                            rgb.save(buf_out, format="JPEG",
                                                     quality=jpeg_quality, optimize=True)
                                        new_blob = buf_out.getvalue()
                                        if len(new_blob) < len(blob):
                                            img_part._blob = new_blob
                                except Exception:
                                    pass
                except ImportError:
                    pass

            # Remove metadata
            if remove_metadata:
                prs.core_properties.title    = ""
                prs.core_properties.author   = ""
                prs.core_properties.subject  = ""
                prs.core_properties.keywords = ""
                prs.core_properties.comments = ""
                prs.core_properties.last_modified_by = ""

            prs.save(output_path)
            return True

        except Exception as e:
            print(f"PowerPoint optimization error {ppt_path}: {e}")
            import traceback; traceback.print_exc()
            return False

    def optimize_av_file(self, src_path, output_path, quality_level, media_type):
        """
        Optimize an audio or video file using ffmpeg.
        quality_level: 0=high quality (less compression), 1=balanced, 2=max compression
        media_type: 'audio' | 'video'
        """
        try:
            import shutil, sys

            # Locate ffmpeg
            ffmpeg_bin = shutil.which("ffmpeg")
            if not ffmpeg_bin:
                candidates = [
                    r"C:\ffmpeg\bin\ffmpeg.exe",
                    r"C:\Program Files\ffmpeg\bin\ffmpeg.exe",
                    r"C:\Program Files (x86)\ffmpeg\bin\ffmpeg.exe",
                    os.path.join(os.environ.get("LOCALAPPDATA", ""), "ffmpeg", "bin", "ffmpeg.exe"),
                    os.path.join(os.environ.get("APPDATA", ""),      "ffmpeg", "bin", "ffmpeg.exe"),
                    os.path.join(getattr(sys, "_MEIPASS", ""),        "ffmpeg.exe"),
                    os.path.join(os.path.dirname(os.path.abspath(__file__)), "ffmpeg.exe"),
                    "/usr/bin/ffmpeg", "/usr/local/bin/ffmpeg",
                    "/opt/homebrew/bin/ffmpeg", "/snap/bin/ffmpeg",
                ]
                for c in candidates:
                    if c and os.path.isfile(c):
                        ffmpeg_bin = c
                        break
            if not ffmpeg_bin:
                print("[optimize_av] can't find ffmpeg — skipping file")
                return False

            ext = Path(output_path).suffix.lower().lstrip(".")

            # Build ffmpeg args based on quality_level and media_type

            if media_type == 'audio':
                AUDIO_PRESETS = {
                    "mp3": {
                        0: ["-codec:a", "libmp3lame", "-q:a", "2",  "-ar", "44100"],
                        1: ["-codec:a", "libmp3lame", "-q:a", "4",  "-ar", "44100"],
                        2: ["-codec:a", "libmp3lame", "-q:a", "7",  "-ar", "44100"],
                    },
                    "aac": {
                        0: ["-codec:a", "aac", "-b:a", "192k", "-ar", "44100"],
                        1: ["-codec:a", "aac", "-b:a", "128k", "-ar", "44100"],
                        2: ["-codec:a", "aac", "-b:a",  "96k", "-ar", "44100"],
                    },
                    "ogg": {
                        0: ["-codec:a", "libvorbis", "-q:a", "6"],
                        1: ["-codec:a", "libvorbis", "-q:a", "4"],
                        2: ["-codec:a", "libvorbis", "-q:a", "2"],
                    },
                    "flac": {
                        0: ["-codec:a", "flac", "-compression_level", "5"],
                        1: ["-codec:a", "flac", "-compression_level", "8"],
                        2: ["-codec:a", "flac", "-compression_level", "12"],
                    },
                    "wav": {
                        0: ["-codec:a", "pcm_s16le", "-ar", "44100"],
                        1: ["-codec:a", "pcm_s16le", "-ar", "44100"],
                        2: ["-codec:a", "pcm_s16le", "-ar", "22050"],
                    },
                }
                default_audio = {
                    0: ["-codec:a", "libmp3lame", "-q:a", "2"],
                    1: ["-codec:a", "libmp3lame", "-q:a", "4"],
                    2: ["-codec:a", "libmp3lame", "-q:a", "7"],
                }
                presets = AUDIO_PRESETS.get(ext, default_audio)
                args = presets.get(quality_level, presets[1])

            else:
                VIDEO_PRESETS = {
                    "mp4": {
                        0: ["-codec:v", "libx264", "-crf", "18", "-preset", "slow",
                            "-codec:a", "aac", "-b:a", "192k", "-movflags", "+faststart"],
                        1: ["-codec:v", "libx264", "-crf", "23", "-preset", "medium",
                            "-codec:a", "aac", "-b:a", "128k", "-movflags", "+faststart"],
                        2: ["-codec:v", "libx264", "-crf", "28", "-preset", "fast",
                            "-codec:a", "aac", "-b:a",  "96k", "-movflags", "+faststart"],
                    },
                    "mkv": {
                        0: ["-codec:v", "libx264", "-crf", "18", "-preset", "slow",
                            "-codec:a", "aac", "-b:a", "192k"],
                        1: ["-codec:v", "libx264", "-crf", "23", "-preset", "medium",
                            "-codec:a", "aac", "-b:a", "128k"],
                        2: ["-codec:v", "libx264", "-crf", "28", "-preset", "fast",
                            "-codec:a", "aac", "-b:a",  "96k"],
                    },
                    "webm": {
                        0: ["-codec:v", "libvpx-vp9", "-crf", "24", "-b:v", "0",
                            "-codec:a", "libopus", "-b:a", "160k"],
                        1: ["-codec:v", "libvpx-vp9", "-crf", "33", "-b:v", "0",
                            "-codec:a", "libopus", "-b:a", "128k"],
                        2: ["-codec:v", "libvpx-vp9", "-crf", "42", "-b:v", "0",
                            "-codec:a", "libopus", "-b:a",  "96k"],
                    },
                    "avi": {
                        0: ["-codec:v", "libxvid", "-qscale:v", "2",
                            "-codec:a", "libmp3lame", "-q:a", "2"],
                        1: ["-codec:v", "libxvid", "-qscale:v", "4",
                            "-codec:a", "libmp3lame", "-q:a", "4"],
                        2: ["-codec:v", "libxvid", "-qscale:v", "8",
                            "-codec:a", "libmp3lame", "-q:a", "7"],
                    },
                }
                default_video = {
                    0: ["-codec:v", "libx264", "-crf", "18", "-preset", "slow",
                        "-codec:a", "aac", "-b:a", "192k"],
                    1: ["-codec:v", "libx264", "-crf", "23", "-preset", "medium",
                        "-codec:a", "aac", "-b:a", "128k"],
                    2: ["-codec:v", "libx264", "-crf", "28", "-preset", "fast",
                        "-codec:a", "aac", "-b:a",  "96k"],
                }
                presets = VIDEO_PRESETS.get(ext, default_video)
                args = presets.get(quality_level, presets[1])

            # Run ffmpeg
            import subprocess
            cmd = [ffmpeg_bin, "-y", "-i", src_path] + args + [output_path]
            result = subprocess.run(cmd, capture_output=True, timeout=3600)

            if result.returncode != 0:
                err = result.stderr.decode(errors="replace")[-500:]
                print(f"[optimize_av] ffmpeg error: {err}")
                return False

            # Accept only if output is smaller than original
            if os.path.exists(output_path):
                orig_size = os.path.getsize(src_path)
                new_size  = os.path.getsize(output_path)
                if new_size >= orig_size * 1.05:
                    import shutil as _sh
                    _sh.copy2(src_path, output_path)
            return True

        except Exception as e:
            if "TimeoutExpired" in type(e).__name__:
                print(f"[optimize_av] Timeout — {src_path}")
            else:
                print(f"[optimize_av] Error {src_path}: {e}")
            return False

    def optimize_web_file(self, src_path, output_path, file_ext):
        """
        Optimize JSON or HTML/HTM files by minification.
        JSON  : re-dump with no indentation/spaces.
        HTML  : strip comments, collapse whitespace between tags.
        """
        try:
            with open(src_path, 'r', encoding='utf-8', errors='replace') as f:
                content = f.read()

            if file_ext == '.json':
                import json as _json
                data = _json.loads(content)
                minified = _json.dumps(data, ensure_ascii=False, separators=(',', ':'))

            else:
                import re as _re
                minified = _re.sub(r'<!--.*?-->', '', content, flags=_re.DOTALL)
                minified = _re.sub(r'>\s+<', '><', minified)
                minified = _re.sub(r'[ 	]{2,}', ' ', minified)
                lines = [l.strip() for l in minified.splitlines()]
                minified = '\n'.join(l for l in lines if l)

            orig_bytes = content.encode('utf-8')
            new_bytes  = minified.encode('utf-8')
            if len(new_bytes) < len(orig_bytes):
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write(minified)
            else:
                import shutil as _sh
                _sh.copy2(src_path, output_path)

            return True

        except Exception as e:
            print(f"Web file optimization error {src_path}: {e}")
            return False

    def optimize_epub_file(self, src_path, output_path, compress_images, quality_level):
        """
        Optimize an EPUB file:
        - Recompress the ZIP with maximum deflate
        - Optionally recompress embedded images with Pillow
        """
        try:
            import zipfile, io as _io

            quality_map = {0: 85, 1: 75, 2: 55}
            img_quality = quality_map.get(quality_level, 75)
            IMAGE_MIMES = {'image/jpeg', 'image/jpg', 'image/png',
                           'image/gif', 'image/webp', 'image/bmp'}
            IMAGE_EXTS  = {'.jpg', '.jpeg', '.png', '.gif', '.webp', '.bmp'}

            buf_out = _io.BytesIO()

            with zipfile.ZipFile(src_path, 'r') as zin, zipfile.ZipFile(buf_out, 'w', compression=zipfile.ZIP_DEFLATED,
                                 compresslevel=9) as zout:

                for item in zin.infolist():
                    data = zin.read(item.filename)
                    ext  = Path(item.filename).suffix.lower()

                    # Recompress embedded images
                    if compress_images and ext in IMAGE_EXTS and len(data) > 2048:
                        try:
                            from PIL import Image as PILImage
                            buf_img = _io.BytesIO(data)
                            buf_new = _io.BytesIO()
                            with PILImage.open(buf_img) as img:
                                if ext in ('.jpg', '.jpeg'):
                                    img.convert('RGB').save(buf_new, format='JPEG',
                                                            quality=img_quality, optimize=True)
                                elif ext == '.png':
                                    compress_lvl = max(1, min(9, int(quality_level * 3) + 1))
                                    img.save(buf_new, format='PNG',
                                             optimize=True, compress_level=compress_lvl)
                                else:
                                    img.convert('RGB').save(buf_new, format='JPEG',
                                                            quality=img_quality, optimize=True)
                            new_data = buf_new.getvalue()
                            if len(new_data) < len(data):
                                data = new_data
                        except Exception:
                            pass

                    # mimetype entry must be stored uncompressed (EPUB spec)
                    if item.filename == 'mimetype':
                        zout.writestr(item, data, compress_type=zipfile.ZIP_STORED)
                    else:
                        zout.writestr(item, data)

            new_bytes  = buf_out.getvalue()
            orig_bytes = os.path.getsize(src_path)

            if len(new_bytes) < orig_bytes:
                with open(output_path, 'wb') as f:
                    f.write(new_bytes)
            else:
                import shutil as _sh
                _sh.copy2(src_path, output_path)

            return True

        except Exception as e:
            print(f"EPUB optimization error {src_path}: {e}")
            return False

    def optimize_excel_file(self, xlsx_path, output_path, compression_level, remove_metadata):
        """Optimize an Excel file — strip metadata and re-save (openpyxl recompresses internally)."""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(xlsx_path)

            if remove_metadata:
                wb.properties.title    = ""
                wb.properties.creator  = ""
                wb.properties.subject  = ""
                wb.properties.keywords = ""
                wb.properties.description = ""
                wb.properties.lastModifiedBy = ""

            # Remove empty sheets
            empty_sheets = [ws.title for ws in wb.worksheets
                            if ws.max_row == 1 and ws.max_column == 1
                            and ws.cell(1, 1).value is None]
            for name in empty_sheets:
                if len(wb.sheetnames) > 1:
                    del wb[name]

            wb.save(output_path)
            return True
        except Exception as e:
            print(f"Excel optimization error {xlsx_path}: {e}")
            return False

    def optimize_image_file(self, img_path, output_path, quality_level):
        """Optimize an image file using Pillow (recompress as JPEG/PNG)."""
        try:
            from PIL import Image as PILImage
            quality_map = {0: 85, 1: 75, 2: 55}
            quality = quality_map.get(quality_level, 75)

            ext_map = {
                ".jpg": "JPEG", ".jpeg": "JPEG",
                ".png": "PNG",  ".bmp":  "PNG",
                ".tiff": "TIFF", ".webp": "WEBP",
                ".gif": "GIF",
            }
            src_ext  = Path(img_path).suffix.lower()
            dst_ext  = Path(output_path).suffix.lower()
            fmt_out  = ext_map.get(dst_ext, ext_map.get(src_ext, "JPEG"))

            with PILImage.open(img_path) as img:
                exif_data = None
                try:
                    exif_data = img.info.get("exif")
                except Exception:
                    pass

                if fmt_out == "JPEG":
                    rgb = img.convert("RGB")
                    save_kwargs = {"quality": quality, "optimize": True}
                    if exif_data:
                        save_kwargs["exif"] = exif_data
                    rgb.save(output_path, format="JPEG", **save_kwargs)
                elif fmt_out == "PNG":
                    compress = max(1, min(9, int(quality_level * 3)))
                    img.save(output_path, format="PNG", optimize=True, compress_level=compress)
                elif fmt_out == "WEBP":
                    img.save(output_path, format="WEBP", quality=quality, method=6)
                else:
                    img.save(output_path, format=fmt_out)

            return True
        except Exception as e:
            print(f"Image optimization error {img_path}: {e}")
            return False

    def clear_file_list(self):
        """Clear the file list (alias for clear_files)"""
        self.clear_files()

    def check_pdf_has_images(self, pdf_path):
        try:
            with open(pdf_path, 'rb') as f:
                pdf_reader = PdfReader(f)
                
                if pdf_reader.is_encrypted:
                    # Do not check images for encrypted PDF
                    # Let conversion handle it after decryption
                    return False
                
                import fitz
                pdf_document = fitz.open(pdf_path)
                has_images = False
                
                for page_num in range(len(pdf_document)):
                    page = pdf_document.load_page(page_num)
                    image_list = page.get_images()
                    
                    if image_list:
                        has_images = True
                        break
                
                pdf_document.close()
                return has_images
        except Exception:
            return False

    def convert_pdf_to_word(self):
        if not (hasattr(self, 'active_templates') and 'pdf_to_word' in self.active_templates):
            _def_id, _ = (self._ensure_template_manager() or object()).get_default_template('Conversion PDF→Word')
            if _def_id:
                (self._ensure_template_manager() or object()).apply_template(_def_id, self)

        if hasattr(self, 'active_templates') and 'word_to_pdf' in self.active_templates:
            template = self.active_templates['word_to_pdf']
            page_format = template.get('page_format_value', 'a4')
            orientation = template.get('orientation_value', 'portrait')
            quality = template.get('quality_value', 150)
            include_metadata = template.get('include_metadata', True)
            compress = template.get('compress', True)
        
        selected_items = self.files_list_widget.selectedItems()
        files_to_process = []
        if selected_items:
            for i in range(self.files_list_widget.count()):
                item = self.files_list_widget.item(i)
                if item.isSelected():
                    files_to_process.append(item.data(Qt.UserRole))
        else:
            files_to_process = [f for f in self.files_list if f.lower().endswith('.pdf')]

        # Filter only PDF files
        pdf_files = [f for f in files_to_process if f.lower().endswith('.pdf')]
        if not pdf_files:
            msg = self.translate_text("Veuillez sélectionner au moins un fichier PDF") if selected_items else self.translate_text("La liste doit contenir au moins un fichier PDF")
            QMessageBox.warning(self, self.translate_text("Avertissement"), msg)
            return
        
        encrypted_pdfs = []
        for pdf_file in pdf_files:
            try:
                with open(pdf_file, 'rb') as f:
                    pdf_reader = PdfReader(f)
                    if pdf_reader.is_encrypted:
                        encrypted_pdfs.append(pdf_file)
            except Exception as e:
                print(f"PDF verification error {pdf_file}: {e}")
        
        if encrypted_pdfs:
            response = self.handle_encrypted_pdfs(encrypted_pdfs, pdf_files)
            if response is False:
                return
            elif response is not None:
                passwords = response
                pdf_files = self.decrypt_pdfs_before_conversion(pdf_files, passwords)
                print(f"[DEBUG] {len(encrypted_pdfs)} password-protected PDF files successfully converted.")
                self.achievement_system.record_protected_file_conversion(len(encrypted_pdfs), "pdf")
                if not pdf_files:
                    return
        
        pdfs_with_images = []
        for pdf_file in pdf_files:
            if self.check_pdf_has_images(pdf_file):
                pdfs_with_images.append(pdf_file)
        
        if not pdfs_with_images:
            self.convert_pdfs_without_images(pdf_files)
            return
        
        if hasattr(self, 'active_templates') and 'pdf_to_word' in self.active_templates:
            conversion_mode = self.active_templates['pdf_to_word'].get('mode', 'with_images')
        else:
            current_mode = self.config.get("pdf_to_word_mode", "with_images")
            dialog = PdfToWordDialog(self, self.current_language, current_mode, has_images=True)
            if dialog.exec() != QDialog.Accepted:
                return
            conversion_mode = dialog.get_conversion_mode()
        
        output_dir = self.get_output_directory()
        if not output_dir:
            return
        
        message = self.translate_text("conversion_pdf_to_word").format(len(pdf_files))
        self.show_progress(True, message)
        self._set_ui_enabled(False)

        _mode   = conversion_mode
        _outdir = output_dir

        def _run_pdf_to_word(task):
            import os, time as _time
            from pathlib import Path as _Path
            from datetime import datetime as _dt
            t0 = _time.perf_counter()
            fp = task["input_path"]
            out = task["output_path"]
            fs  = os.path.getsize(fp) if os.path.exists(fp) else 0
            if _mode == "with_images":
                self.convert_pdf_to_docx_improved(fp, out)
            elif _mode == "text_only":
                self.convert_pdf_to_docx_text_only(fp, out)
            else:
                self.convert_pdf_to_docx_with_image_text(fp, out)
            return {"success": True, "error": "",
                    "file_size": fs, "operation_time": _time.perf_counter() - t0}

        tasks = [
            {"index": i, "total": len(pdf_files),
             "input_path": fp,
             "output_path": os.path.join(_outdir, f"{Path(fp).stem}.docx")}
            for i, fp in enumerate(pdf_files)
        ]

        def _on_file_done(result):
            if result.get("success"):
                fp  = result["input_path"]
                out = result["output_path"]
                fs  = result.get("file_size", 0)
                ot  = result.get("operation_time", 0)
                self.achievement_system.record_conversion("pdf_to_word", fs, True)
                self.achievement_system.mark_format_as_used("pdf")
                self.achievement_system.mark_format_as_used("docx")
                self.db_manager.add_conversion_record(
                    source_file=fp, source_format="PDF",
                    target_file=out, target_format="DOCX",
                    operation_type="pdf_to_word", file_size=fs,
                    conversion_time=ot, success=True,
                    notes=f"Mode: {_mode}"
                )
            else:
                fp  = result["input_path"]
                self.db_manager.add_conversion_record(
                    source_file=fp, source_format="PDF",
                    target_file="", target_format="DOCX",
                    operation_type="pdf_to_word", file_size=0,
                    conversion_time=0, success=False,
                    notes=f"Error: {result.get('error','')}"
                )

        def _on_finished(summary):
            self.show_progress(False)
            self._set_ui_enabled(True)
            sc = summary["success_count"]
            total = summary["total"]
            total_time = summary["total_time"]
            if sc > 0 and 0 <= datetime.now().hour < 6:
                self.achievement_system.increment_stat("night_conversions", sc)
                self.achievement_system.check_achievement("night_owl")
            formatted_time = f"{total_time:.1f}"
            msg = self.translate_text("pdf_to_word_success_sum").format(sc, total, formatted_time)
            if self.config.get("enable_system_notifications", True):
                self.system_notifier.send("pdf_to_word")
            if sc >= 50 and total_time <= 300:
                self.achievement_system.update_stat("recent_batch_files", sc)
                self.achievement_system.update_stat("recent_batch_time", total_time)
                self.achievement_system.check_speed_conversion(sc, total_time)
            QMessageBox.information(self, self.translate_text("Succès"), msg)

        self._worker = ConversionWorker(tasks, _run_pdf_to_word)
        self._worker.progress.connect(self.progress_bar.setValue)
        self._worker.file_done.connect(_on_file_done)
        self._worker.finished.connect(_on_finished)
        self._worker.start()

    def handle_encrypted_pdfs(self, encrypted_pdfs, all_pdfs):
        """Handle encrypted PDFs before conversion"""
        dialog = QDialog(self)
        dialog.setWindowTitle(self.translate_text("PDFs protégés par mot de passe"))
        dialog.setMinimumWidth(500)
        
        layout = QVBoxLayout(dialog)
        
        # Information message
        if len(encrypted_pdfs) == 1:
            message = self.translate_text(f"Le fichier '{Path(encrypted_pdfs[0]).name}' est protégé par mot de passe.")
        else:
            message = self.translate_text(f"{len(encrypted_pdfs)} fichiers PDF sont protégés par mot de passe.")
        
        message += "\n\n" + self.translate_text("Pour convertir ces fichiers en Word, vous devez fournir les mots de passe.\n\n")
        message += self.translate_text("Options :")
        
        info_label = QLabel(message)
        info_label.setWordWrap(True)
        info_label.setStyleSheet("font-weight: bold; color: #d35400; padding: 10px;")
        layout.addWidget(info_label)
        
        tab_widget = QTabWidget()
        
        single_tab = QWidget()
        single_layout = QVBoxLayout(single_tab)
        
        if len(encrypted_pdfs) == 1:
            single_layout.addWidget(QLabel(self.translate_text(f"Fichier : {Path(encrypted_pdfs[0]).name}")))
            single_password_input = QLineEdit()
            single_password_input.setEchoMode(QLineEdit.Password)
            single_password_input.setPlaceholderText(self.translate_text("Mot de passe du PDF"))
            
            single_layout.addWidget(QLabel(self.translate_text("Mot de passe :")))
            single_layout.addWidget(single_password_input)
            
            dialog.single_password_input = single_password_input
            dialog.single_pdf = encrypted_pdfs[0]
        
        multiple_tab = QWidget()
        multiple_layout = QVBoxLayout(multiple_tab)
        
        if len(encrypted_pdfs) > 1:
            scroll_area = QScrollArea()
            scroll_widget = QWidget()
            scroll_layout = QVBoxLayout(scroll_widget)
            
            password_inputs = {}
            
            for pdf_file in encrypted_pdfs:
                group = QGroupBox(Path(pdf_file).name)
                group_layout = QVBoxLayout(group)
                
                password_input = QLineEdit()
                password_input.setEchoMode(QLineEdit.Password)
                password_input.setPlaceholderText(self.translate_text("Mot de passe (laisser vide si aucun)"))
                
                group_layout.addWidget(QLabel(self.translate_text("Mot de passe :")))
                group_layout.addWidget(password_input)
                
                password_inputs[pdf_file] = password_input
                scroll_layout.addWidget(group)
            
            scroll_widget.setLayout(scroll_layout)
            scroll_area.setWidget(scroll_widget)
            scroll_area.setWidgetResizable(True)
            multiple_layout.addWidget(scroll_area)
            
            dialog.password_inputs = password_inputs
        
        same_tab = QWidget()
        same_layout = QVBoxLayout(same_tab)
        
        same_password_input = QLineEdit()
        same_password_input.setEchoMode(QLineEdit.Password)
        same_password_input.setPlaceholderText(self.translate_text("Mot de passe commun à tous les PDF"))
        
        same_layout.addWidget(QLabel(self.translate_text("Ce mot de passe sera essayé sur tous les PDF chiffrés.")))
        same_layout.addWidget(QLabel(self.translate_text("Mot de passe commun :")))
        same_layout.addWidget(same_password_input)
        
        dialog.same_password_input = same_password_input
        
        if len(encrypted_pdfs) == 1:
            tab_widget.addTab(single_tab, self.translate_text("PDF unique"))
        else:
            tab_widget.addTab(multiple_tab, self.translate_text("PDF multiples"))
            tab_widget.addTab(same_tab, self.translate_text("Mot de passe commun"))
        
        layout.addWidget(tab_widget)
        
        options_group = QGroupBox(self.translate_text("Options de sortie"))
        options_layout = QVBoxLayout(options_group)
        
        keep_encrypted_check = AnimatedCheckBox(
            self.translate_text("Conserver les versions originales chiffrées après conversion")
        )
        keep_encrypted_check.setChecked(True)
        
        remove_passwords_check = AnimatedCheckBox(
            self.translate_text("Supprimer la protection par mot de passe des PDF convertis (recommandé)")
        )
        remove_passwords_check.setChecked(True)
        
        options_layout.addWidget(keep_encrypted_check)
        options_layout.addWidget(remove_passwords_check)
        
        layout.addWidget(options_group)
        
        button_box = QDialogButtonBox(QDialogButtonBox.Ok | QDialogButtonBox.Cancel)

        def collect_passwords():
            passwords = {}
            
            if len(encrypted_pdfs) == 1:
                # One encrypted PDF
                password = dialog.single_password_input.text()
                if password:
                    passwords[dialog.single_pdf] = password
            else:
                if tab_widget.currentIndex() == 0:
                    for pdf_file, password_input in dialog.password_inputs.items():
                        password = password_input.text()
                        if password:
                            passwords[pdf_file] = password
                else:
                    common_password = dialog.same_password_input.text()
                    if common_password:
                        for pdf_file in encrypted_pdfs:
                            passwords[pdf_file] = common_password
            
            dialog.passwords = passwords
            dialog.keep_originals = keep_encrypted_check.isChecked()
            dialog.remove_passwords = remove_passwords_check.isChecked()
            dialog.accept()
        
        button_box.accepted.connect(collect_passwords)
        button_box.rejected.connect(dialog.reject)
        
        layout.addWidget(button_box)
        
        if dialog.exec() == QDialog.Accepted:
            if hasattr(dialog, 'passwords') and dialog.passwords:
                return {
                    'passwords': dialog.passwords,
                    'keep_originals': dialog.keep_originals,
                    'remove_passwords': dialog.remove_passwords
                }
            else:
                QMessageBox.warning(self, self.translate_text("Avertissement"),
                                self.translate_text("Aucun mot de passe fourni. La conversion sera annulée."))
                return False
        else:
            return False

    def decrypt_pdfs_before_conversion(self, pdf_files, passwords_info):
        """Decrypt PDFs before conversion and return new paths"""
        temp_files = []
        success_count = 0
        failed_files = []
        
        passwords = passwords_info.get('passwords', {})
        keep_originals = passwords_info.get('keep_originals', True)
        remove_passwords = passwords_info.get('remove_passwords', True)
        
        self.show_progress(True, self.translate_text(f"Déchiffrement des PDF..."))
        
        for i, pdf_file in enumerate(pdf_files):
            try:
                # Check if the PDF is encrypted
                with open(pdf_file, 'rb') as f:
                    pdf_reader = PdfReader(f)
                    
                    if pdf_reader.is_encrypted:
                        password = passwords.get(pdf_file)
                        
                        if not password:
                            try:
                                pdf_reader.decrypt('')
                            except:
                                failed_files.append(f"{Path(pdf_file).name} - {self.translate_text('Mot de passe requis')}")
                                continue
                        else:
                            success = pdf_reader.decrypt(password)
                            
                            if not success:
                                for pwd_variation in [password, password.lower(), password.upper()]:
                                    try:
                                        success = pdf_reader.decrypt(pwd_variation)
                                        if success:
                                            break
                                    except:
                                        continue
                                
                                if not success:
                                    failed_files.append(f"{Path(pdf_file).name} - {self.translate_text('Mot de passe incorrect')}")
                                    continue
                        
                        pdf_writer = PdfWriter()
                        
                        for page in pdf_reader.pages:
                            pdf_writer.add_page(page)
                        
                        original_stem = Path(pdf_file).stem
                        temp_file = os.path.join(
                            tempfile.gettempdir(),
                            f"{original_stem}_decrypted.pdf"
                        )
                        
                        counter = 1
                        while os.path.exists(temp_file):
                            temp_file = os.path.join(
                                tempfile.gettempdir(),
                                f"{original_stem}_decrypted_{counter}.pdf"
                            )
                            counter += 1
                        self.temp_files.append(temp_file)
                        
                        with open(temp_file, 'wb') as output_file:
                            if remove_passwords:
                                pdf_writer.write(output_file)
                            else:
                                pdf_writer.encrypt(password)
                                pdf_writer.write(output_file)
                        
                        temp_files.append(temp_file)
                        success_count += 1
                        
                        if not keep_originals:
                            backup_file = pdf_file + ".backup"
                            shutil.copy2(pdf_file, backup_file)
                            shutil.copy2(temp_file, pdf_file)
                            temp_files.append(backup_file)
                    else:
                        temp_files.append(pdf_file)
                        success_count += 1
            
            except Exception as e:
                error_msg = f"{Path(pdf_file).name} - {str(e)}"
                failed_files.append(error_msg)
                print(f"Decryption error {pdf_file}: {e}")
            
            self.progress_bar.setValue(int((i + 1) / len(pdf_files) * 100))
        
        self.show_progress(False)
        
        if failed_files:
            error_message = self.translate_text(f"{success_count}/{len(pdf_files)} PDF(s) déchiffré(s) avec succès.\n\n")
            error_message += self.translate_text(f"Échecs ({len(failed_files)}):\n")
            error_message += "\n".join(failed_files[:3])  # Display only the first three
            
            if len(failed_files) > 3:
                error_message += f"\n... et {len(failed_files) - 3} autres"
            
            QMessageBox.warning(self, self.translate_text("Résultat déchiffrement"), error_message)
        
        return temp_files

    def decrypt_single_pdf(self, pdf_path, password, output_path=None, remove_password=True):
        """Decrypt a single PDF"""
        try:
            with open(pdf_path, 'rb') as f:
                pdf_reader = PdfReader(f)
                
                if not pdf_reader.is_encrypted:
                    return pdf_path
                
                success = pdf_reader.decrypt(password)
                
                if not success:
                    for pwd_variation in [password, password.lower(), password.upper()]:
                        try:
                            success = pdf_reader.decrypt(pwd_variation)
                            if success:
                                break
                        except:
                            continue
                    
                    if not success:
                        raise Exception(self.translate_text("Mot de passe incorrect"))
                
                pdf_writer = PdfWriter()
                
                for page in pdf_reader.pages:
                    pdf_writer.add_page(page)
                
                if not output_path:
                    # Create temp file
                    output_path = self.create_temp_file(suffix="_decrypted.pdf")
                
                with open(output_path, 'wb') as output_file:
                    if remove_password:
                        pdf_writer.write(output_file)
                    else:
                        pdf_writer.encrypt(password)
                        pdf_writer.write(output_file)
                
                return output_path
        
        except Exception as e:
            print(f"Decryption error {pdf_path}: {e}")
            raise

    def convert_pdfs_without_images(self, pdf_files):
        output_dir = self.get_output_directory()
        if not output_dir:
            return
        
        message = self.translate_text("conversion_pdf_to_word").format(len(pdf_files))
        self.show_progress(True, message)
        
        success_count = 0
        start_time = datetime.now()
        
        for i, file_path in enumerate(pdf_files):
            try:
                output_file = os.path.join(output_dir, f"{Path(file_path).stem}.docx")
                file_size = os.path.getsize(file_path)
                
                operation_start = datetime.now()
                
                current_mode = self.config.get("pdf_to_word_mode", "with_images")
                if current_mode == "with_images":
                    self.convert_pdf_to_docx_improved(file_path, output_file)
                elif current_mode == "text_only":
                    self.convert_pdf_to_docx_text_only(file_path, output_file)
                else:
                    self.convert_pdf_to_docx_text_only(file_path, output_file)
                
                operation_time = (datetime.now() - operation_start).total_seconds()
                self.achievement_system.record_conversion("pdf_to_word", file_size, True)
                self.achievement_system.mark_format_as_used("pdf")
                self.achievement_system.mark_format_as_used("docx")
                
                self.db_manager.add_conversion_record(
                    source_file=file_path,
                    source_format="PDF",
                    target_file=output_file,
                    target_format="DOCX",
                    operation_type="pdf_to_word",
                    file_size=file_size,
                    conversion_time=operation_time,
                    success=True,
                    notes=f"Mode: {current_mode}"
                )
                
                success_count += 1
                self.progress_bar.setValue(int((i + 1) / len(pdf_files) * 100))
            
            except Exception as e:
                self.db_manager.add_conversion_record(
                    source_file=file_path,
                    source_format="PDF",
                    target_file="",
                    target_format="DOCX",
                    operation_type="pdf_to_word",
                    file_size=0,
                    conversion_time=0,
                    success=False,
                    notes=f"Error: {str(e)}"
                )
                QMessageBox.critical(self, self.translate_text("Erreur"), self.translate_text(f"Erreur avec {Path(file_path).name}: {str(e)}"))
        
        total_time = (datetime.now() - start_time).total_seconds()
        self.show_progress(False)

        if success_count > 0 and 0 <= datetime.now().hour < 6:
            self.achievement_system.increment_stat("night_conversions", success_count)
            self.achievement_system.check_achievement("night_owl")
        
        if success_count >= 50 and total_time <= 300:
            self.achievement_system.update_stat("recent_batch_files", success_count)
            self.achievement_system.update_stat("recent_batch_time", total_time)
            self.achievement_system.check_speed_conversion(success_count, total_time)
        
        formatted_time = f"{total_time:.1f}"
        message = self.translate_text("pdf_to_word_success_sum").format(success_count,len(pdf_files),formatted_time)
        if self.config.get("enable_system_notifications", True):
            self.system_notifier.send("pdf_to_word")  
        QMessageBox.information(self, self.translate_text("Succès"), self.translate_text(message))


    def _get_chart_regions(self, page):
        """
        Return a list of fitz.Rect for regions that look like charts/graphs.

        Detection heuristic:
          - fitz.Page.get_drawings() returns every vector path on the page.
          - A cluster of paths that spans a bounding box >= MIN_AREA and
            contains >= MIN_PATHS elements is considered a chart region.
          - We inflate each bbox slightly and clip it to the page rect.
          - Overlapping bboxes are merged so a multi-series chart becomes
            one region rather than dozens of tiny fragments.
        """
        MIN_PATHS   = 8
        MIN_AREA    = 8000
        PADDING     = 10
        MERGE_GAP   = 20

        import fitz
        drawings = page.get_drawings()
        if not drawings:
            return []

        path_rects = [fitz.Rect(d["rect"]) for d in drawings if d.get("rect")]
        if len(path_rects) < MIN_PATHS:
            return []

        path_rects.sort(key=lambda r: (r.y0, r.x0))
        clusters = []
        cur = fitz.Rect(path_rects[0])
        count = 1

        for r in path_rects[1:]:
            expanded = fitz.Rect(cur.x0 - MERGE_GAP, cur.y0 - MERGE_GAP,
                                 cur.x1 + MERGE_GAP, cur.y1 + MERGE_GAP)
            if expanded.intersects(r):
                cur = fitz.Rect(
                    min(cur.x0, r.x0), min(cur.y0, r.y0),
                    max(cur.x1, r.x1), max(cur.y1, r.y1)
                )
                count += 1
            else:
                if count >= MIN_PATHS and cur.get_area() >= MIN_AREA:
                    clusters.append(cur)
                cur = fitz.Rect(r)
                count = 1

        if count >= MIN_PATHS and cur.get_area() >= MIN_AREA:
            clusters.append(cur)

        # Pad and clip to page
        page_rect = page.rect
        result = []
        for c in clusters:
            padded = fitz.Rect(
                max(0,            c.x0 - PADDING),
                max(0,            c.y0 - PADDING),
                min(page_rect.x1, c.x1 + PADDING),
                min(page_rect.y1, c.y1 + PADDING),
            )
            result.append(padded)

        return result

    def _rasterize_region(self, page, rect, dpi=150):
        """
        Rasterize a rectangular region of a PDF page to a PNG bytes object.
        dpi=150 gives a good balance between quality and file size.
        """
        import fitz
        zoom   = dpi / 72.0
        matrix = fitz.Matrix(zoom, zoom)
        clip   = fitz.IRect(
            int(rect.x0 * zoom), int(rect.y0 * zoom),
            int(rect.x1 * zoom), int(rect.y1 * zoom),
        )
        pix = page.get_pixmap(matrix=matrix, clip=rect, alpha=False)
        return pix.tobytes("png")

    def _try_word_com_pdf_to_docx(self, pdf_path, docx_path):
        """
        Use Microsoft Word COM to open a PDF natively and save as DOCX.
        Word 2013+ shows a hidden "converting PDF..." confirmation dialog
        even with DisplayAlerts=0. A background thread finds that dialog
        and sends Enter every 500 ms to dismiss it automatically.
        Hard timeout of 60 s prevents infinite blocking.
        Returns True on success, False otherwise.
        Only attempted on Windows when Word is installed.
        """
        import sys
        if sys.platform != "win32":
            return False

        # Fast check: is Word actually installed?
        try:
            import winreg
            winreg.OpenKey(winreg.HKEY_CLASSES_ROOT, "Word.Application")
        except (OSError, ImportError):
            print("[PDF→DOCX] Word not found in registry — skipping COM tier")
            return False

        import os
        import threading
        pdf_path  = os.path.abspath(pdf_path)
        docx_path = os.path.abspath(docx_path)

        result      = {"ok": False}
        stop_event  = threading.Event()

        def _dialog_dismisser():
            """
            Polls every 300 ms for the Word PDF-conversion confirmation dialog
            ("Microsoft Word" title, #32770 class) and clicks its OK button.
            Uses EnumChildWindows to find the actual OK button handle so the
            click is reliable regardless of button position.
            """
            import ctypes
            import ctypes.wintypes
            user32   = ctypes.windll.user32
            BM_CLICK = 0x00F5

            DialogProc = ctypes.WINFUNCTYPE(
                ctypes.c_bool,
                ctypes.wintypes.HWND,
                ctypes.wintypes.LPARAM,
            )

            def _find_ok_button(hwnd_dialog):
                """Return the HWND of the OK button inside hwnd_dialog, or None."""
                found = ctypes.wintypes.HWND(0)

                @DialogProc
                def _enum(hwnd_child, _lp):
                    buf = ctypes.create_unicode_buffer(64)
                    user32.GetWindowTextW(hwnd_child, buf, 64)
                    # Match "OK" in any language variant (Ok, ok, OK)
                    if buf.value.strip().upper() in ("OK", "O&K"):
                        found.value = hwnd_child
                        return False
                    # Also match button with ID=1 (IDOK)
                    ctrl_id = user32.GetDlgCtrlID(hwnd_child)
                    if ctrl_id == 1:
                        found.value = hwnd_child
                        return False
                    return True

                user32.EnumChildWindows(hwnd_dialog, _enum, 0)
                return found.value or None

            dialog_classes = ["#32770", "bosa_sdm_msword"]
            target_title   = "Microsoft Word"

            while not stop_event.is_set():
                # Search by class name
                for cls in dialog_classes:
                    hwnd = user32.FindWindowW(cls, None)
                    if hwnd:
                        ok_btn = _find_ok_button(hwnd)
                        if ok_btn:
                            user32.SendMessageW(ok_btn, BM_CLICK, 0, 0)
                        else:
                            # Fallback: WM_COMMAND IDOK on the dialog itself
                            user32.SendMessageW(hwnd, 0x0111, 1, 0)
                # Also search by title
                hwnd = user32.FindWindowW(None, target_title)
                if hwnd:
                    ok_btn = _find_ok_button(hwnd)
                    if ok_btn:
                        user32.SendMessageW(ok_btn, BM_CLICK, 0, 0)
                    else:
                        user32.SendMessageW(hwnd, 0x0111, 1, 0)
                import time as _t
                _t.sleep(0.3)

        def _worker():
            word = None
            doc  = None
            try:
                import pythoncom
                import comtypes.client
                pythoncom.CoInitialize()
                word = comtypes.client.CreateObject('Word.Application')
                word.Visible        = False
                word.DisplayAlerts  = 0
                word.AutomationSecurity = 3

                doc = word.Documents.Open(
                    pdf_path,
                    ConfirmConversions = False,
                    ReadOnly           = True,
                    AddToRecentFiles   = False,
                    NoEncodingDialog   = True,
                )
                doc.SaveAs2(docx_path, FileFormat=16)
                result["ok"] = True
                print(f"[PDF→DOCX] Word COM success: {os.path.basename(pdf_path)}")
            except Exception as e:
                print(f"[PDF→DOCX] Word COM failed: {e}")
            finally:
                if doc is not None:
                    try: doc.Close(False)
                    except: pass
                if word is not None:
                    try: word.Quit()
                    except: pass
                try:
                    import pythoncom
                    pythoncom.CoUninitialize()
                except: pass

        dismisser = threading.Thread(target=_dialog_dismisser, daemon=True)
        dismisser.start()

        t = threading.Thread(target=_worker, daemon=True)
        t.start()
        t.join(timeout=60)

        stop_event.set()

        if t.is_alive():
            print("[PDF→DOCX] Word COM timed out — falling back to pdf2docx")
            return False

        return result["ok"]

    def _try_pdf2docx(self, pdf_path, docx_path):
        """
        Convert PDF → DOCX using pdf2docx.
        Returns True on success, False otherwise.
        """
        if Converter is None:
            print("[PDF→DOCX] pdf2docx not available — skipping Tier 2")
            return False
        try:
            cv = Converter(pdf_path)
            cv.convert(docx_path, start=0, end=None, parse_drawing=True)
            cv.close()
            print(f"[PDF→DOCX] pdf2docx success: {Path(pdf_path).name}")
            return True
        except Exception as e:
            print(f"[PDF→DOCX] pdf2docx failed: {e}")
            return False

    def _solid_fitz_fallback(self, pdf_path, docx_path):
        """
        High-quality fallback using fitz get_text("dict") to reconstruct:
          - Text blocks with approximate font size → heading/body heuristic
          - Bold/italic detection via font flags
          - Tables detected by fitz find_tables() (PyMuPDF 1.23+)
          - Images extracted and placed at their real vertical position
          - Multi-column layout handled by sorting blocks top-to-bottom
          - Page breaks between pages
        """
        import io as _io
        import fitz
        from docx.shared import Pt, Inches
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        doc = Document()

        from docx.shared import Cm
        for section in doc.sections:
            section.top_margin    = Cm(2.0)
            section.bottom_margin = Cm(2.0)
            section.left_margin   = Cm(2.5)
            section.right_margin  = Cm(2.5)

        pdf_doc = fitz.open(pdf_path)

        def _flag_to_style(flags, size, median_size):
            """Heuristic: map font flags + size to a paragraph style label."""
            is_bold   = bool(flags & 16)
            is_italic = bool(flags & 2)
            ratio     = size / median_size if median_size else 1.0
            if ratio >= 1.6 and is_bold:
                return "heading1"
            if ratio >= 1.3 and is_bold:
                return "heading2"
            if ratio >= 1.1 and is_bold:
                return "heading3"
            return ("bold_body" if is_bold else
                    "italic_body" if is_italic else "body")

        def _apply_run_style(run, flags, size):
            run.bold   = bool(flags & 16)
            run.italic = bool(flags & 2)
            if size:
                run.font.size = Pt(round(size, 1))

        def _add_table_to_doc(doc, tab):
            """
            Render a fitz TableFinder table into a python-docx Table.
            tab.extract() returns list-of-lists of cell text.
            """
            try:
                rows_data = tab.extract()
                if not rows_data or not rows_data[0]:
                    return
                n_rows = len(rows_data)
                n_cols = max(len(r) for r in rows_data)
                tbl = doc.add_table(rows=n_rows, cols=n_cols)
                tbl.style = "Table Grid"
                for r_idx, row_data in enumerate(rows_data):
                    for c_idx, cell_text in enumerate(row_data):
                        if c_idx < n_cols:
                            cell = tbl.cell(r_idx, c_idx)
                            cell.text = (cell_text or "").strip()
                doc.add_paragraph()
            except Exception as te:
                print(f"[fallback table] {te}")

        for page_num in range(len(pdf_doc)):
            if page_num > 0:
                doc.add_page_break()

            page     = pdf_doc.load_page(page_num)
            page_h   = page.rect.height
            page_w   = page.rect.width
            text_dict = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)

            # Collect all font sizes to compute median
            all_sizes = []
            for blk in text_dict.get("blocks", []):
                if blk["type"] != 0:
                    continue
                for line in blk.get("lines", []):
                    for span in line.get("spans", []):
                        if span.get("size", 0) > 0:
                            all_sizes.append(span["size"])
            all_sizes.sort()
            median_size = all_sizes[len(all_sizes) // 2] if all_sizes else 11.0

            # Detect tables
            table_rects = []
            tables_on_page = []
            try:
                finder = page.find_tables()
                tables_on_page = finder.tables if finder else []
                for t in tables_on_page:
                    table_rects.append(fitz.Rect(t.bbox))
            except Exception:
                pass 

            def _rect_in_table(rect):
                """True if this rect overlaps significantly with a known table."""
                for tr in table_rects:
                    inter = fitz.Rect(rect).intersect(tr)
                    if not inter.is_empty and inter.get_area() > 100:
                        return True
                return False

            # Collect images with their y-position
            images_on_page = []
            for img in page.get_images(full=True):
                xref = img[0]
                try:
                    bbox_list = page.get_image_rects(xref)
                    if not bbox_list:
                        continue
                    bbox = bbox_list[0]
                    if _rect_in_table(bbox):
                        continue
                    pix = fitz.Pixmap(pdf_doc, xref)
                    if pix.colorspace and pix.colorspace.n > 3:
                        pix = fitz.Pixmap(fitz.csRGB, pix)
                    if pix.width < 10 or pix.height < 10:
                        continue   # skip tiny decorative pixels
                    img_bytes = pix.tobytes("png")
                    images_on_page.append({
                        "y0": bbox.y0,
                        "y1": bbox.y1,
                        "bytes": img_bytes,
                        "width_pt": bbox.width,
                    })
                    pix = None
                except Exception as ie:
                    print(f"[fallback img xref={xref}] {ie}")

            images_on_page.sort(key=lambda x: x["y0"])

            stream = []
            for blk in text_dict.get("blocks", []):
                y0 = blk.get("bbox", [0, 0])[1]
                if blk["type"] == 0:
                    if not _rect_in_table(fitz.Rect(blk["bbox"])):
                        stream.append(("text", y0, blk))
            for img in images_on_page:
                stream.append(("image", img["y0"], img))
            for tab in tables_on_page:
                try:
                    stream.append(("table", tab.bbox[1], tab))
                except Exception:
                    pass

            stream.sort(key=lambda x: x[1])

            img_inserted_xranges = set()

            for item_type, y0, obj in stream:

                if item_type == "image":
                    key = round(y0 / 5) * 5
                    if key in img_inserted_xranges:
                        continue
                    img_inserted_xranges.add(key)
                    try:
                        buf = _io.BytesIO(obj["bytes"])
                        w_pt  = obj["width_pt"]
                        max_w = Inches(5.5)
                        width = min(Pt(w_pt), max_w)
                        p = doc.add_paragraph()
                        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                        run = p.add_run()
                        run.add_picture(buf, width=width)
                    except Exception as ie:
                        print(f"[fallback img render] {ie}")

                elif item_type == "table":
                    _add_table_to_doc(doc, obj)

                elif item_type == "text":
                    blk = obj
                    for line in blk.get("lines", []):
                        spans = line.get("spans", [])
                        if not spans:
                            continue

                        line_text = "".join(s.get("text", "") for s in spans)
                        if not line_text.strip():
                            continue
                        dom = max(spans, key=lambda s: len(s.get("text", "")))
                        flags = dom.get("flags", 0)
                        size  = dom.get("size", median_size)
                        style_label = _flag_to_style(flags, size, median_size)

                        if style_label == "heading1":
                            p = doc.add_heading(line_text.strip(), level=1)
                        elif style_label == "heading2":
                            p = doc.add_heading(line_text.strip(), level=2)
                        elif style_label == "heading3":
                            p = doc.add_heading(line_text.strip(), level=3)
                        else:
                            p = doc.add_paragraph()
                            for span in spans:
                                span_text = span.get("text", "")
                                if not span_text:
                                    continue
                                run = p.add_run(span_text)
                                _apply_run_style(run, span.get("flags", 0),
                                                 span.get("size", None))

        pdf_doc.close()
        doc.save(docx_path)
        print(f"[PDF→DOCX] solid fitz fallback done: {Path(docx_path).name}")

    def convert_pdf_to_docx_with_charts(self, pdf_path, docx_path):
        """
        PDF → DOCX, three-tier strategy:
          1. Microsoft Word COM  (best fidelity — native conversion)
          2. pdf2docx            (good layout + chart injection)
          3. Solid fitz fallback (tables, images at position, bold/italic/headings)
        """
        import io as _io
        import fitz

        if self._try_word_com_pdf_to_docx(pdf_path, docx_path):
            return

        base_ok = self._try_pdf2docx(pdf_path, docx_path)

        if base_ok:
            try:
                pdf_doc = fitz.open(pdf_path)
                docx    = Document(docx_path)

                charts_per_page = {}
                for page_num in range(len(pdf_doc)):
                    page  = pdf_doc.load_page(page_num)
                    rects = self._get_chart_regions(page)
                    if not rects:
                        continue
                    imgs = []
                    for rect in rects:
                        png_bytes = self._rasterize_region(page, rect, dpi=150)
                        imgs.append(_io.BytesIO(png_bytes))
                    charts_per_page[page_num] = imgs

                pdf_doc.close()

                if charts_per_page:
                    from docx.shared import Inches as _Inches
                    import copy

                    def _add_picture_paragraph(docx_doc, img_buf, max_width_inches=5.5):
                        p   = docx_doc.add_paragraph()
                        p.alignment = 1
                        run = p.add_run()
                        run.add_picture(img_buf, width=_Inches(max_width_inches))
                        return p

                    def _insert_para_after(ref_para, new_para_xml):
                        ref_para._element.addnext(new_para_xml)

                    body_paras = docx.paragraphs
                    page_break_indices = []
                    for idx, para in enumerate(body_paras):
                        for run in para.runs:
                            if (run._element.xml.find("w:lastRenderedPageBreak") != -1 or
                                    run._element.xml.find("w:br") != -1):
                                page_break_indices.append(idx)
                                break

                    for page_num in sorted(charts_per_page.keys()):
                        imgs = charts_per_page[page_num]
                        if page_num == 0 or page_num - 1 >= len(page_break_indices):
                            for img_buf in imgs:
                                _add_picture_paragraph(docx, img_buf)
                        else:
                            ref_idx  = page_break_indices[page_num - 1]
                            ref_para = body_paras[ref_idx]
                            for img_buf in reversed(imgs):
                                tmp_doc  = Document()
                                pic_para = _add_picture_paragraph(tmp_doc, img_buf)
                                _insert_para_after(ref_para, copy.deepcopy(pic_para._element))

                    docx.save(docx_path)
            except Exception as e:
                print(f"[chart injection] {e}")
            return

        print(f"[PDF→DOCX] falling back to solid fitz for {Path(pdf_path).name}")
        self._solid_fitz_fallback(pdf_path, docx_path)

    def convert_pdf_to_docx_improved(self, pdf_path, docx_path):
        """PDF → DOCX — 3-tier: Word COM → pdf2docx → solid fitz fallback."""
        self.convert_pdf_to_docx_with_charts(pdf_path, docx_path)

    def convert_pdf_to_docx_text_only(self, pdf_path, docx_path):
        """Text-only extraction, no images."""
        try:
            import fitz
            doc = Document()
            pdf_document = fitz.open(pdf_path)
            for page_num in range(len(pdf_document)):
                page = pdf_document.load_page(page_num)
                text = page.get_text("text")
                if text.strip():
                    lines = text.split('\n')
                    for line in lines:
                        line = line.strip()
                        if line:
                            doc.add_paragraph(line)
                if page_num < len(pdf_document) - 1:
                    doc.add_page_break()
            pdf_document.close()
            doc.save(docx_path)
        except Exception as e:
            self._solid_fitz_fallback(pdf_path, docx_path)

    def convert_pdf_to_docx_fallback(self, pdf_path, docx_path):
        """Legacy alias — now delegates to the solid fitz fallback."""
        self._solid_fitz_fallback(pdf_path, docx_path)

    def convert_pdf_to_docx_basic(self, pdf_path, docx_path):
        """Last-resort basic extraction (kept for external callers)."""
        try:
            self._solid_fitz_fallback(pdf_path, docx_path)
        except Exception:
            import fitz
            doc = Document()
            pdf_document = fitz.open(pdf_path)
            doc.add_heading(f"Conversion de: {Path(pdf_path).name}", level=1)
            for page_num in range(len(pdf_document)):
                page = pdf_document.load_page(page_num)
                text = page.get_text()
                doc.add_paragraph(f"--- Page {page_num + 1} ---")
                doc.add_paragraph(text if text.strip() else "[Graphical content]")
                if page_num < len(pdf_document) - 1:
                    doc.add_page_break()
            pdf_document.close()
            doc.save(docx_path)

    def convert_word_to_pdf(self):
        if not (hasattr(self, 'active_templates') and 'word_to_pdf' in self.active_templates):
            _def_id, _ = (self._ensure_template_manager() or object()).get_default_template('Conversion Word→PDF')
            if _def_id:
                (self._ensure_template_manager() or object()).apply_template(_def_id, self)

        selected_items = self.files_list_widget.selectedItems()
        files_to_process = []

        if selected_items:
            for i in range(self.files_list_widget.count()):
                item = self.files_list_widget.item(i)
                if item.isSelected():
                    files_to_process.append(item.data(Qt.UserRole))
        else:
            files_to_process = [f for f in self.files_list if f.lower().endswith(('.docx', '.doc'))]
        
        word_files = [f for f in files_to_process if f.lower().endswith(('.docx', '.doc'))]
        
        if not word_files:
            msg = self.translate_text("Veuillez sélectionner au moins un fichier Word") if selected_items else self.translate_text("La liste doit contenir au moins un fichier Word")
            QMessageBox.warning(self, self.translate_text("Avertissement"), msg)
            return
        
        has_formatted_content = False
        for word_file in word_files[:1]:
            try:
                doc = Document(word_file)
                if len(doc.inline_shapes) > 0:
                    has_formatted_content = True
                    break
                if len(doc.tables) > 0:
                    has_formatted_content = True
                    break
                for paragraph in doc.paragraphs[:10]:
                    if paragraph.style.name not in ['Normal', 'Default Paragraph Font']:
                        has_formatted_content = True
                        break
            except:
                pass
        
        if hasattr(self, 'active_templates') and 'word_to_pdf' in self.active_templates:
            _tpl = self.active_templates['word_to_pdf']
            _quality_map = {
                "Haute (300 DPI)":    "Haute qualité (300 DPI)",
                "Standard (150 DPI)": "Qualité standard (150 DPI)",
                "Basse (72 DPI)":     "Optimisé (96 DPI)",
            }
            options = {
                'mode':             _tpl.get('mode', 'preserve_all'),
                'quality':          _quality_map.get(_tpl.get('quality', ''), 'Qualité standard (150 DPI)'),
                'compress_images':  _tpl.get('compress_images', True),
                'include_metadata': _tpl.get('include_metadata', True),
            }
        else:
            dialog = WordToPdfOptionsDialog(self, self.current_language, has_formatted_content)
            if self.config.get('word_to_pdf_mode') == 'text_only':
                dialog.text_only_radio.setChecked(True)
            if dialog.exec() != QDialog.Accepted:
                return
            options = dialog.get_conversion_mode()
        
        output_dir = self.get_output_directory()
        if not output_dir:
            return
        
        message = self.translate_text("conversion_word_to_pdf").format(len(word_files))
        self.show_progress(True, message)
        self._set_ui_enabled(False)

        _options = options
        _outdir  = output_dir

        def _run_word_to_pdf(task):
            import os, time as _time
            from pathlib import Path as _Path
            t0  = _time.perf_counter()
            fp  = task["input_path"]
            out = task["output_path"]
            fs  = os.path.getsize(fp) if os.path.exists(fp) else 0
            if _options['mode'] == 'preserve_all':
                self.convert_docx_to_pdf_preserve_all(fp, out, _options)
            else:
                self.convert_docx_to_pdf_text_only(fp, out)
            return {"success": True, "error": "",
                    "file_size": fs, "operation_time": _time.perf_counter() - t0}

        tasks = [
            {"index": i, "total": len(word_files),
             "input_path": fp,
             "output_path": os.path.join(_outdir, f"{Path(fp).stem}.pdf")}
            for i, fp in enumerate(word_files)
        ]

        def _on_file_done(result):
            if result.get("success"):
                fp  = result["input_path"]
                out = result["output_path"]
                fs  = result.get("file_size", 0)
                ot  = result.get("operation_time", 0)
                self.achievement_system.record_conversion("word_to_pdf", fs, True)
                self.achievement_system.mark_format_as_used("docx")
                self.achievement_system.mark_format_as_used("pdf")
                if 0 <= datetime.now().hour < 6:
                    self.achievement_system.increment_stat("night_conversions", 1)
                self.db_manager.add_conversion_record(
                    source_file=fp, source_format="DOCX",
                    target_file=out, target_format="PDF",
                    operation_type="word_to_pdf", file_size=fs,
                    conversion_time=ot, success=True,
                    notes=f"Mode: {_options['mode']}"
                )
            else:
                fp = result["input_path"]
                self.db_manager.add_conversion_record(
                    source_file=fp, source_format="DOCX",
                    target_file="", target_format="PDF",
                    operation_type="word_to_pdf", file_size=0,
                    conversion_time=0, success=False,
                    notes=f"Error: {result.get('error','')}"
                )

        def _on_finished(summary):
            self.show_progress(False)
            self._set_ui_enabled(True)
            sc         = summary["success_count"]
            total      = summary["total"]
            total_time = summary["total_time"]
            self.achievement_system.check_achievement("night_owl")
            mode_name = (self.translate_text("Conserver toute la mise en forme")
                         if _options['mode'] == 'preserve_all'
                         else self.translate_text("Texte seulement"))
            formatted_time = f"{total_time:.1f}"
            msg = self.translate_text("word_to_pdf_success_sum").format(
                sc, total, formatted_time, mode_name)
            if self.config.get("enable_system_notifications", True):
                self.system_notifier.send("word_to_pdf")
            if sc >= 50 and total_time <= 300:
                self.achievement_system.update_stat("recent_batch_files", sc)
                self.achievement_system.update_stat("recent_batch_time", total_time)
                self.achievement_system.check_speed_conversion(sc, total_time)
            QMessageBox.information(self, self.translate_text("Succès"),
                                    self.translate_text(msg))

        self._worker = ConversionWorker(tasks, _run_word_to_pdf)
        self._worker.progress.connect(self.progress_bar.setValue)
        self._worker.file_done.connect(_on_file_done)
        self._worker.finished.connect(_on_finished)
        self._worker.start()

    def convert_word_to_pdf_com(self, input_path, output_path, progress_callback=None):
        """Conversion Word → PDF by COM direct (compatible with pyinstaller)"""
        import os
        import pythoncom
        
        input_path = os.path.abspath(input_path)
        output_path = os.path.abspath(output_path)
        
        word = None
        doc = None
        
        if progress_callback:
            progress_callback(40)
            
        try:
            pythoncom.CoInitialize()
            import comtypes.client

            if progress_callback:
                progress_callback(50)

            word = comtypes.client.CreateObject('Word.Application')
            word.Visible = False
            word.DisplayAlerts = 0
            
            if progress_callback:
                progress_callback(60)
            
            doc = word.Documents.Open(input_path)
            
            if progress_callback:
                progress_callback(80)
            
            doc.SaveAs2(output_path, FileFormat=17)
            
            if progress_callback:
                progress_callback(95)
            
            print(f"[SUCCESS] COM conversion successful: {output_path}")
            return True
            
        except Exception as e:
            print(f"[ERROR] COM conversion failed: {e}")
            import traceback
            traceback.print_exc()
            return False
        
        finally:
            if doc is not None:
                try:
                    doc.Close(False)
                except:
                    pass
            if word is not None:
                try:
                    word.Quit()
                except:
                    pass
            try:
                pythoncom.CoUninitialize()
            except:
                pass
            
            if progress_callback:
                progress_callback(100)

    def convert_docx_to_pdf_preserve_all(self, docx_path, pdf_path, options, progress_callback=None):
        """Convert Word to PDF preserving all formatting with progress tracking"""
        
        if progress_callback:
            progress_callback(0)

        try:
            from docx2pdf import convert
            if progress_callback: progress_callback(30)
            convert(docx_path, pdf_path)
            if progress_callback: progress_callback(100)
            print(f"[SUCCESS] Conversion successful with docx2pdf: {pdf_path}")
            return True
        except ImportError:
            print("[INFO] docx2pdf not installed, trying next method...")
        except Exception as e:
            print(f"[WARNING] docx2pdf failed: {e}")

        try:
            if self.convert_word_to_pdf_com(docx_path, pdf_path, progress_callback=progress_callback):
                return True
        
        except ImportError as e:
            print(f"[INFO] COM libraries not available: {e}")
        except Exception as e:
            print(f"[WARNING] COM conversion failed: {e}")
            import traceback
            traceback.print_exc()
        
        print("[INFO] COM/Word not available, trying LibreOffice...")
        try:
            if self._convert_docx_to_pdf_libreoffice(docx_path, pdf_path, progress_callback=progress_callback):
                return True
        except Exception as e:
            print(f"[WARNING] LibreOffice conversion failed: {e}")

        print("[INFO] Using fallback method with reportlab...")
        if progress_callback: progress_callback(10)

        return self._convert_docx_to_pdf_fallback_reportlab(
            docx_path,
            pdf_path,
            options,
            progress_callback=progress_callback
        )

    def _convert_docx_to_pdf_libreoffice(self, docx_path, pdf_path, progress_callback=None):
        """
        Convert DOCX → PDF using LibreOffice headless.
        Tries common LibreOffice executable names on Windows, macOS and Linux.
        Returns True on success, False/raises on failure.
        """
        import os, shutil, subprocess, tempfile
        from pathlib import Path

        candidates = [
            "libreoffice", "soffice",
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
            "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        ]

        soffice = None
        for c in candidates:
            if shutil.which(c) or os.path.isfile(c):
                soffice = c
                break

        if soffice is None:
            print("[INFO] LibreOffice not found on this system.")
            return False

        if progress_callback:
            progress_callback(20)

        with tempfile.TemporaryDirectory() as tmp_dir:
            cmd = [
                soffice,
                "--headless",
                "--convert-to", "pdf",
                "--outdir", tmp_dir,
                os.path.abspath(docx_path),
            ]
            print(f"[INFO] LibreOffice command: {' '.join(cmd)}")

            if progress_callback:
                progress_callback(40)

            result = subprocess.run(
                cmd,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                timeout=120,
            )

            if progress_callback:
                progress_callback(85)

            if result.returncode != 0:
                err = result.stderr.decode(errors="replace")
                print(f"[WARNING] LibreOffice exited with code {result.returncode}: {err}")
                return False

            stem = Path(docx_path).stem
            generated = Path(tmp_dir) / f"{stem}.pdf"
            if not generated.exists():
                matches = list(Path(tmp_dir).glob("*.pdf"))
                if not matches:
                    print("[WARNING] LibreOffice produced no PDF output.")
                    return False
                generated = matches[0]

            shutil.move(str(generated), pdf_path)

        if progress_callback:
            progress_callback(100)

        print(f"[SUCCESS] LibreOffice conversion successful: {pdf_path}")
        return True

    def _convert_docx_to_pdf_fallback_reportlab(self, docx_path, pdf_path, options, progress_callback=None):
        """
        Fallback conversion using reportlab.
        Walks the XML body in document order so images appear at their real
        position and tables are rendered with borders and cell text.
        """
        try:
            from docx import Document
            from docx.oxml.ns import qn
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import A4
            from reportlab.lib import colors
            from PIL import Image as PILImage
            import tempfile, shutil, os, io

            temp_dir = tempfile.mkdtemp()

            doc      = Document(docx_path)
            c        = canvas.Canvas(pdf_path, pagesize=A4)
            W, H     = A4
            margin   = 72
            usable_w = W - 2 * margin
            y        = H - margin
            lh       = 14

            # helpers
            def ensure_space(needed):
                """New page if not enough vertical room."""
                nonlocal y
                if y - needed < margin:
                    c.showPage()
                    y = H - margin

            def wrap_text(text, font_name, font_size, max_w):
                """Word-wrap a string, return list of lines."""
                c.setFont(font_name, font_size)
                words  = text.split()
                lines  = []
                cur    = []
                for w in words:
                    test = ' '.join(cur + [w])
                    if c.stringWidth(test, font_name, font_size) <= max_w:
                        cur.append(w)
                    else:
                        if cur:
                            lines.append(' '.join(cur))
                        cur = [w]
                if cur:
                    lines.append(' '.join(cur))
                return lines or ['']

            def draw_paragraph(para):
                """Render a docx Paragraph object onto the canvas."""
                nonlocal y
                text = para.text.strip()
                if not text:
                    y -= lh // 2
                    return

                style_name = (para.style.name or '') if para.style else ''
                is_heading = any(k in style_name for k in
                                 ['Heading', 'Titre', 'Title', 'heading', 'titre', 'title'])

                if is_heading:
                    font_name, font_size = 'Helvetica-Bold', 14
                    c.setFillColorRGB(0.1, 0.1, 0.4)
                    extra_after = 8
                else:
                    font_name, font_size = 'Helvetica', 11
                    c.setFillColorRGB(0, 0, 0)
                    extra_after = 4

                lines = wrap_text(text, font_name, font_size, usable_w)
                needed = len(lines) * lh + extra_after
                ensure_space(needed)

                c.setFont(font_name, font_size)
                if is_heading:
                    c.setFillColorRGB(0.1, 0.1, 0.4)
                else:
                    c.setFillColorRGB(0, 0, 0)

                for line in lines:
                    ensure_space(lh)
                    c.drawString(margin, y - lh + 2, line)
                    y -= lh
                y -= extra_after

            def draw_image_from_blob(blob):
                """Render an image blob onto the canvas at current y position."""
                nonlocal y
                try:
                    tmp = os.path.join(temp_dir, f"img_{id(blob)}.png")
                    with open(tmp, 'wb') as f:
                        f.write(blob)
                    with PILImage.open(tmp) as img:
                        iw, ih = img.size
                    max_w, max_h = usable_w, 400
                    if iw > max_w:
                        ih = int(ih * max_w / iw); iw = max_w
                    if ih > max_h:
                        iw = int(iw * max_h / ih); ih = max_h
                    ensure_space(ih + 10)
                    x_pos = margin + (usable_w - iw) / 2
                    c.drawImage(tmp, x_pos, y - ih, width=iw, height=ih)
                    y -= ih + 10
                except Exception as e:
                    print(f"[WARN] Image render error: {e}")

            def draw_table(tbl):
                """Render a docx Table onto the canvas."""
                nonlocal y
                rows = tbl.rows
                if not rows:
                    return
                n_cols = max(len(r.cells) for r in rows)
                col_w  = usable_w / n_cols
                row_h  = lh + 8

                c.setFont('Helvetica', 10)
                for r_idx, row in enumerate(rows):
                    ensure_space(row_h)
                    rx = margin
                    for cell in row.cells:
                        # cell border
                        c.setStrokeColorRGB(0.6, 0.6, 0.6)
                        c.setFillColorRGB(1, 1, 1)
                        c.rect(rx, y - row_h, col_w, row_h, stroke=1, fill=1)
                        cell_text = cell.text.strip()
                        if cell_text:
                            c.setFillColorRGB(0, 0, 0)
                            c.setFont('Helvetica-Bold' if r_idx == 0 else 'Helvetica', 9)
                            max_chars = int(col_w / 5.5)
                            display = cell_text[:max_chars] + ('…' if len(cell_text) > max_chars else '')
                            c.drawString(rx + 4, y - row_h + 4, display)
                        rx += col_w
                    y -= row_h
                y -= 6

            # build a mapping rId → blob for inline images
            def _get_inline_image_blob(inline_elem):
                """Extract image blob from an <a:blip> inside an inline/anchor element."""
                try:
                    blip = inline_elem.find('.//' + qn('a:blip'))
                    if blip is None:
                        return None
                    r_embed = blip.get(qn('r:embed'))
                    if r_embed and r_embed in doc.part.rels:
                        return doc.part.rels[r_embed].target_part.blob
                except Exception:
                    pass
                return None

            body_children = list(doc.element.body)
            total = len(body_children)

            for idx, child in enumerate(body_children):
                if progress_callback and total > 0:
                    progress_callback(10 + int(idx / total * 85))

                tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag

                if tag == 'p':
                    inlines = child.findall('.//' + qn('wp:inline')) + \
                              child.findall('.//' + qn('wp:anchor'))
                    if inlines:
                        for inline in inlines:
                            blob = _get_inline_image_blob(inline)
                            if blob:
                                draw_image_from_blob(blob)
                    else:
                        para_obj = None
                        for p in doc.paragraphs:
                            if p._element is child:
                                para_obj = p
                                break
                        if para_obj is not None:
                            draw_paragraph(para_obj)
                        else:
                            raw = ''.join(t.text or '' for t in child.iter(qn('w:t')))
                            if raw.strip():
                                c.setFont('Helvetica', 11)
                                c.setFillColorRGB(0, 0, 0)
                                lines = wrap_text(raw.strip(), 'Helvetica', 11, usable_w)
                                for line in lines:
                                    ensure_space(lh)
                                    c.drawString(margin, y - lh + 2, line)
                                    y -= lh

                elif tag == 'tbl':
                    tbl_obj = None
                    for t in doc.tables:
                        if t._element is child:
                            tbl_obj = t
                            break
                    if tbl_obj is not None:
                        draw_table(tbl_obj)

            try:
                shutil.rmtree(temp_dir)
            except Exception:
                pass

            c.save()
            if progress_callback:
                progress_callback(100)
            print(f"[SUCCESS] Fallback conversion successful: {pdf_path}")
            return True

        except Exception as e:
            print(f"[ERROR] Fallback conversion error: {e}")
            import traceback
            traceback.print_exc()
            if progress_callback:
                progress_callback(100)
            try:
                return self.convert_docx_to_pdf_simple(docx_path, pdf_path)
            except Exception as e2:
                print(f"[ERROR] Simple method failed: {e2}")
                return False

    def convert_docx_to_pdf_text_only(self, docx_path, pdf_path):
        """Convert Word to PDF with text only, well formatted"""
        try:
            from reportlab.lib.enums import TA_JUSTIFY, TA_CENTER, TA_RIGHT
            
            print(f"[DEBUG] Starting text-only conversion: {docx_path}")
            
            doc = Document(docx_path)
            
            buffer = io.BytesIO()
            
            doc_template = SimpleDocTemplate(
                buffer, 
                pagesize=A4,
                rightMargin=72, 
                leftMargin=72,
                topMargin=72, 
                bottomMargin=72
            )
            
            story = []
            styles = getSampleStyleSheet()
            
            style_title = ParagraphStyle(
                'TitleStyle',
                parent=styles['Heading1'],
                fontSize=18,
                leading=22,
                spaceAfter=20,
                alignment=TA_CENTER,
                fontName='Helvetica-Bold',
                textColor='#2c3e50'
            )
            
            style_heading = ParagraphStyle(
                'HeadingStyle',
                parent=styles['Heading2'],
                fontSize=14,
                leading=18,
                spaceAfter=10,
                spaceBefore=15,
                fontName='Helvetica-Bold',
                textColor='#34495e'
            )
            
            style_normal = ParagraphStyle(
                'CleanNormal',
                parent=styles['Normal'],
                fontSize=11,
                leading=14,
                spaceAfter=8,
                alignment=TA_JUSTIFY,
                wordWrap='LTR'
            )
            
            style_bullet = ParagraphStyle(
                'BulletStyle',
                parent=styles['Normal'],
                fontSize=11,
                leading=14,
                spaceAfter=6,
                leftIndent=20,
                firstLineIndent=-15
            )
            
            for i, paragraph in enumerate(doc.paragraphs):
                text = paragraph.text.strip()
                if text:
                    style_name = str(paragraph.style.name) if paragraph.style else "Normal"
                    
                    # Style according to type
                    if i == 0 and len(text) < 100:
                        style = style_title
                    elif 'Heading' in style_name or 'Titre' in style_name:
                        style = style_heading
                    elif text.startswith('•') or text.startswith('-') or text.startswith('*'):
                        style = style_bullet
                        text = "• " + text.lstrip('•-* ')
                    else:
                        style = style_normal
                    
                    if hasattr(paragraph, 'alignment'):
                        if paragraph.alignment == 1:
                            style = ParagraphStyle(
                                'CenterText',
                                parent=style,
                                alignment=TA_CENTER
                            )
                        elif paragraph.alignment == 2:
                            style = ParagraphStyle(
                                'RightText',
                                parent=style,
                                alignment=TA_RIGHT
                            )
                    
                    p = Paragraph(_sanitize_xml(text), style)
                    story.append(p)
                    story.append(Spacer(1, 6))
            
            if story:
                story.append(Spacer(1, 20))
                note_style = ParagraphStyle(
                    'NoteStyle',
                    parent=styles['Normal'],
                    fontSize=9,
                    leading=12,
                    alignment=TA_CENTER,
                    textColor='#7f8c8d',
                    fontName='Helvetica-Oblique'
                )
                note_text = "Note : Version texte seulement - Les images et tableaux ont été omis pour une meilleure lisibilité"
                story.append(Paragraph(note_text, note_style))
            
            if story:
                doc_template.build(story)
            else:
                empty_style = ParagraphStyle(
                    'EmptyStyle',
                    parent=styles['Normal'],
                    fontSize=12,
                    leading=16,
                    alignment=TA_CENTER,
                    textColor='#95a5a6'
                )
                story.append(Paragraph("Document vide", empty_style))
                doc_template.build(story)
            
            with open(pdf_path, 'wb') as f:
                f.write(buffer.getvalue())
            
            buffer.close()
            
            print(f"[DEBUG] Text-only conversion complete: {pdf_path}")
        
        except Exception as e:
            print(f"[ERROR] Text-only conversion error: {e}")
            # Fallback to simple version
            self.convert_docx_to_pdf_simple(docx_path, pdf_path)

    def analyze_word_content(self, docx_path):
        """Analyze the content of a Word document to determine its complexity"""
        try:
            
            doc = Document(docx_path)
            
            content_info = {
                'paragraphs': len(doc.paragraphs),
                'tables': len(doc.tables),
                'images': len(doc.inline_shapes),
                'sections': len(doc.sections),
                'has_formatting': False,
                'complexity': 'simple'
            }
            
            # Check formatting
            for paragraph in doc.paragraphs[:20]:
                if paragraph.style and paragraph.style.name not in ['Normal', 'Default Paragraph Font']:
                    content_info['has_formatting'] = True
                    break
            
            if content_info['images'] > 0 or content_info['tables'] > 2:
                content_info['complexity'] = 'complex'
            elif content_info['has_formatting'] or content_info['paragraphs'] > 50:
                content_info['complexity'] = 'medium'
            else:
                content_info['complexity'] = 'simple'
            
            return content_info
        
        except Exception as e:
            print(f"Word document analysis error: {e}")
            return {
                'paragraphs': 0,
                'tables': 0,
                'images': 0,
                'sections': 0,
                'has_formatting': False,
                'complexity': 'unknown'
            }

    def convert_docx_to_pdf_advanced(self, docx_path, pdf_path):
        """Improved Word to PDF conversion with complete text extraction"""
        try:
            from reportlab.lib.enums import TA_JUSTIFY, TA_LEFT, TA_CENTER, TA_RIGHT
            
            # Load word doc
            doc = Document(docx_path)
            
            # Create PDF
            buffer = io.BytesIO()
            
            doc_template = SimpleDocTemplate(
                buffer, 
                pagesize=A4,
                rightMargin=72, 
                leftMargin=72,
                topMargin=72, 
                bottomMargin=72,
                title=f"Converted from: {Path(docx_path).name}"
            )
            
            story = []
            styles = getSampleStyleSheet()
            
            style_normal = ParagraphStyle(
                'CustomNormal',
                parent=styles['Normal'],
                fontSize=11,
                leading=14,
                spaceAfter=6,
                alignment=TA_JUSTIFY,
                wordWrap='LTR'
            )
            
            style_heading1 = ParagraphStyle(
                'CustomHeading1',
                parent=styles['Heading1'],
                fontSize=16,
                leading=20,
                spaceAfter=12,
                spaceBefore=24,
                alignment=TA_LEFT
            )
            
            style_heading2 = ParagraphStyle(
                'CustomHeading2',
                parent=styles['Heading2'],
                fontSize=14,
                leading=18,
                spaceAfter=8,
                spaceBefore=16,
                alignment=TA_LEFT
            )
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text = paragraph.text
                    
                    style = style_normal
                    if paragraph.style and paragraph.style.name:
                        if 'Heading 1' in paragraph.style.name or 'Titre 1' in paragraph.style.name:
                            style = style_heading1
                        elif 'Heading 2' in paragraph.style.name or 'Titre 2' in paragraph.style.name:
                            style = style_heading2
                        elif 'Title' in paragraph.style.name or 'Titre' in paragraph.style.name:
                            style = style_heading1
                    
                    if paragraph.alignment:
                        if paragraph.alignment == 1:
                            style = ParagraphStyle(
                                'CenterAlign',
                                parent=style,
                                alignment=TA_CENTER
                            )
                        elif paragraph.alignment == 2:
                            style = ParagraphStyle(
                                'RightAlign',
                                parent=style,
                                alignment=TA_RIGHT
                            )
                        elif paragraph.alignment == 3:
                            style = ParagraphStyle(
                                'JustifyAlign',
                                parent=style,
                                alignment=TA_JUSTIFY
                            )
                    
                    p = Paragraph(_sanitize_xml(text), style)
                    story.append(p)
                    story.append(Spacer(1, 6))
            
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            for para in cell.paragraphs:
                                if para.text.strip():
                                    p = Paragraph(para.text, style_normal)
                                    story.append(p)
                            story.append(Spacer(1, 3))
                story.append(Spacer(1, 12))
            
            for section in doc.sections:
                # Header
                if section.header:
                    for para in section.header.paragraphs:
                        if para.text.strip():
                            header_style = ParagraphStyle(
                                'Header',
                                parent=style_normal,
                                fontSize=9,
                                alignment=TA_CENTER
                            )
                            p = Paragraph(para.text, header_style)
                            story.append(p)
                
                if section.footer:
                    for para in section.footer.paragraphs:
                        if para.text.strip():
                            footer_style = ParagraphStyle(
                                'Footer',
                                parent=style_normal,
                                fontSize=9,
                                alignment=TA_CENTER
                            )
                            p = Paragraph(para.text, footer_style)
                            story.append(p)
            
            if story:
                doc_template.build(story)
            else:
                doc_template.build([Paragraph("Document vide ou conversion non supportée", style_normal)])
            
            with open(pdf_path, 'wb') as f:
                f.write(buffer.getvalue())
            
            buffer.close()
            
            if os.path.getsize(pdf_path) < 1024:
                self.convert_docx_to_pdf_fallback(docx_path, pdf_path)
                
        except Exception as e:
            print(f"Advanced Word to PDF conversion error: {e}")
            try:
                self.convert_docx_to_pdf_simple(docx_path, pdf_path)
            except Exception as e2:
                print(f"Simple method failed: {e2}")
                self.create_empty_pdf_with_message(pdf_path, f"Conversion error: {str(e)}")

    def convert_docx_to_pdf_simple(self, docx_path, pdf_path):
        """Simple Word to PDF conversion method"""
        try:
            
            # Load word doc
            doc = Document(docx_path)
            
            # Create PDF
            c = canvas.Canvas(pdf_path, pagesize=A4)
            width, height = A4
            
            margin_left = 72
            margin_right = 72
            margin_top = 72
            margin_bottom = 72
            line_height = 14
            y_position = height - margin_top
            
            # Usable width
            usable_width = width - margin_left - margin_right
            
            # Font and size
            c.setFont("Helvetica", 11)
            
            # Process each paragraph
            for paragraph in doc.paragraphs:
                text = paragraph.text.strip()
                if text:
                    is_heading = False
                    if paragraph.style and paragraph.style.name:
                        heading_keywords = ['Heading', 'Titre', 'Title', 'heading', 'titre', 'title']
                        is_heading = any(keyword in paragraph.style.name for keyword in heading_keywords)
                    
                    if is_heading:
                        # Style for headings
                        c.setFont("Helvetica-Bold", 14)
                        y_position -= 20  # Space before heading
                        
                        if y_position < margin_bottom:
                            c.showPage()
                            y_position = height - margin_top
                            c.setFont("Helvetica-Bold", 14)
                    else:
                        c.setFont("Helvetica", 11)
                        y_position -= line_height
                    
                    if y_position < margin_bottom:
                        c.showPage()
                        y_position = height - margin_top
                        if is_heading:
                            c.setFont("Helvetica-Bold", 14)
                        else:
                            c.setFont("Helvetica", 11)
                    
                    words = text.split()
                    lines = []
                    current_line = []
                    
                    for word in words:
                        test_line = ' '.join(current_line + [word])
                        estimated_width = len(test_line) * 6.5
                        
                        if estimated_width <= usable_width:
                            current_line.append(word)
                        else:
                            if current_line:
                                lines.append(' '.join(current_line))
                            current_line = [word]
                    
                    if current_line:
                        lines.append(' '.join(current_line))
                    
                    for line in lines:
                        if y_position < margin_bottom:
                            c.showPage()
                            y_position = height - margin_top
                            if is_heading:
                                c.setFont("Helvetica-Bold", 14)
                            else:
                                c.setFont("Helvetica", 11)
                        
                        c.drawString(margin_left, y_position, line)
                        y_position -= line_height
                    
                    if is_heading:
                        y_position -= 8
            
            c.save()
        
        except Exception as e:
            print(f"Simple Word to PDF conversion error: {e}")
            raise

    def extract_table_from_element(self, element, doc):
        """Extract table data from XML element"""
        try:
            from docx.table import Table as DocxTable
            
            table_data = []
            
            for row in element.xpath('.//w:tr'):
                row_data = []
                for cell in row.xpath('.//w:tc'):
                    cell_text = ""
                    for para in cell.xpath('.//w:p'):
                        for text_elem in para.xpath('.//w:t'):
                            cell_text += text_elem.text if text_elem.text else ""
                    row_data.append(cell_text.strip())
                
                if row_data:
                    table_data.append(row_data)
            
            return table_data if table_data else None
        
        except Exception as e:
            print(f"Table extraction error: {e}")
            return None

    def convert_docx_to_pdf_with_images(self, docx_path, pdf_path):
        try:
            self.cleanup_temp_files()
            
            doc = Document(docx_path)
            c = canvas.Canvas(pdf_path, pagesize=letter)
            
            width, height = letter
            y_position = height - 50
            line_height = 14
            margin = 50
            
            for paragraph in doc.paragraphs:
                text = paragraph.text.strip()
                if text:
                    if y_position < margin:
                        c.showPage()
                        y_position = height - 50
                    
                    words = text.split()
                    lines = []
                    current_line = []
                    
                    for word in words:
                        test_line = ' '.join(current_line + [word])
                        if len(test_line) * 7 <= (width - 2 * margin):
                            current_line.append(word)
                        else:
                            if current_line:
                                lines.append(' '.join(current_line))
                            current_line = [word]
                    
                    if current_line:
                        lines.append(' '.join(current_line))
                    
                    for line in lines:
                        if y_position < margin:
                            c.showPage()
                            y_position = height - 50
                        
                        c.drawString(margin, y_position, line)
                        y_position -= line_height
            
            try:
                for i, rel in enumerate(doc.part.rels.values()):
                    if "image" in rel.target_ref:
                        try:
                            image_data = rel.target_part.blob
                            
                            temp_img_path = self.create_temp_file(suffix=f"_img_{i}.png")
                            with open(temp_img_path, 'wb') as f:
                                f.write(image_data)
                            
                            if y_position < 200:
                                c.showPage()
                                y_position = height - 50
                            
                            with Image.open(temp_img_path) as img:
                                img_width, img_height = img.size
                            
                            max_width = width - 2 * margin
                            max_height = 300
                            
                            if img_width > max_width:
                                ratio = max_width / img_width
                                img_width = max_width
                                img_height = int(img_height * ratio)
                            
                            if img_height > max_height:
                                ratio = max_height / img_height
                                img_height = max_height
                                img_width = int(img_width * ratio)
                            
                            x_pos = (width - img_width) / 2
                            
                            c.drawImage(temp_img_path, x_pos, y_position - img_height, 
                                        width=img_width, height=img_height)
                            
                            y_position -= img_height + 20
                        
                        except Exception as img_error:
                            print(f"Image error {i}: {img_error}")
                            continue
            
            except Exception as img_error:
                print(f"General image error: {img_error}")
            
            c.save()
        
        except Exception as e:
            try:
                self.convert_docx_to_pdf_fallback(docx_path, pdf_path)
            except Exception as e2:
                self.create_empty_pdf_with_message(pdf_path, f"Conversion error: {str(e)}")

    def convert_docx_to_pdf_fallback(self, docx_path, pdf_path):
        """Fallback method with better text handling"""
        try:
            
            doc = Document(docx_path)
            c = canvas.Canvas(pdf_path, pagesize=A4)
            
            width, height = A4
            y_position = height - 72
            line_height = 14
            margin = 72
            
            max_width = width - 2 * margin
            
            for paragraph in doc.paragraphs:
                text = paragraph.text.strip()
                if text:
                    lines = simpleSplit(text, "Helvetica", 11, max_width)
                    
                    needed_space = len(lines) * line_height
                    if y_position - needed_space < margin:
                        c.showPage()
                        y_position = height - margin
                    
                    for line in lines:
                        c.drawString(margin, y_position, line)
                        y_position -= line_height
                    
                    y_position -= 6
            
            c.save()
        
        except Exception as e:
            print(f"Fallback Word to PDF conversion error: {e}")
            self.create_minimal_pdf_from_docx(docx_path, pdf_path)

    def create_minimal_pdf_from_docx(self, docx_path, pdf_path):
        """Create a minimal PDF with all text from the Word document"""
        try:
            
            doc = Document(docx_path)
            c = canvas.Canvas(pdf_path, pagesize=A4)
            
            width, height = A4
            y_position = height - 50
            line_height = 12
            margin = 50
            
            all_text = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    all_text.append(paragraph.text)
            
            for text in all_text:
                words = text.split()
                lines = []
                current_line = []
                
                for word in words:
                    test_line = ' '.join(current_line + [word])
                    if len(test_line) * 7 <= (width - 2 * margin):
                        current_line.append(word)
                    else:
                        if current_line:
                            lines.append(' '.join(current_line))
                        current_line = [word]
                
                if current_line:
                    lines.append(' '.join(current_line))
                
                for line in lines:
                    if y_position < margin:
                        c.showPage()
                        y_position = height - 50
                    
                    c.drawString(margin, y_position, line)
                    y_position -= line_height
                
                y_position -= 6
            
            c.save()
        
        except Exception as e:
            print(f"Minimal PDF creation error: {e}")
            self.create_empty_pdf_with_message(pdf_path, f"Word document: {Path(docx_path).name}")

    def create_empty_pdf_with_message(self, pdf_path, message):
        c = canvas.Canvas(pdf_path, pagesize=letter)
        width, height = letter
        
        c.setFont("Helvetica", 12)
        c.drawString(100, height - 150, message)
        c.drawString(100, height - 170, "Le document original n'a pas pu être converti correctement.")
        
        c.save()

    def convert_images_to_pdf(self):
        if not (hasattr(self, 'active_templates') and 'images_to_pdf' in self.active_templates):
            _def_id, _ = (self._ensure_template_manager() or object()).get_default_template('Conversion Images→PDF')
            if _def_id:
                (self._ensure_template_manager() or object()).apply_template(_def_id, self)

        if hasattr(self, 'active_templates') and 'images_to_pdf' in self.active_templates:
            self.config['separate_image_pdfs'] =                 self.active_templates['images_to_pdf'].get('separate', False)

        selected_items = self.files_list_widget.selectedItems()
        files_to_process = []
        if selected_items:
            for i in range(self.files_list_widget.count()):
                item = self.files_list_widget.item(i)
                if item.isSelected():
                    files_to_process.append(item.data(Qt.UserRole))
        else:
            files_to_process = self.files_list

        IMAGE_EXTENSIONS = ('.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.webp')
        image_files = [f for f in files_to_process if f.lower().endswith(IMAGE_EXTENSIONS)]
        num_images = len(image_files)
        if num_images == 0:
            msg = self.translate_text("Aucun fichier image compatible sélectionné ou dans la liste.")
            QMessageBox.warning(self, self.translate_text("Avertissement"), msg)
            return

        for file_path in image_files:
            ext = Path(file_path).suffix.lower().lstrip('.')
            if ext == 'jpeg':
                ext = 'jpg'
            if ext in ['jpg', 'png']:
                self.achievement_system.mark_format_as_used(ext)

        self.achievement_system.mark_format_as_used("pdf")

        separate_mode = self.config.get("separate_image_pdfs", False)
        current_hour = datetime.now().hour
        is_night_time = 0 <= current_hour < 6
        if separate_mode:
            self.convert_images_to_separate_pdfs(image_files, selected_items)
        else:
            self.convert_images_to_merged_pdf(image_files, selected_items)

    def convert_images_to_separate_pdfs(self, image_files, selected_items):
        """Convert each image into a separate PDF"""
        num_images = len(image_files)
        
        output_dir = QFileDialog.getExistingDirectory(
            self,
            self.translate_text("Sélectionner un dossier pour les PDFs")
        )
        if not output_dir:
            return
        
        message = self.translate_text("conversion_images_to_separate_pdfs").format(num_images)
        self.show_progress(True, message)
        
        success_count = 0
        start_time = datetime.now()
        total_size = 0
        
        current_hour = datetime.now().hour
        is_night_time = 0 <= current_hour < 6
        
        import fitz
        for i, file_path in enumerate(image_files):
            try:
                base_name = Path(file_path).stem
                output_file = os.path.join(output_dir, f"{base_name}.pdf")
                
                pdf_document = fitz.open()
                img = fitz.open(file_path)
                rect = img[0].rect
                page = pdf_document.new_page(width=rect.width, height=rect.height)
                page.insert_image(rect, filename=file_path)
                img.close()
                pdf_document.save(output_file)
                pdf_document.close()
                
                file_size = os.path.getsize(file_path)
                total_size += file_size
                
                self.db_manager.add_conversion_record(
                    source_file=file_path,
                    source_format=Path(file_path).suffix.upper().replace('.', ''),
                    target_file=output_file,
                    target_format="PDF",
                    operation_type="image_to_pdf_s",
                    file_size=file_size,
                    conversion_time=(datetime.now() - start_time).total_seconds(),
                    success=True,
                    notes=f"Image {i+1}/{num_images} converted to separate PDF"
                )
                
                self.achievement_system.record_conversion("image_to_pdf", file_size, True)
                self.achievement_system.mark_format_as_used("pdf")
                success_count += 1
                
                progress = int((i + 1) / num_images * 100)
                self.progress_bar.setValue(progress)
                
            except Exception as e:
                self.db_manager.add_conversion_record(
                    source_file=file_path,
                    source_format=Path(file_path).suffix.upper().replace('.', ''),
                    target_file="",
                    target_format="PDF",
                    operation_type="image_to_pdf_s",
                    file_size=0,
                    conversion_time=0,
                    success=False,
                    notes=f"Error: {str(e)}"
                )
                try:
                    self.achievement_system.record_conversion("image_to_pdf", 0, False)
                except:
                    pass
                print(f"Individual conversion error for {file_path}: {e}")
        
        total_time = (datetime.now() - start_time).total_seconds()
        self.show_progress(False)
        
        if is_night_time and success_count > 0:
            self.achievement_system.increment_stat("night_conversions", success_count)
            self.achievement_system.check_achievement("night_owl")
        
        if success_count >= 50 and total_time <= 300:
            print(f"[DEBUG] Flash Gordon attempt: {success_count} images to separate PDFs in {total_time:.3f}s")
            self.achievement_system.update_stat("recent_batch_files", success_count)
            self.achievement_system.update_stat("recent_batch_time", total_time)
            self.achievement_system.check_speed_conversion(success_count, total_time)
        
        message = self.translate_text("images_converted_separate").format(
            success_count=success_count,
            num_images=num_images,
            output_dir=output_dir
        )
        if self.config.get("enable_system_notifications", True):
            self.system_notifier.send("image_to_pdf_s")
        QMessageBox.information(self, self.translate_text("Succès"), message)

    def convert_images_to_merged_pdf(self, image_files, selected_items):
        """Convert images into a single merged PDF"""
        num_images = len(image_files)
        is_merge_operation = num_images >= 2
        _quality  = self.config.get('images_to_pdf_quality', 150)
        _compress = self.config.get('images_to_pdf_compress', True)

        current_hour = datetime.now().hour
        is_night_time = 0 <= current_hour < 6
        
        try:
            self.show_progress(True, self.translate_text(f"Traitement de {num_images} image(s)..."))
            
            if is_merge_operation:
                default_filename = f"fusion_images_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
                default_dir = self.config.get("default_output_folder")
                if default_dir and os.path.exists(default_dir):
                    _start = os.path.join(default_dir, default_filename)
                else:
                    _start = default_filename
                output_file, _ = QFileDialog.getSaveFileName(
                    self, self.translate_text("Save file"), _start, self.translate_text("PDF files (*.pdf)")
                )
                if not output_file:
                    self.show_progress(False)
                    return
                
                start_time = datetime.now()
                total_size = sum(os.path.getsize(f) for f in image_files if os.path.exists(f))
                self.progress_bar.setValue(0)
                
                images = []
                for i, file_path in enumerate(image_files):
                    try:
                        img = Image.open(file_path).convert('RGB')
                        images.append(img)
                        self.progress_bar.setValue(int((i + 1) / num_images * 50))
                    except Exception as e:
                        print(f"Image loading error {file_path}: {e}")
                        self.progress_bar.setValue(int((i + 1) / num_images * 50))
                        continue
                
                if not images:
                    self.show_progress(False)
                    QMessageBox.warning(self, self.translate_text("Erreur"), 
                                    self.translate_text("Aucune image valide à traiter."))
                    return
                
                try:
                    first_image = images[0]
                    if len(images) > 1:
                        first_image.save(
                            output_file,
                            format='PDF',
                            save_all=True,
                            append_images=images[1:],
                            resolution=100.0
                        )
                    else:
                        first_image.save(output_file, format='PDF', resolution=100.0)
                    
                    conversion_time = (datetime.now() - start_time).total_seconds()
                    
                    self.achievement_system.update_stat("recent_batch_files", num_images)
                    self.achievement_system.update_stat("recent_batch_time", conversion_time)
                    
                    for img_file in image_files:
                        img_size = os.path.getsize(img_file) if os.path.exists(img_file) else 0
                        self.achievement_system.record_conversion("image_to_pdf", img_size, True)
                        self.achievement_system.mark_format_as_used("pdf")
                    
                    if is_night_time:
                        self.achievement_system.increment_stat("night_conversions", num_images)
                        self.achievement_system.check_achievement("night_owl")
                    
                    self.db_manager.add_conversion_record(
                        source_file=", ".join([Path(f).name for f in image_files]),
                        source_format="Image",
                        target_file=output_file,
                        target_format="PDF",
                        operation_type="image_to_pdf",
                        file_size=total_size,
                        conversion_time=conversion_time,
                        success=True,
                        notes=f"{num_images} images merged into 1 PDF"
                    )
                    
                    self.show_progress(False)
                    message = self.translate_text("images_merged_success").format(
                        num_images=num_images,
                        conversion_time=conversion_time
                    )
                    if self.config.get("enable_system_notifications", True):
                        self.system_notifier.send("image_to_pdf")
                    QMessageBox.information(self, self.translate_text("Succès"), message)
                    
                    if num_images >= 50 and conversion_time <= 300:
                        print(f"[DEBUG] Flash Gordon attempt: {num_images} images merged in {conversion_time:.3f}s")
                        self.achievement_system.check_speed_conversion(num_images, conversion_time)
                        
                except Exception as e:
                    self.show_progress(False)
                    raise e
                    
            elif num_images == 1:
                file_path = image_files[0]
                default_filename = f"{Path(file_path).stem}.pdf"
                output_file = self.get_output_directory(default_filename)
                if not output_file:
                    self.show_progress(False)
                    return
                
                start_time = datetime.now()
                total_size = os.path.getsize(file_path) if os.path.exists(file_path) else 0
                self.progress_bar.setValue(0)
                
                try:
                    img = Image.open(file_path).convert('RGB')
                    img.save(output_file, format='PDF')
                    conversion_time = (datetime.now() - start_time).total_seconds()
                    
                    self.achievement_system.update_stat("recent_batch_files", 1)
                    self.achievement_system.update_stat("recent_batch_time", conversion_time)
                    self.achievement_system.record_conversion("image_to_pdf", total_size, True)
                    self.achievement_system.mark_format_as_used("pdf")
                    
                    if is_night_time:
                        self.achievement_system.increment_stat("night_conversions", 1)
                        self.achievement_system.check_achievement("night_owl")
                    
                    self.db_manager.add_conversion_record(
                        source_file=Path(file_path).name,
                        source_format="Image",
                        target_file=output_file,
                        target_format="PDF",
                        operation_type="image_to_pdf",
                        file_size=total_size,
                        conversion_time=conversion_time,
                        success=True,
                        notes="1 image converted to 1 PDF"
                    )
                    
                    self.show_progress(False)
                    QMessageBox.information(
                        self,
                        self.translate_text("Succès"),
                        self.translate_text("image_to_pdf_success").format(time=conversion_time)
                    )
                    if self.config.get("enable_system_notifications", True):
                        self.system_notifier.send("image_to_pdf")
                except Exception as e:
                    self.show_progress(False)
                    raise e
        
        except Exception as e:
            self.show_progress(False)
            self.db_manager.add_conversion_record(
                source_file=", ".join([Path(f).name for f in image_files]) if image_files else "Unknown",
                source_format="Image",
                target_file="",
                target_format="PDF",
                operation_type="image_to_pdf",
                file_size=0,
                conversion_time=0,
                success=False,
                notes=f"Error: {str(e)}"
            )
            QMessageBox.critical(
                self,
                self.translate_text("Erreur"),
                self.translate_text("error_conversion_fusion").format(error=str(e))
            )

    def merge_pdfs(self):
        from app.ui import MergeOrderDialog
        if not (hasattr(self, 'active_templates') and 'pdf_merge' in self.active_templates):
            _def_id, _ = (self._ensure_template_manager() or object()).get_default_template('Fusion PDF')
            if _def_id:
                (self._ensure_template_manager() or object()).apply_template(_def_id, self)

        selected_items = self.files_list_widget.selectedItems()
        files_to_process = []
        
        if selected_items:
            for i in range(self.files_list_widget.count()):
                item = self.files_list_widget.item(i)
                if item.isSelected():
                    files_to_process.append(item.data(Qt.UserRole))
        else:
            files_to_process = self.files_list
        
        pdf_files = [f for f in files_to_process if f.lower().endswith('.pdf')]
        
        if len(pdf_files) < 2:
            msg = self.translate_text("Veuillez sélectionner au moins 2 fichiers PDF") if selected_items else self.translate_text("La liste doit contenir au moins 2 fichiers PDF")
            QMessageBox.warning(self, self.translate_text("Avertissement"), msg)
            return

        _pre_key = None
        if hasattr(self, 'active_templates') and 'pdf_merge' in self.active_templates:
            _pre_key = self.active_templates['pdf_merge'].get('merge_order_key')

        if _pre_key and _pre_key != 'manual':
            import re as _re
            if _pre_key == 'alpha_az':
                pdf_files.sort(key=lambda f: Path(f).name.lower())
            elif _pre_key == 'alpha_za':
                pdf_files.sort(key=lambda f: Path(f).name.lower(), reverse=True)
            elif _pre_key == 'num_asc':
                def _nk(f):
                    nums = _re.findall(r'\d+', Path(f).stem)
                    return [int(n) for n in nums] if nums else [0]
                pdf_files.sort(key=_nk)
            elif _pre_key == 'num_desc':
                def _nkd(f):
                    nums = _re.findall(r'\d+', Path(f).stem)
                    return [int(n) for n in nums] if nums else [0]
                pdf_files.sort(key=_nkd, reverse=True)
            elif _pre_key == 'date_asc':
                pdf_files.sort(key=lambda f: os.path.getmtime(f))
            elif _pre_key == 'date_desc':
                pdf_files.sort(key=lambda f: os.path.getmtime(f), reverse=True)
        else:
            order_dlg = MergeOrderDialog(pdf_files, "PDF", self, self.current_language,
                                         pre_select_key=_pre_key)
            if order_dlg.exec() != QDialog.Accepted:
                return
            pdf_files = order_dlg.get_ordered_files()

        if hasattr(self, 'active_templates') and 'pdf_merge' in self.active_templates:
            _rn = self.active_templates['pdf_merge'].get('resolved_name', '')
            default_filename = _rn if _rn else f"fusion_pdf_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
        else:
            default_filename = f"fusion_pdf_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
        output_file = self.get_output_directory(default_filename)

        if not output_file:
            return

        try:
            start_time = datetime.now()
            total_size = sum(os.path.getsize(f) for f in pdf_files)

            from PyPDF2 import PdfMerger
            merger = PdfMerger()

            for file_path in pdf_files:
                merger.append(file_path)
            
            output_file_path = output_file
            with open(output_file_path, 'wb') as output_pdf:
                merger.write(output_pdf)
            
            merger.close()
            
            time.sleep(0.2) 
            
            import fitz
            pages_count = 0
            try:
                pdf_document = fitz.open(output_file_path)
                pages_count = len(pdf_document)
                pdf_document.close()
                print(f"[DEBUG MERGE] Pages counted in final PDF: {pages_count}")
            except Exception as e:
                print(f"[WARNING MERGE] Error reading pages from final PDF: {e}")
                pages_count = sum(fitz.open(f).page_count for f in pdf_files)
                print(f"[DEBUG MERGE] Pages counted via sources (fallback): {pages_count}")
            
            self.achievement_system.mark_format_as_used("pdf")
            
            conversion_time = (datetime.now() - start_time).total_seconds()
            
            self.achievement_system.record_pdf_merge(pages_count)
            
            self.db_manager.add_conversion_record(
                source_file=", ".join([Path(f).name for f in pdf_files]),
                source_format="PDF",
                target_file=output_file,
                target_format="PDF",
                operation_type="merge_pdf",
                file_size=total_size,
                conversion_time=conversion_time,
                success=True,
                notes=f"{len(pdf_files)} files merged, Pages: {pages_count}"
            )
            
            message = self.translate_text("pdf_merge_success").format(count=len(pdf_files), time=conversion_time)
            if pages_count > 0:
                message += self.translate_text("pdf_merge_pages_info").format(pages=pages_count)
            if self.config.get("enable_system_notifications", True):
                self.system_notifier.send("merge_pdf")
            QMessageBox.information(self, self.translate_text("Succès"), message)
        
        except Exception as e:
            self.db_manager.add_conversion_record(
                source_file=", ".join([Path(f).name for f in pdf_files]),
                source_format="PDF",
                target_file="",
                target_format="PDF",
                operation_type="merge_pdf",
                file_size=0,
                conversion_time=0,
                success=False,
                notes=f"Error: {str(e)}"
            )
            QMessageBox.critical(
                self,
                self.translate_text("Erreur"),
                self.translate_text("error_merge").format(error=str(e))
            )

    def merge_word_docs(self):
        from app.ui import MergeOrderDialog
        if not (hasattr(self, 'active_templates') and 'word_merge' in self.active_templates):
            _def_id, _ = (self._ensure_template_manager() or object()).get_default_template('Fusion Word')
            if _def_id:
                (self._ensure_template_manager() or object()).apply_template(_def_id, self)

        selected_items = self.files_list_widget.selectedItems()
        files_to_process = []
        
        if selected_items:
            for i in range(self.files_list_widget.count()):
                item = self.files_list_widget.item(i)
                if item.isSelected():
                    files_to_process.append(item.data(Qt.UserRole))
        else:
            files_to_process = self.files_list
        
        word_files = [f for f in files_to_process if f.lower().endswith(('.docx', '.doc'))]
        
        if len(word_files) < 2:
            msg = self.translate_text("Veuillez sélectionner au moins 2 fichiers Word") if selected_items else self.translate_text("La liste doit contenir au moins 2 fichiers Word")
            QMessageBox.warning(self, self.translate_text("Avertissement"), msg)
            return

        _pre_key = None
        if hasattr(self, 'active_templates') and 'word_merge' in self.active_templates:
            _pre_key = self.active_templates['word_merge'].get('merge_order_key')

        if _pre_key and _pre_key != 'manual':
            import re as _re
            if _pre_key == 'alpha_az':
                word_files.sort(key=lambda f: Path(f).name.lower())
            elif _pre_key == 'alpha_za':
                word_files.sort(key=lambda f: Path(f).name.lower(), reverse=True)
            elif _pre_key == 'num_asc':
                def _nk(f):
                    nums = _re.findall(r'\d+', Path(f).stem)
                    return [int(n) for n in nums] if nums else [0]
                word_files.sort(key=_nk)
            elif _pre_key == 'num_desc':
                def _nkd(f):
                    nums = _re.findall(r'\d+', Path(f).stem)
                    return [int(n) for n in nums] if nums else [0]
                word_files.sort(key=_nkd, reverse=True)
            elif _pre_key == 'date_asc':
                word_files.sort(key=lambda f: os.path.getmtime(f))
            elif _pre_key == 'date_desc':
                word_files.sort(key=lambda f: os.path.getmtime(f), reverse=True)
        else:
            order_dlg = MergeOrderDialog(word_files, "Word", self, self.current_language,
                                         pre_select_key=_pre_key)
            if order_dlg.exec() != QDialog.Accepted:
                return
            word_files = order_dlg.get_ordered_files()

        if hasattr(self, 'active_templates') and 'word_merge' in self.active_templates:
            _rn = self.active_templates['word_merge'].get('resolved_name', '')
            default_filename = _rn if _rn else f"fusion_word_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        else:
            default_filename = f"fusion_word_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        default_dir = self.config.get("default_output_folder")
        if default_dir and os.path.exists(default_dir):
            _start = os.path.join(default_dir, default_filename)
        else:
            _start = default_filename
        output_file, _ = QFileDialog.getSaveFileName(
            self, self.translate_text("Save file"), _start, "Word Files (*.docx)"
        )
        if not output_file:
            return

        try:
            start_time = datetime.now()
            total_size = sum(os.path.getsize(f) for f in word_files)

            from docx import Document
            merged_doc = Document(word_files[0])
            
            if len(word_files) > 1:
                merged_doc.add_page_break()
            
            for i, file_path in enumerate(word_files[1:]):
                sub_doc = Document(file_path)
                
                for element in sub_doc.element.body:
                    merged_doc.element.body.append(element)
                
                if i < len(word_files) - 2: 
                    merged_doc.add_page_break()
            
            merged_doc.save(output_file)
            self.achievement_system.mark_format_as_used("docx")
            
            conversion_time = (datetime.now() - start_time).total_seconds()
            
            self.db_manager.add_conversion_record(
                source_file=", ".join([Path(f).name for f in word_files]),
                source_format="DOCX",
                target_file=output_file,
                target_format="DOCX",
                operation_type="merge_word",
                file_size=total_size,
                conversion_time=conversion_time,
                success=True,
                notes=f"{len(word_files)} documents merged"
            )
            
            message = self.translate_text("word_merge_success").format(
                count=len(word_files),
                time=conversion_time
            )
            if self.config.get("enable_system_notifications", True):
                self.system_notifier.send("merge_word")
            QMessageBox.information(self, self.translate_text("Succès"), message)
        
        except Exception as e:
            self.db_manager.add_conversion_record(
                source_file=", ".join([Path(f).name for f in word_files]),
                source_format="DOCX",
                target_file="",
                target_format="DOCX",
                operation_type="merge_word",
                file_size=0,
                conversion_time=0,
                success=False,
                notes=f"Error: {str(e)}"
            )
            QMessageBox.critical(
                self,
                self.translate_text("Erreur"),
                self.translate_text("error_merge").format(error=str(e))
            )

    def split_pdf(self):
        if not (hasattr(self, 'active_templates') and 'pdf_split' in self.active_templates):
            _def_id, _ = (self._ensure_template_manager() or object()).get_default_template('Division PDF')
            if _def_id:
                (self._ensure_template_manager() or object()).apply_template(_def_id, self)

        all_pdf_files = [f for f in self.files_list if f.lower().endswith('.pdf')]
        if not all_pdf_files:
            QMessageBox.warning(
                self,
                self.translate_text("Avertissement"),
                self.translate_text("Veuillez sélectionner au moins un fichier PDF")
            )
            return

        selected_paths = []
        for i in range(self.files_list_widget.count()):
            item = self.files_list_widget.item(i)
            if item.isSelected():
                path = item.data(Qt.UserRole)
                if path and path.lower().endswith('.pdf'):
                    selected_paths.append(path)

        pdf_files = selected_paths if selected_paths else all_pdf_files

        try:
            with open(pdf_files[0], 'rb') as f:
                total_pages_first = len(PdfReader(f).pages)
        except Exception as e:
            QMessageBox.critical(
                self,
                self.translate_text("Erreur"),
                self.translate_text("split_error").format(error=str(e))
            )
            return

        dialog = SplitDialog(total_pages_first, self, self.current_language)

        if len(pdf_files) > 1:
            dialog.setWindowTitle(
                self.translate_text("Diviser PDF") + f"  —  {len(pdf_files)} " + self.translate_text("fichiers")
            )

        _bypass_dialog = False
        if hasattr(self, 'active_templates') and 'pdf_split' in self.active_templates:
            _st = self.active_templates['pdf_split']
            _method_label = _st.get('split_method_label', 'Par pages')
            _ppf = _st.get('pages_per_file', 1)
            for i in range(dialog.split_method.count()):
                if dialog.split_method.itemText(i) == self.translate_text(_method_label):
                    dialog.split_method.setCurrentIndex(i)
                    break
            dialog.page_interval.setValue(_ppf)
            if _method_label != 'Plage de pages':
                _bypass_dialog = True

        if not (_bypass_dialog or dialog.exec() == QDialog.Accepted):
            return

        output_dir = self.get_output_directory()
        if not output_dir:
            return

        suffix        = self.translate_text("split_folder_suffix")
        success_count = 0
        failed_files  = []

        for pdf_file in pdf_files:
            stem        = Path(pdf_file).stem
            file_outdir = os.path.join(output_dir, f"{stem}{suffix}")
            os.makedirs(file_outdir, exist_ok=True)

            start_time = datetime.now()
            file_size  = os.path.getsize(pdf_file)
            try:
                self.split_pdf_file(pdf_file, file_outdir, dialog, silent=len(pdf_files) > 1)
                self.achievement_system.mark_format_as_used("pdf")
                conversion_time = (datetime.now() - start_time).total_seconds()
                self.db_manager.add_conversion_record(
                    source_file=pdf_file,
                    source_format="PDF",
                    target_file=file_outdir,
                    target_format="PDF",
                    operation_type="split_pdf",
                    file_size=file_size,
                    conversion_time=conversion_time,
                    success=True,
                    notes=f"Method: {dialog.split_method.currentText()} | outdir: {file_outdir}"
                )
                success_count += 1
            except Exception as e:
                failed_files.append((Path(pdf_file).name, str(e)))
                self.db_manager.add_conversion_record(
                    source_file=pdf_file,
                    source_format="PDF",
                    target_file="",
                    target_format="PDF",
                    operation_type="split_pdf",
                    file_size=file_size,
                    conversion_time=0,
                    success=False,
                    notes=f"Error: {str(e)}"
                )

        if len(pdf_files) > 1:
            lines = [
                self.translate_text("split_result_summary").format(
                    success=success_count,
                    total=len(pdf_files),
                    output_dir=output_dir
                )
            ]
            for pdf_file in pdf_files[:success_count]:
                lines.append(f"   • {Path(pdf_file).stem}{suffix}/")
            if failed_files:
                lines.append(
                    self.translate_text("split_result_failed").format(count=len(failed_files))
                )
                for name, err in failed_files:
                    lines.append(f"   • {name} : {err}")
            QMessageBox.information(
                self,
                self.translate_text("Résultat"),
                "\n".join(lines)
            )

        if self.config.get("enable_system_notifications", True):
            self.system_notifier.send("split_pdf")

    def split_pdf_file(self, pdf_path, output_dir, dialog, silent=False):
        """
        Splits a single PDF file according to the dialog settings.
        When silent=True (batch mode), suppresses the per-file success popup
        and notification — the caller (split_pdf) handles those instead.
        """
        try:
            with open(pdf_path, 'rb') as file:
                pdf_reader = PdfReader(file)
                total_pages = len(pdf_reader.pages)

                method = dialog.split_method.currentText()

                print(f"[DEBUG SPLIT PDF] File: {Path(pdf_path).name}, Pages: {total_pages}, Method: {method}")
                self.achievement_system.record_pdf_split(total_pages)

                if method == dialog.translate_text("Toutes les pages"):
                    for page_num in range(total_pages):
                        pdf_writer = PdfWriter()
                        pdf_writer.add_page(pdf_reader.pages[page_num])
                        output_file = os.path.join(output_dir, f"{Path(pdf_path).stem}_page_{page_num + 1}.pdf")
                        with open(output_file, 'wb') as output_pdf:
                            pdf_writer.write(output_pdf)
                    message = self.translate_text("pdf_split_into_parts").format(file_count=total_pages)

                elif method == dialog.translate_text("Par pages"):
                    interval   = dialog.page_interval.value()
                    file_count = 0
                    for i in range(0, total_pages, interval):
                        pdf_writer = PdfWriter()
                        end = min(i + interval, total_pages)
                        for page_num in range(i, end):
                            pdf_writer.add_page(pdf_reader.pages[page_num])
                        file_count += 1
                        output_file = os.path.join(output_dir, f"{Path(pdf_path).stem}_part_{file_count}.pdf")
                        with open(output_file, 'wb') as output_pdf:
                            pdf_writer.write(output_pdf)
                    message = self.translate_text("pdf_split_into_files").format(total_pages=total_pages)

                elif method == dialog.translate_text("Plage de pages"):
                    start = dialog.start_page.value() - 1
                    end   = dialog.end_page.value()
                    pdf_writer = PdfWriter()
                    for page_num in range(start, end):
                        pdf_writer.add_page(pdf_reader.pages[page_num])
                    output_file = os.path.join(output_dir, f"{Path(pdf_path).stem}_pages_{start+1}-{end}.pdf")
                    with open(output_file, 'wb') as output_pdf:
                        pdf_writer.write(output_pdf)
                    extracted_pages = end - start
                    current_max = self.achievement_system.stats.get("max_pdf_split_pages", 0)
                    if extracted_pages > current_max:
                        self.achievement_system.update_stat("max_pdf_split_pages", extracted_pages)
                    message = self.translate_text("pages_extracted_success").format(
                        start_page=start + 1,
                        end_page=end
                    )

                if not silent:
                    if self.config.get("enable_system_notifications", True):
                        self.system_notifier.send("split_pdf")
                    QMessageBox.information(self, self.translate_text("Succès"), self.translate_text(message))

        except Exception as e:
            message = self.translate_text("split_error").format(error=str(e))
            QMessageBox.critical(self, self.translate_text("Erreur"), message)

    def protect_pdf(self):
        from app.ui import PdfProtectionDialog 
        if not (hasattr(self, 'active_templates') and 'pdf_protection' in self.active_templates):
            _def_id, _ = (self._ensure_template_manager() or object()).get_default_template('Protection PDF')
            if _def_id:
                (self._ensure_template_manager() or object()).apply_template(_def_id, self)

        selected_items = self.files_list_widget.selectedItems()
        
        if selected_items:
            files_to_process = []
            for i in range(self.files_list_widget.count()):
                item = self.files_list_widget.item(i)
                if item.isSelected():
                    file_path = item.data(Qt.UserRole)
                    # Filter only PDFs
                    if file_path.lower().endswith('.pdf'):
                        files_to_process.append(file_path)
        else:
            files_to_process = [f for f in self.files_list if f.lower().endswith('.pdf')]
        
        if not files_to_process:
            if selected_items:
                msg = self.translate_text("Aucun fichier PDF sélectionné. Veuillez sélectionner au moins un fichier PDF.")
            else:
                msg = self.translate_text("Aucun fichier PDF dans la liste. Ajoutez des fichiers PDF d'abord.")
            QMessageBox.warning(self, self.translate_text("Avertissement"), msg)
            return
        
        if hasattr(self, 'active_templates') and 'pdf_protection' in self.active_templates:
            tpl = self.active_templates['pdf_protection']
            _mode = tpl.get('mode', 'basic')

            from pypdf.constants import UserAccessPermissions
            _perms = UserAccessPermissions(0)
            if tpl.get('allow_printing', True):
                _perms |= UserAccessPermissions.PRINT
                _perms |= UserAccessPermissions.PRINT_TO_REPRESENTATION
            if tpl.get('allow_copying', True):
                _perms |= UserAccessPermissions.EXTRACT
            if tpl.get('allow_copy_accessibility', True):
                _perms |= UserAccessPermissions.EXTRACT_TEXT_AND_GRAPHICS
            if tpl.get('allow_modifications', False):
                _perms |= UserAccessPermissions.MODIFY
            if tpl.get('allow_annotations', False):
                _perms |= UserAccessPermissions.ADD_OR_MODIFY
            if tpl.get('allow_forms', False):
                _perms |= UserAccessPermissions.FILL_FORM_FIELDS
            if tpl.get('allow_assemble', False):
                _perms |= UserAccessPermissions.ASSEMBLE_DOC

            if _mode == 'advanced':
                dialog = PdfProtectionDialog(self, self.current_language)
                dialog.mode_combo.setCurrentIndex(1)
                dialog._on_mode_changed(1)
                dialog.allow_print_check.setChecked(tpl.get('allow_printing', True))
                dialog.allow_copy_check.setChecked(tpl.get('allow_copying', True))
                dialog.allow_copy_accessibility_check.setChecked(tpl.get('allow_copy_accessibility', True))
                dialog.allow_modify_check.setChecked(tpl.get('allow_modifications', False))
                dialog.allow_annotations_check.setChecked(tpl.get('allow_annotations', False))
                dialog.allow_forms_check.setChecked(tpl.get('allow_forms', False))
                dialog.allow_assemble_check.setChecked(tpl.get('allow_assemble', False))
                if dialog.exec() != QDialog.Accepted:
                    return
                password    = dialog.get_password()
                permissions = dialog.get_permissions()
            else:
                password    = None
                permissions = _perms

        else:
            dialog = PdfProtectionDialog(self, self.current_language)
            if dialog.exec() != QDialog.Accepted:
                return
            password    = dialog.get_password()
            permissions = dialog.get_permissions()

        output_dir = self.get_output_directory()
        if not output_dir:
            return

        success_count = 0
        start_time = datetime.now()

        for pdf_file in files_to_process:
            try:
                if selected_items:
                    output_file = os.path.join(output_dir, f"protected_{Path(pdf_file).name}")
                else:
                    output_file = os.path.join(output_dir, Path(pdf_file).name)

                counter = 1
                base_name = Path(output_file).stem
                extension = Path(output_file).suffix
                while os.path.exists(output_file):
                    output_file = os.path.join(output_dir, f"{base_name}_{counter}{extension}")
                    counter += 1

                file_size = os.path.getsize(pdf_file)
                operation_start = datetime.now()

                pdf_reader = PdfReader(pdf_file)
                pdf_writer = PdfWriter()
                for page in pdf_reader.pages:
                    pdf_writer.add_page(page)
                
                import secrets
                owner_password = secrets.token_hex(24)

                if password:
                    pdf_writer.encrypt(password, owner_password=owner_password,
                                       permissions_flag=permissions)
                else:
                    pdf_writer.encrypt("", owner_password=owner_password,
                                       permissions_flag=permissions)

                with open(output_file, 'wb') as output_pdf:
                    pdf_writer.write(output_pdf)

                operation_time = (datetime.now() - operation_start).total_seconds()
                password_length = len(password) if password else 0
                self.achievement_system.record_pdf_protection(1, password_length)

                self.db_manager.add_conversion_record(
                    source_file=pdf_file,
                    source_format="PDF",
                    target_file=output_file,
                    target_format="PDF",
                    operation_type="protect_pdf",
                    file_size=file_size,
                    conversion_time=operation_time,
                    success=True,
                    notes="Password protected"
                )
                success_count += 1

            except Exception as e:
                self.db_manager.add_conversion_record(
                    source_file=pdf_file,
                    source_format="PDF",
                    target_file="",
                    target_format="PDF",
                    operation_type="protect_pdf",
                    file_size=0,
                    conversion_time=0,
                    success=False,
                    notes=f"Error: {str(e)}"
                )
                QMessageBox.warning(
                    self,
                    self.translate_text("Erreur"),
                    self.translate_text("error_with_file").format(
                        filename=Path(pdf_file).name,
                        error=str(e)
                    )
                )

        total_time = (datetime.now() - start_time).total_seconds()

        if selected_items:
            message = self.translate_text("selected_pdfs_protected").format(
                success_count=success_count,
                total_time=total_time
            )
        else:
            message = self.translate_text("all_pdfs_protected").format(
                success_count=success_count,
                total_time=total_time
            )

        QMessageBox.information(self, self.translate_text("Succès"), self.translate_text(message))

        if success_count > 0 and self.config.get("enable_notifications", True):
            self.achievement_system.mark_format_as_used("pdf")
            reply = QMessageBox.question(
                self,
                self.translate_text("Ouverture du dossier"),
                self.translate_text("Voulez-vous ouvrir le dossier contenant les PDF protégés ?"),
                QMessageBox.Yes | QMessageBox.No
            )
            if reply == QMessageBox.Yes:
                try:
                    if sys.platform == "win32":
                        os.startfile(output_dir)
                    elif sys.platform == "darwin":
                        os.system(f"open '{output_dir}'")
                    else:
                        os.system(f"xdg-open '{output_dir}'")
                except Exception as e:
                    print(f"Cannot open folder: {e}")

    def compress_files(self):
        if not (hasattr(self, 'active_templates') and 'compression' in self.active_templates):
            _def_id, _ = (self._ensure_template_manager() or object()).get_default_template('Compression')
            if _def_id:
                (self._ensure_template_manager() or object()).apply_template(_def_id, self)

        selected_items = self.files_list_widget.selectedItems()
        files_to_process = []
        
        if selected_items:
            for i in range(self.files_list_widget.count()):
                item = self.files_list_widget.item(i)
                if item.isSelected():
                    files_to_process.append(item.data(Qt.UserRole))
        else:
            files_to_process = self.files_list
        
        files_to_process = [f for f in files_to_process if f is not None]
        
        if not files_to_process:
            if selected_items:
                msg = self.translate_text("Aucun fichier sélectionné")
            else:
                msg = self.translate_text("La liste de fichiers est vide")
            QMessageBox.warning(self, self.translate_text("Avertissement"), msg)
            return
        
        folders_to_compress = []
        files_to_compress = []
        
        for item_path in files_to_process:
            if item_path is not None:
                if os.path.isdir(item_path):
                    folders_to_compress.append(item_path)
                else:
                    files_to_compress.append(item_path)
        
        compression_mode = "files"
        
        if folders_to_compress:
            dialog = QDialog(self)
            dialog.setWindowTitle(self.translate_text("Mode de compression des dossiers"))
            dialog.setMinimumWidth(500)
            
            layout = QVBoxLayout(dialog)
            
            # Folders information
            folder_info = QLabel(self.translate_text("Dossiers sélectionnés pour compression:"))
            folder_info.setStyleSheet("font-weight: bold;")
            layout.addWidget(folder_info)
            
            folder_list = QTextEdit()
            folder_list.setReadOnly(True)
            folder_list.setMaximumHeight(150)
            
            folder_text = ""
            for folder in folders_to_compress:
                folder_name = Path(folder).name
                file_count = sum(len(files) for _, _, files in os.walk(folder))
                folder_size = self.calculate_folder_size(folder)
                folder_text += self.translate_text("fld_txt").format(folder_name, file_count, self.format_size(folder_size))
            
            folder_list.setText(folder_text)
            layout.addWidget(folder_list)
            
            options_group = QGroupBox(self.translate_text("Options de compression"))
            options_layout = QVBoxLayout(options_group)
            
            option1 = QRadioButton(self.translate_text("📦 Compresser les dossiers avec leur structure (recommandé)"))
            option1.setChecked(True)
            option1.setToolTip(self.translate_text("Crée des archives avec la structure complète des dossiers"))
            
            option2 = QRadioButton(self.translate_text("📄 Traiter les dossiers comme des fichiers individuels"))
            option2.setToolTip(self.translate_text("Ajoute tous les fichiers des dossiers sans conserver la structure"))
            
            options_layout.addWidget(option1)
            options_layout.addWidget(option2)
            layout.addWidget(options_group)
            
            button_box = QDialogButtonBox(QDialogButtonBox.Ok | QDialogButtonBox.Cancel)
            
            def set_mode_and_close():
                nonlocal compression_mode
                if option1.isChecked():
                    compression_mode = "folders_with_structure"
                else:
                    compression_mode = "files_only"
                dialog.accept()
            
            button_box.accepted.connect(set_mode_and_close)
            button_box.rejected.connect(dialog.reject)
            
            layout.addWidget(button_box)
            
            if dialog.exec() != QDialog.Accepted:
                return
        
        if hasattr(self, 'active_templates') and 'compression' in self.active_templates:
            tpl = self.active_templates['compression']
            _fmt_map = {
                'ZIP': self.translate_text('ZIP'), 'RAR': self.translate_text('RAR'),
                'TAR.GZ': self.translate_text('TAR.GZ'), 'TAR': self.translate_text('TAR'),
            }
            _lvl_map = {
                'Normal': self.translate_text('Normal'),
                'Haute compression': self.translate_text('Haute compression'),
                'Compression maximale': self.translate_text('Compression maximale'),
            }
            _fmt = tpl.get('format', 'ZIP')
            _lvl = tpl.get('compression_level', 'Normal')
            _split = tpl.get('split_archive', False)
            _split_size = tpl.get('split_size', 0) if _split else 0
            _encrypt = tpl.get('encrypt', False)
            _delete = tpl.get('delete_originals', False)

            if compression_mode == "folders_with_structure" and folders_to_compress:
                _name = Path(folders_to_compress[0]).name if len(folders_to_compress) == 1 else f"folders_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
            else:
                _name = f"archive_{datetime.now().strftime('%Y%m%d_%H%M%S')}"

            settings = {
                'format':          _fmt_map.get(_fmt, _fmt),
                'level':           _lvl_map.get(_lvl, _lvl),
                'name':            _name,
                'password':        _encrypt,
                'split':           _split,
                'split_size':      _split_size,
                'delete_originals': _delete,
            }

            if _encrypt:
                pwd_dialog = PasswordDialog(self, self.current_language)
                if pwd_dialog.exec() != QDialog.Accepted:
                    return
                password = pwd_dialog.get_password()
                if not password:
                    QMessageBox.warning(self, self.translate_text("Avertissement"),
                                        self.translate_text("Veuillez entrer un mot de passe"))
                    return
                if pwd_dialog.password_input.text() != pwd_dialog.confirm_input.text():
                    QMessageBox.warning(self, self.translate_text("Erreur"),
                                        self.translate_text("Les mots de passe ne correspondent pas"))
                    return

            output_dir = self.get_output_directory()
            if not output_dir:
                return

            archive_format    = settings['format']
            compression_level = settings['level']
            archive_name      = settings['name']
            use_password      = settings['password']
            split_archive     = settings['split']
            split_size        = settings['split_size']
            delete_originals  = settings['delete_originals']
            password          = password if _encrypt else None

        else:
            dialog = CompressionDialog(self, self.current_language)
            dialog.split_checkbox.setEnabled(True)
            if compression_mode == "folders_with_structure" and folders_to_compress:
                _dn = Path(folders_to_compress[0]).name if len(folders_to_compress) == 1                     else f"folders_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
                dialog.filename_input.setText(_dn)

            if dialog.exec() != QDialog.Accepted:
                return

            settings          = dialog.get_compression_settings()
            archive_format    = settings['format']
            compression_level = settings['level']
            archive_name      = settings['name']
            use_password      = settings['password']
            split_archive     = settings['split']
            split_size        = settings['split_size']
            delete_originals  = settings['delete_originals']

            if not archive_name:
                QMessageBox.warning(self, self.translate_text("Erreur"),
                                    self.translate_text("Veuillez entrer un nom pour l'archive"))
                return

            output_dir = self.get_output_directory()
            if not output_dir:
                return

            password = None
            if use_password:
                pwd_dialog = PasswordDialog(self, self.current_language)
                if pwd_dialog.exec() == QDialog.Accepted:
                    password = pwd_dialog.get_password()
                    if not password:
                        QMessageBox.warning(self, self.translate_text("Avertissement"),
                                            self.translate_text("Veuillez entrer un mot de passe"))
                        return
                    if pwd_dialog.password_input.text() != pwd_dialog.confirm_input.text():
                        QMessageBox.warning(self, self.translate_text("Erreur"),
                                            self.translate_text("Les mots de passe ne correspondent pas"))
                        return
                else:
                    return

        archive_format = archive_format.lower()
        norm_map = {
            'tar.gz': 'gz',
            'tar': 'tar',
            'zip': 'zip',
            'rar': 'rar'
        }
        fmt_to_record = norm_map.get(archive_format, archive_format)
        self.achievement_system.mark_format_as_used(fmt_to_record)
        archive_format = settings['format']
        compression_level = settings['level']
        archive_name = settings['name']
        use_password = settings['password']
        split_archive = settings['split']
        split_size = settings['split_size']
        delete_originals = settings['delete_originals']
        
        if not archive_name:
            QMessageBox.warning(self, self.translate_text("Erreur"), 
                            self.translate_text("Veuillez entrer un nom pour l'archive"))
            return
        
        output_dir = self.get_output_directory()
        if not output_dir:
            return
        
        password = None
        if use_password:
            pwd_dialog = PasswordDialog(self, self.current_language)
            if pwd_dialog.exec() == QDialog.Accepted:
                password = pwd_dialog.get_password()
                if not password:
                    QMessageBox.warning(self, self.translate_text("Avertissement"), 
                                    self.translate_text("Veuillez entrer un mot de passe"))
                    return
                if pwd_dialog.password_input.text() != pwd_dialog.confirm_input.text():
                    QMessageBox.warning(self, self.translate_text("Erreur"), 
                                    self.translate_text("Les mots de passe ne correspondent pas"))
                    return
            else:
                return
        
        start_time = datetime.now()
        
        if compression_mode == "folders_with_structure":
            success = self.compress_folders_with_structure(
                folders_to_compress, 
                files_to_compress,
                output_dir, 
                archive_name, 
                archive_format, 
                compression_level, 
                password, 
                delete_originals, 
                split_size
            )
        else:
            all_files = files_to_compress.copy()
            for folder in folders_to_compress:
                for root, dirs, files in os.walk(folder):
                    for file in files:
                        all_files.append(os.path.join(root, file))
            
            success = self.process_compression(
                all_files, 
                output_dir, 
                archive_name, 
                archive_format, 
                compression_level, 
                password, 
                delete_originals, 
                split_size
            )
        
        if success:
            conversion_time = (datetime.now() - start_time).total_seconds()
            
            operation_type = "Compression (selection)" if selected_items else "Compression (all)"
            
            notes = f"Format: {archive_format}, Level: {compression_level}"
            
            if compression_mode == "folders_with_structure":
                notes += self.translate_text("nt_cmp1").format(len(folders_to_compress))
                if files_to_compress:
                    notes += self.translate_text("nt_cmp2").format(len(files_to_compress))
            else:
                notes += self.translate_text("nt_cmp3").format(len(files_to_process))
            total_size = 0
            for item_path in files_to_process:
                if os.path.isfile(item_path):
                    total_size += os.path.getsize(item_path)
                elif os.path.isdir(item_path):
                    total_size += self.calculate_folder_size(item_path) 
            
            self.db_manager.add_conversion_record(
                source_file=", ".join([Path(f).name for f in files_to_process[:3]]),
                source_format="Various",
                target_file=os.path.join(output_dir, f"{archive_name}.{self.get_archive_extension(archive_format)}"),
                target_format=archive_format,
                operation_type=operation_type,
                file_size=total_size,
                conversion_time=conversion_time,
                success=True,
                notes=notes
            )

    def compress_folders_with_structure(self, folders, additional_files, output_dir, archive_name, 
                                        archive_format, compression_level, password, delete_originals, split_size):
        """Compress folders with their complete structure"""
        try:
            print(f"[DEBUG] Compressing folders with structure: {len(folders)} folders, {len(additional_files)} additional files")
            
            total_size = 0
            for folder in folders:
                total_size += self.calculate_folder_size(folder)
            for file in additional_files:
                if os.path.exists(file):
                    total_size += os.path.getsize(file)
            
            total_size_gb = total_size / (1024**3) 
            self.achievement_system.record_compression(total_size_gb)
            
            folder_names = ", ".join([Path(f).name for f in folders])
            message = self.translate_text(f"Compression of {len(folders)} folder(s) with structure: {folder_names}")
            if additional_files:
                message += f" and {len(additional_files)} additional file(s)"
            
            self.show_progress(True, message)
            
            extension = self.get_archive_extension(archive_format)
            archive_path = os.path.join(output_dir, f"{archive_name}.{extension}")
            
            counter = 1
            base_name = Path(archive_path).stem
            while os.path.exists(archive_path):
                archive_path = os.path.join(output_dir, f"{base_name}_{counter}.{extension}")
                counter += 1
            
            print(f"[DEBUG] Final archive: {archive_path}")
            
            if archive_format in ["ZIP", self.translate_text("ZIP")]:
                success = self.create_structured_zip_archive(
                    archive_path, folders, additional_files, compression_level, password, split_size
                )
            elif archive_format in ["RAR", self.translate_text("RAR")]:
                success = self.create_structured_rar_archive(
                    archive_path, folders, additional_files, compression_level, password, split_size
                )
            else:
                all_files = []
                for folder in folders:
                    for root, dirs, files in os.walk(folder):
                        for file in files:
                            all_files.append(os.path.join(root, file))
                
                all_files.extend(additional_files)
                success = self.process_compression(
                    all_files, output_dir, archive_name, archive_format, 
                    compression_level, password, delete_originals, split_size
                )
            
            self.show_progress(False)
            
            if success:
                compressed_size = os.path.getsize(archive_path) if os.path.exists(archive_path) else 0
                
                message = self.translate_text("creat_succ").format(len(folders))
                
                for i, folder in enumerate(folders):
                    folder_name = Path(folder).name
                    file_count = sum(len(files) for _, _, files in os.walk(folder))
                    message += self.translate_text("fl_nc").format(folder_name, file_count)
                
                if additional_files:
                    message += self.translate_text("fl_ad").format(len(additional_files))
                
                message += self.translate_text("fmt_ar").format(archive_format, Path(archive_path).name, self.format_size(compressed_size))
                
                QMessageBox.information(self, self.translate_text("Succès"), self.translate_text(message))
                
                if self.config.get("enable_system_notifications", True):
                    self.system_notifier.send("file_compression")
                
                if delete_originals:
                    deleted_count = 0
                    for item in folders + additional_files:
                        try:
                            if os.path.exists(item):
                                if os.path.isdir(item):
                                    import shutil
                                    shutil.rmtree(item)
                                else:
                                    os.remove(item)
                                deleted_count += 1
                        except Exception as e:
                            print(f"[ERROR] Cannot delete {item}: {e}")
                    
                    if deleted_count > 0:
                        self.status_bar.showMessage(self.translate_text("org_el_del").format(deleted_count))
                
                return True
            else:
                QMessageBox.critical(self, self.translate_text("Erreur"), 
                                self.translate_text("Compression failed"))
                return False
        
        except Exception as e:
            self.show_progress(False)
            print(f"[ERROR] Error compressing folders with structure: {e}")
            import traceback
            traceback.print_exc()
            QMessageBox.critical(self, self.translate_text("Erreur"), 
                            self.translate_text(f"Error during compression: {str(e)}"))
            return False

    def create_structured_zip_archive(self, archive_path, folders, additional_files, compression_level, password, split_size):
        """Create a ZIP archive with folder structure"""
        try:
            print(f"[DEBUG] Creating structured ZIP: {archive_path}")
            
            compression_map = {
                self.translate_text("Normal"): zipfile.ZIP_STORED,
                self.translate_text("Haute compression"): zipfile.ZIP_DEFLATED,
                self.translate_text("Compression maximale"): zipfile.ZIP_LZMA
            }
            
            compression_method = compression_map.get(compression_level, zipfile.ZIP_DEFLATED)
            
            os.makedirs(os.path.dirname(archive_path), exist_ok=True)
            
            if password:
                try:
                    import pyzipper
                    print("[DEBUG] Using pyzipper with AES-256 encryption and structure")
                    
                    with pyzipper.AESZipFile(
                        archive_path, 
                        'w', 
                        compression=compression_method,
                        encryption=pyzipper.WZ_AES
                    ) as zipf:
                        zipf.setpassword(password.encode('utf-8'))
                        
                        for folder in folders:
                            folder_name = Path(folder).name
                            print(f"[DEBUG] Adding folder: {folder_name}")
                            
                            for root, dirs, files in os.walk(folder):
                                for file in files:
                                    full_path = os.path.join(root, file)
                                    rel_path = os.path.relpath(full_path, os.path.dirname(folder))
                                    arcname = os.path.join(folder_name, rel_path)
                                    
                                    try:
                                        zipf.write(full_path, arcname)
                                        print(f"[DEBUG] Added: {arcname}")
                                    except Exception as e:
                                        print(f"[WARNING] Cannot add {full_path}: {e}")
                        
                        for file_path in additional_files:
                            if os.path.exists(file_path):
                                arcname = Path(file_path).name
                                zipf.write(file_path, arcname)
                                print(f"[DEBUG] Additional file added: {arcname}")
                    
                    print(f"[SUCCESS] Structured ZIP archive created: {archive_path}")
                    return True
                
                except ImportError:
                    print("[WARNING] pyzipper not installed, using standard zipfile")
                    password = None
            
            with zipfile.ZipFile(archive_path, 'w', compression=compression_method) as zipf:
                for folder in folders:
                    folder_name = Path(folder).name
                    print(f"[DEBUG] Adding folder: {folder_name}")
                    
                    for root, dirs, files in os.walk(folder):
                        for file in files:
                            full_path = os.path.join(root, file)
                            rel_path = os.path.relpath(full_path, os.path.dirname(folder))
                            arcname = os.path.join(folder_name, rel_path)
                            
                            try:
                                zipf.write(full_path, arcname)
                                print(f"[DEBUG] Added: {arcname}")
                            except Exception as e:
                                print(f"[WARNING] Cannot add {full_path}: {e}")
                
                for file_path in additional_files:
                    if os.path.exists(file_path):
                        arcname = Path(file_path).name
                        zipf.write(file_path, arcname)
                        print(f"[DEBUG] Additional file added: {arcname}")
            
            print(f"[SUCCESS] Structured ZIP archive created (without encryption): {archive_path}")
            return True
        
        except Exception as e:
            print(f"[ERROR] Error creating structured ZIP: {e}")
            import traceback
            traceback.print_exc()
            return False

    def process_compression(self, files_to_compress, output_dir, archive_name, 
                        archive_format, compression_level, password, delete_originals=False, split_size=0):
        """Process file compression with splitting"""
        
        folder_count = sum(1 for f in files_to_compress if os.path.isdir(f))
        
        if folder_count > 0:
            message = self.translate_text(f"Compressing {len(files_to_compress)} items (including {folder_count} folders) to {archive_format}...")
        else:
            message = self.translate_text(f"Compressing {len(files_to_compress)} file(s) to {archive_format}...")
        
        self.show_progress(True, message)
        
        try:
            self.cleanup_temp_files()
            
            extension = self.get_archive_extension(archive_format)
            archive_path = os.path.join(output_dir, f"{archive_name}.{extension}")
            
            counter = 1
            base_name = Path(archive_path).stem
            while os.path.exists(archive_path):
                archive_path = os.path.join(output_dir, f"{base_name}_{counter}.{extension}")
                counter += 1
            
            print(f"[DEBUG PROCESS] Final archive: {archive_path}")
            print(f"[DEBUG PROCESS] Split size: {split_size} MB")
            print(f"[DEBUG PROCESS] Format: {archive_format}")
            
            total_size = sum(os.path.getsize(f) for f in files_to_compress if os.path.exists(f))
            total_size_gb = total_size / (1024**3)
            
            if len(files_to_compress) >= 500:
                print(f"[DEBUG] Attempting Dragon's Breath achievement: {len(files_to_compress)} files")
                self.achievement_system.record_batch_conversion(len(files_to_compress))
            
            success = False
            password_length = len(password) if password else 0
            fmt_lower = archive_format.lower()
            
            if password and password_length >= 12 and fmt_lower in ['zip', 'rar']:
                print(f"[DEBUG] Attempting Impenetrable Fortress achievement: 1 archive, Pwd len: {password_length}")
                self.achievement_system.record_archive_protection(1, password_length, fmt_lower)
            
            if split_size > 0 and archive_format in ["ZIP", self.translate_text("ZIP")]:
                print(f"[DEBUG] Creating split ZIP with WinRAR - max size: {split_size}MB")
                success = self.create_split_zip_archive(archive_path, files_to_compress, compression_level, password, split_size)
            elif split_size > 0 and archive_format in ["RAR", self.translate_text("RAR")]:
                print(f"[DEBUG] Creating split RAR - max size: {split_size}MB")
                success = self.create_rar_archive(archive_path, files_to_compress, compression_level, password, split_size)
            else:
                # case, without splitting (simple ZIP)
                if archive_format in ["RAR", self.translate_text("RAR")]:
                    print(f"[DEBUG] Creating simple RAR (without splitting)")
                    success = self.create_rar_archive(archive_path, files_to_compress, compression_level, password, 0)
                elif archive_format in ["ZIP", self.translate_text("ZIP")]:
                    print(f"[DEBUG] Creating simple ZIP (without splitting)")
                    compression_method = zipfile.ZIP_DEFLATED
                    success = self.create_single_zip_archive(archive_path, files_to_compress, compression_method, password)
                elif archive_format in ["TAR", "TAR.GZ", self.translate_text("TAR"), self.translate_text("TAR.GZ")]:
                    print(f"[DEBUG] Creating simple TAR (without splitting)")
                    success = self.create_tar_archive(archive_path, files_to_compress, archive_format, compression_level)
                else:
                    # Fallback to ZIP by default
                    print(f"[DEBUG] Unknown format, using default ZIP (without splitting)")
                    success = self.create_single_zip_archive(archive_path, files_to_compress, zipfile.ZIP_DEFLATED, password)
            
            self.show_progress(False)
            
            if success:
                parts_created = []
                
                if split_size > 0:
                    base_path = Path(archive_path)
                    base_dir = base_path.parent
                    base_stem = base_path.stem
                    
                    if archive_format in ["ZIP", self.translate_text("ZIP")]:
                        # 1st Format: .zip (main file)
                        if os.path.exists(archive_path):
                            parts_created.append(Path(archive_path))
                        
                        # 2nd Format: .z01, .z02, .z03, etc...
                        pattern_z = f"{base_stem}.z*"
                        z_parts = sorted(base_dir.glob(pattern_z))
                        if z_parts:
                            parts_created.extend(z_parts)
                        
                        # 3rd Format: .zip.001, .zip.002, etc...
                        pattern_zip_num = f"{base_stem}.zip.*"
                        zip_num_parts = sorted(base_dir.glob(pattern_zip_num))
                        if zip_num_parts:
                            parts_created.extend(zip_num_parts)
                        
                        # 4th Format: .part01.zip, .part02.zip, etc...
                        pattern_part = f"{base_stem}.part*.zip"
                        part_parts = sorted(base_dir.glob(pattern_part))
                        if part_parts:
                            parts_created.extend(part_parts)
                    
                    elif archive_format in ["RAR", self.translate_text("RAR")]:
                        # 1st Format: .rar (main file)
                        if os.path.exists(archive_path):
                            parts_created.append(Path(archive_path))
                        
                        # 2nd Format: .r00, .r01, .r02, etc...
                        pattern_r = f"{base_stem}.r*"
                        r_parts = sorted(base_dir.glob(pattern_r))
                        if r_parts:
                            parts_created.extend(r_parts)
                        
                        # 3rd Format: .part01.rar, .part02.rar, etc...
                        pattern_part = f"{base_stem}.part*.rar"
                        part_parts = sorted(base_dir.glob(pattern_part))
                        if part_parts:
                            parts_created.extend(part_parts)
                    
                    # Remove duplicates and sort
                    parts_created = sorted(set(parts_created))
                
                if not parts_created and os.path.exists(archive_path):
                    parts_created = [Path(archive_path)]
                
                compressed_size = 0
                for part in parts_created:
                    if os.path.exists(part):
                        compressed_size += os.path.getsize(part)
                
                self.achievement_system.record_compression(total_size_gb)
                
                fmt_input = archive_format.lower()
                
                if "zip" in fmt_input or fmt_input == self.translate_text("ZIP").lower():
                    self.achievement_system.mark_format_as_used('zip')
                elif "rar" in fmt_input or fmt_input == self.translate_text("RAR").lower():
                    self.achievement_system.mark_format_as_used('rar')
                elif "tar" in fmt_input:
                    if ".gz" in fmt_input or fmt_input == self.translate_text("TAR.GZ").lower():
                        self.achievement_system.mark_format_as_used('gz')
                        self.achievement_system.mark_format_as_used('tar') 
                    else:
                        self.achievement_system.mark_format_as_used('tar')
                
                deleted_count = 0
                if delete_originals:
                    for file_path in files_to_compress:
                        try:
                            if os.path.exists(file_path):
                                os.remove(file_path)
                                deleted_count += 1
                        except Exception as e:
                            print(f"[ERROR] Cannot delete {file_path}: {e}")
                
                num_parts = len(parts_created)
                
                message = (self.translate_text("arch_created").format(archive_format, len(files_to_compress)))
                
                compressed_size_mb = compressed_size / (1024**2)
                
                if split_size > 0 and num_parts > 1:
                    message += self.translate_text("part_created").format(num_parts, split_size)
                    
                    main_file = None
                    for part in parts_created:
                        part_str = str(part)
                        if part_str.endswith('.zip') or part_str.endswith('.rar'):
                            main_file = part
                            break
                    
                    if main_file:
                        message += self.translate_text("main_arch").format(main_file.name)
                    
                    message += self.translate_text("tot_comp").format(f"{compressed_size_mb:.2f}")
                else:
                    message += self.translate_text("arch_success").format(Path(archive_path).name, f"{compressed_size_mb:.2f}")
                
                message += self.translate_text("org_size").format( f"{total_size_gb:.2f}")
                
                if total_size > 0:
                    compression_rate = ((total_size - compressed_size) / total_size * 100)
                    message += self.translate_text("comp_rate").format(f"{compression_rate:.1f}")
                
                if deleted_count > 0:
                    message += self.translate_text("org_del").format(deleted_count)
                
                message += self.translate_text("place_").format(output_dir)
                
                if split_size > 0 and num_parts > 1 and num_parts <= 15:
                    message += self.translate_text("cr_fl")
                    for i, part in enumerate(parts_created[:15]):
                        if os.path.exists(part):
                            size_mb = os.path.getsize(part) / (1024 * 1024)
                            message += self.translate_text("prt").format(part.name, f"{size_mb:.1f}")
                    if num_parts > 15:
                        message += self.translate_text("prt_t").format(num_parts - 15)
                
                QMessageBox.information(self, self.translate_text("Succès"), self.translate_text(message))
                if self.config.get("enable_system_notifications", True):
                    self.system_notifier.send("file_compression")
                return True
            else:
                QMessageBox.critical(self, self.translate_text("Erreur"), 
                                self.translate_text("Compression failed. Check that:\n"
                                                "1. The files exist\n"
                                                "2. You have the permissions\n"
                                                "3. WinRAR is installed (for RAR)"))
                return False
        
        except Exception as e:
            self.show_progress(False)
            print(f"[ERROR] Exception in process_compression: {e}")
            import traceback
            traceback.print_exc()
            QMessageBox.critical(self, self.translate_text("Erreur"), 
                            self.translate_text(f"Error during compression: {str(e)}"))
            return False

    def find_split_archive_parts(self, base_archive_path, archive_format):
        """Find all parts of a split archive"""
        base_path = Path(base_archive_path)
        base_dir = base_path.parent
        base_stem = base_path.stem
        extension = base_path.suffix.lower()
        
        parts_created = []
        
        if archive_format in ["ZIP", self.translate_text("ZIP")]:
            # Patterns for ZIP split
            patterns = [
                f"{base_stem}{extension}",
                f"{base_stem}.z*",
                f"{base_stem}{extension}.*",
                f"{base_stem}.part*{extension}",
            ]
        elif archive_format in ["RAR", self.translate_text("RAR")]:
            patterns = [
                f"{base_stem}{extension}",
                f"{base_stem}.r*",
                f"{base_stem}.part*{extension}",
            ]
        else:
            patterns = [f"{base_stem}{extension}"]
        
        for pattern in patterns:
            try:
                files = list(base_dir.glob(pattern))
                for file in files:
                    if file not in parts_created:
                        parts_created.append(file)
            except:
                continue
        
        parts_created.sort()
        
        return parts_created

    def create_split_zip_archive(self, base_archive_path, files_to_compress, compression_level, password, split_size_mb):
        """Create a split ZIP archive in multiple parts using WinRAR"""
        try:
            print(f"[DEBUG SPLIT ZIP] Starting - max size: {split_size_mb}MB, files: {len(files_to_compress)}, password: {'Yes' if password else 'No'}")
            
            # Search WinRAR in the common locations
            winrar_paths = [
                r"C:\Program Files\WinRAR\WinRAR.exe",
                r"C:\Program Files (x86)\WinRAR\WinRAR.exe",
                r"C:\Program Files\WinRAR\Rar.exe",
                "rar",
                "winrar"
            ]
            
            winrar_exe = None
            for path in winrar_paths:
                if path in ["rar", "winrar"]:
                    try:
                        import subprocess
                        result = subprocess.run([path, "--version"], capture_output=True, shell=True)
                        if result.returncode == 0 or result.returncode == 1:
                            winrar_exe = path
                            break
                    except:
                        continue
                elif os.path.exists(path):
                    winrar_exe = path
                    break
            
            if not winrar_exe:
                QMessageBox.warning(self, self.translate_text("Information"), 
                                self.translate_text("WinRAR not found for ZIP splitting.\n"
                                                    "Installation required for splitting."))
                return False
            
            print(f"[DEBUG] WinRAR found: {winrar_exe}")
            
            import tempfile
            with tempfile.NamedTemporaryFile(mode='w', delete=False, suffix='.txt', encoding='utf-8') as f:
                for file_path in files_to_compress:
                    if os.path.exists(file_path):
                        # Escape paths with spaces
                        escaped_path = file_path.replace('"', '\\"')
                        f.write(f'"{escaped_path}"\n')
                list_file = f.name
            
            print(f"[DEBUG] List file created: {list_file}")
            
            try:
                compression_map = {
                    self.translate_text("Normal"): "-m3",
                    self.translate_text("Haute compression"): "-m5", 
                    self.translate_text("Compression maximale"): "-m5 -md128M"
                }
                
                compression_args = compression_map.get(compression_level, "-m3")
                
                cmd = [winrar_exe, 'a']
                
                cmd.append(compression_args)
                
                cmd.append("-afzip")
                
                cmd.append(f"-v{split_size_mb}M")
                print(f"[DEBUG] Splitting enabled: {split_size_mb}MB per part")
                
                if password:
                    cmd.append(f"-p{password}")
                    cmd.append("-hp")
                    print("[DEBUG] Using password with header encryption")
                else:
                    print("[DEBUG] No password, no encryption options")
                
                cmd.append("-ep1")
                cmd.append("-idq")
                cmd.append("-r")
                
                cmd.append(base_archive_path)
                cmd.append(f"@{list_file}")
                
                print(f"[DEBUG] WinRAR command for split ZIP: {' '.join(cmd)}")
                
                import subprocess
                result = subprocess.run(cmd, capture_output=True, text=True, shell=True)
                
                try:
                    os.unlink(list_file)
                except:
                    pass
                
                if result.returncode == 0:
                    print(f"[DEBUG] Split ZIP archive successfully created: {base_archive_path}")
                    
                    base_path = Path(base_archive_path)
                    base_dir = base_path.parent
                    base_stem = base_path.stem
                    
                    parts_created = []
                    
                    if os.path.exists(base_archive_path):
                        parts_created.append(Path(base_archive_path))
                    
                    pattern_z = f"{base_stem}.z*"
                    z_parts = sorted(base_dir.glob(pattern_z))
                    if z_parts:
                        parts_created.extend(z_parts)
                    
                    pattern_zip_num = f"{base_stem}.zip.*"
                    zip_num_parts = sorted(base_dir.glob(pattern_zip_num))
                    if zip_num_parts:
                        parts_created.extend(zip_num_parts)
                    
                    pattern_part = f"{base_stem}.part*.zip"
                    part_parts = sorted(base_dir.glob(pattern_part))
                    if part_parts:
                        parts_created.extend(part_parts)
                    
                    parts_created = sorted(set(parts_created))
                    
                    if parts_created:
                        print(f"[DEBUG] Split ZIP archive created: {len(parts_created)} parts")
                        for part in sorted(parts_created):
                            size_mb = os.path.getsize(part) / (1024 * 1024)
                            print(f"[DEBUG] Part: {part.name} - {size_mb:.1f}MB")
                        return True
                    else:
                        # Check if at least one file has been created
                        if os.path.exists(base_archive_path):
                            size_mb = os.path.getsize(base_archive_path) / (1024 * 1024)
                            print(f"[DEBUG] Single ZIP archive created: {base_archive_path} - {size_mb:.1f}MB")
                            return True
                        else:
                            print("[ERROR] No archive created")
                            return False
                else:
                    print(f"[ERROR] WinRAR error (code {result.returncode}):")
                    print(f"[ERROR] stdout: {result.stdout}")
                    print(f"[ERROR] stderr: {result.stderr}")
                    
                    try:
                        base_path = Path(base_archive_path)
                        base_dir = base_path.parent
                        base_stem = base_path.stem
                        
                        patterns_to_clean = [
                            f"{base_stem}.zip",
                            f"{base_stem}.z*",
                            f"{base_stem}.zip.*",
                            f"{base_stem}.part*.zip"
                        ]
                        
                        for pattern in patterns_to_clean:
                            for file in base_dir.glob(pattern):
                                try:
                                    os.remove(file)
                                    print(f"[DEBUG] Cleaning: {file.name}")
                                except:
                                    pass
                    except:
                        pass
                    
                    return False
            
            except Exception as e:
                print(f"[ERROR] Exception creating split ZIP with WinRAR: {e}")
                
                try:
                    if os.path.exists(list_file):
                        os.unlink(list_file)
                except:
                    pass
                
                return False
        
        except Exception as e:
            print(f"[ERROR] General error creating split ZIP: {e}")
            import traceback
            traceback.print_exc()
            return False

    def create_single_zip_archive(self, archive_path, files_to_compress, compression_method, password):
        """Create a single ZIP archive with improved password handling"""
        try:
            print(f"[DEBUG CREATE ZIP] Creating: {archive_path}, files: {len(files_to_compress)}, password: {'Yes' if password else 'No'}")
            
            os.makedirs(os.path.dirname(archive_path), exist_ok=True)
            
            if password:
                try:
                    import pyzipper
                    print("[DEBUG] Using pyzipper with AES-256 encryption")
                    
                    with pyzipper.AESZipFile(
                        archive_path, 
                        'w', 
                        compression=compression_method,
                        encryption=pyzipper.WZ_AES
                    ) as zipf:
                        zipf.setpassword(password.encode('utf-8'))
                        
                        for i, file_path in enumerate(files_to_compress):
                            try:
                                if os.path.exists(file_path):
                                    arcname = Path(file_path).name
                                    zipf.write(file_path, arcname)
                                    
                                    progress = int((i + 1) / len(files_to_compress) * 100)
                                    self.progress_bar.setValue(progress)
                                    print(f"[DEBUG] Added to ZIP: {arcname}")
                                else:
                                    print(f"[WARNING] File not found: {file_path}")
                            except Exception as e:
                                print(f"[ERROR] Error adding {file_path}: {e}")
                                return False
                    
                    print(f"[SUCCESS] ZIP archive successfully created: {archive_path}")
                    return True
                
                except ImportError:
                    print("[WARNING] pyzipper not installed, using standard zipfile")
                    QMessageBox.warning(self, self.translate_text("Information"), 
                                    self.translate_text("pyzipper is not installed. Encryption not available."))
                    password = None
                
                except Exception as e:
                    print(f"[ERROR] pyzipper error: {e}")
                    password = None
            
            try:
                with zipfile.ZipFile(archive_path, 'w', compression=compression_method) as zipf:
                    for i, file_path in enumerate(files_to_compress):
                        try:
                            if os.path.exists(file_path):
                                arcname = Path(file_path).name
                                zipf.write(file_path, arcname)
                                
                                progress = int((i + 1) / len(files_to_compress) * 100)
                                self.progress_bar.setValue(progress)
                                print(f"[DEBUG] Added to ZIP: {arcname}")
                            else:
                                print(f"[WARNING] File not found: {file_path}")
                        except Exception as e:
                            print(f"[ERROR] Error adding {file_path}: {e}")
                            return False
                
                print(f"[SUCCESS] ZIP archive created (without encryption): {archive_path}")
                return True
            
            except Exception as e:
                print(f"[ERROR] Error creating ZIP: {e}")
                return False
        
        except Exception as e:
            print(f"[ERROR] Error creating ZIP: {e}")
            import traceback
            traceback.print_exc()
            return False

    def get_archive_extension(self, archive_format):
        extensions = {
            "ZIP": "zip",
            "RAR": "rar", 
            "TAR.GZ": "tar.gz",
            "TAR": "tar",
            self.translate_text("ZIP"): "zip",
            self.translate_text("TAR.GZ"): "tar.gz",
            self.translate_text("TAR"): "tar",
            self.translate_text("RAR"): "rar"
        }
        return extensions.get(archive_format, "zip")

    def create_zip_archive(self, archive_path, files_to_compress, compression_level, password):
        """Create a ZIP archive with encryption support"""
        try:
            compression_map = {
                self.translate_text("Normal"): zipfile.ZIP_STORED,
                self.translate_text("Haute compression"): zipfile.ZIP_DEFLATED,
                self.translate_text("Compression maximale"): zipfile.ZIP_LZMA
            }
            
            compression_method = compression_map.get(compression_level, zipfile.ZIP_DEFLATED)
            
            print(f"[DEBUG] Creating ZIP: {archive_path}")
            print(f"[DEBUG] Compression method: {compression_method}")
            print(f"[DEBUG] Number of files: {len(files_to_compress)}")
            print(f"[DEBUG] Password: {'Yes' if password else 'No'}")
            
            if password:
                try:
                    import pyzipper
                    print("[DEBUG] Using pyzipper with AES encryption")
                    
                    with pyzipper.AESZipFile(
                        archive_path, 
                        'w', 
                        compression=compression_method,
                        encryption=pyzipper.WZ_AES
                    ) as zipf:
                        zipf.setpassword(password.encode('utf-8'))
                        
                        for i, file_path in enumerate(files_to_compress):
                            try:
                                if os.path.exists(file_path):
                                    arcname = Path(file_path).name
                                    zipf.write(file_path, arcname)
                                    
                                    progress = int((i + 1) / len(files_to_compress) * 100)
                                    self.progress_bar.setValue(progress)
                                    print(f"[DEBUG] Added: {arcname}")
                            except Exception as e:
                                print(f"[ERROR] Error adding {file_path}: {e}")
                    
                    print(f"[DEBUG] ZIP archive successfully created: {archive_path}")
                    return True
                
                except ImportError:
                    print("[WARNING] pyzipper not installed, using standard zipfile")
                    QMessageBox.warning(self, self.translate_text("Information"), 
                                        self.translate_text("pyzipper is not installed. Encryption not available."))
                    password = None 
            
            try:
                with zipfile.ZipFile(archive_path, 'w', compression=compression_method) as zipf:
                    for i, file_path in enumerate(files_to_compress):
                        try:
                            if os.path.exists(file_path):
                                arcname = Path(file_path).name
                                zipf.write(file_path, arcname)
                                
                                progress = int((i + 1) / len(files_to_compress) * 100)
                                self.progress_bar.setValue(progress)
                                print(f"[DEBUG] Added: {arcname}")
                        except Exception as e:
                            print(f"[ERROR] Error adding {file_path}: {e}")
                
                print(f"[DEBUG] ZIP archive successfully created: {archive_path}")
                return True
            
            except Exception as e:
                print(f"[ERROR] Error creating ZIP: {e}")
                return False
        
        except Exception as e:
            print(f"[ERROR] Error creating ZIP: {e}")
            return False

    def create_rar_archive(self, archive_path, files_to_compress, compression_level, password, split_size=0):
        """Create a RAR archive with splitting support"""
        try:
            print(f"[DEBUG] Creating RAR: {archive_path}")
            print(f"[DEBUG] Split size: {split_size}MB")
            
            winrar_paths = [
                r"C:\Program Files\WinRAR\WinRAR.exe",
                r"C:\Program Files (x86)\WinRAR\WinRAR.exe",
                r"C:\Program Files\WinRAR\Rar.exe",
                "rar",
                "winrar"
            ]
            
            winrar_exe = None
            for path in winrar_paths:
                if path in ["rar", "winrar"]:
                    try:
                        import subprocess
                        result = subprocess.run([path, "--version"], capture_output=True, shell=True)
                        if result.returncode == 0 or result.returncode == 1:
                            winrar_exe = path
                            break
                    except:
                        continue
                elif os.path.exists(path):
                    winrar_exe = path
                    break
            
            if not winrar_exe:
                QMessageBox.warning(self, self.translate_text("Information"), 
                                self.translate_text("WinRAR not found. Installation required:\n"
                                                    "1. Download WinRAR from win-rar.com\n"
                                                    "2. Install it\n"
                                                    "3. Add WinRAR to PATH or restart the application"))
                return False
            
            print(f"[DEBUG] WinRAR found: {winrar_exe}")
            
            import tempfile
            with tempfile.NamedTemporaryFile(mode='w', delete=False, suffix='.txt', encoding='utf-8') as f:
                for file_path in files_to_compress:
                    if os.path.exists(file_path):
                        escaped_path = file_path.replace('"', '\\"')
                        f.write(f'"{escaped_path}"\n')
                list_file = f.name
            
            print(f"[DEBUG] List file created: {list_file}")
            
            try:
                compression_map = {
                    self.translate_text("Normal"): "-m3",
                    self.translate_text("Haute compression"): "-m5", 
                    self.translate_text("Compression maximale"): "-m5 -md128M"
                }
                
                compression_args = compression_map.get(compression_level, "-m3")
                
                cmd = [winrar_exe, 'a']
                
                cmd.append(compression_args)
                
                if split_size > 0:
                    cmd.append(f"-v{split_size}M")
                    print(f"[DEBUG] Splitting enabled: {split_size}MB per part")
                
                if password:
                    cmd.append(f"-p{password}")
                    cmd.append("-hp")
                    print("[DEBUG] Using password with header encryption")
                else:
                    cmd.append("-p-")
                
                cmd.append("-ep1")
                cmd.append("-idq")
                
                cmd.append(archive_path)
                cmd.append(f"@{list_file}")
                
                print(f"[DEBUG] WinRAR command: {' '.join(cmd)}")
                
                import subprocess
                result = subprocess.run(cmd, capture_output=True, text=True, shell=True)
                
                try:
                    os.unlink(list_file)
                except:
                    pass
                
                if result.returncode == 0:
                    print(f"[DEBUG] RAR archive successfully created: {archive_path}")
                    
                    if split_size > 0:
                        base_name = Path(archive_path).stem
                        base_dir = Path(archive_path).parent
                        parts = list(base_dir.glob(f"{base_name}.part*.rar"))
                        if parts:
                            print(f"[DEBUG] Split archive created: {len(parts)} parts")
                            return True
                        else:
                            if os.path.exists(archive_path):
                                print(f"[DEBUG] Single archive created: {archive_path}")
                                return True
                            else:
                                print("[ERROR] No archive created")
                                return False
                    else:
                        if os.path.exists(archive_path):
                            print(f"[DEBUG] Single archive created: {archive_path}")
                            return True
                        else:
                            print("[ERROR] Archive not created")
                            return False
                else:
                    print(f"[ERROR] WinRAR error (code {result.returncode}):")
                    print(f"[ERROR] stdout: {result.stdout}")
                    print(f"[ERROR] stderr: {result.stderr}")
                    
                    try:
                        if os.path.exists(archive_path):
                            os.remove(archive_path)
                        if split_size > 0:
                            base_name = Path(archive_path).stem
                            base_dir = Path(archive_path).parent
                            for part in base_dir.glob(f"{base_name}.part*.rar"):
                                try:
                                    os.remove(part)
                                except:
                                    pass
                    except:
                        pass
                    
                    return False
            
            except Exception as e:
                print(f"[ERROR] Exception creating RAR: {e}")
                
                try:
                    if os.path.exists(list_file):
                        os.unlink(list_file)
                except:
                    pass
                
                return False
        
        except Exception as e:
            print(f"[ERROR] General error creating RAR: {e}")
            return False

    def create_tar_archive(self, archive_path, files_to_compress, archive_format, compression_level):
        """Create a TAR or TAR.GZ archive"""
        import tarfile
        compression_map = {
            "TAR.GZ": "gz",
            "TAR": None
        }
        
        compression_type = compression_map[archive_format]
        
        mode = "w"
        if compression_type == "gz":
            mode = "w:gz"
        
        try:
            with tarfile.open(archive_path, mode) as tar:
                for i, file_path in enumerate(files_to_compress):
                    try:
                        if os.path.exists(file_path):
                            tar.add(file_path, arcname=Path(file_path).name)
                            
                            progress = int((i + 1) / len(files_to_compress) * 100)
                            self.progress_bar.setValue(progress)
                    except Exception as e:
                        print(f"Error adding {file_path}: {e}")
            
            return True
            
        except Exception as e:
            print(f"Error creating TAR: {e}")
            return False

    def batch_convert(self):
        if not self.files_list:
            QMessageBox.warning(self, self.translate_text("Avertissement"), self.translate_text("Aucun fichier sélectionné"))
            return
        
        dialog = BatchConvertDialog(self, self.current_language)
        if dialog.exec() == QDialog.Accepted:
            target_format = dialog.format_combo.currentText()
            output_dir = self.get_output_directory()
            if output_dir:
                start_time = datetime.now()
                total_size = sum(os.path.getsize(f) for f in self.files_list if os.path.exists(f))
                
                self.process_batch_conversion(output_dir, target_format)
                
                conversion_time = (datetime.now() - start_time).total_seconds()
                
                self.db_manager.add_conversion_record(
                    source_file="Batch of files",
                    source_format="Various",
                    target_file=output_dir,
                    target_format=target_format,
                    operation_type="batch_conversion",
                    file_size=total_size,
                    conversion_time=conversion_time,
                    success=True,
                    notes=f"Target format: {target_format}"
                )

    def process_batch_conversion(self, output_dir, target_format):
        self.show_progress(True, self.translate_text(f"Batch conversion to {target_format}..."))
        
        start_time = datetime.now()
        success_count = 0
        total_files = len(self.files_list)
        failed_files = []
        total_pages_converted = 0
        
        if total_files >= 500:
            print(f"[DEBUG] Attempting Dragon's Breath achievement: {total_files} files")
            self.achievement_system.record_batch_conversion(total_files)
        
        self.achievement_system.record_batch_conversion(total_files)
        
        # Batch conversion  (async)
        self._set_ui_enabled(False)

        _target = target_format
        _outdir = output_dir
        _files  = list(self.files_list)

        def _run_batch_file(task):
            import os, shutil as _shutil
            from pathlib import Path as _Path
            fp       = task["input_path"]
            file_ext = _Path(fp).suffix.lower()
            target   = task["target_format"]
            outdir   = task["output_dir"]
            pages_converted = 0

            if target == self.translate_text("PDF"):
                out = os.path.join(outdir, f"{_Path(fp).stem}.pdf")
                if file_ext in ['.docx', '.doc']:
                    self.convert_docx_to_pdf_advanced(fp, out)
                elif file_ext in ['.jpg', '.jpeg', '.png', '.bmp', '.tiff']:
                    self.convert_single_image_to_pdf(fp, out)
                elif file_ext == '.pdf':
                    _shutil.copy2(fp, out)
                else:
                    return {"success": False,
                            "error": self.translate_text('format not supported for PDF'),
                            "pages_converted": 0}

            elif target == self.translate_text("DOCX"):
                out = os.path.join(outdir, f"{_Path(fp).stem}.docx")
                if file_ext == '.pdf':
                    mode = self.config.get("pdf_to_word_mode", "with_images")
                    if mode == "with_images":
                        self.convert_pdf_to_docx_improved(fp, out)
                    else:
                        self.convert_pdf_to_docx_text_only(fp, out)
                elif file_ext in ['.docx', '.doc']:
                    _shutil.copy2(fp, out)
                else:
                    return {"success": False,
                            "error": self.translate_text('format not supported for DOCX'),
                            "pages_converted": 0}

            elif target == self.translate_text("Images PNG"):
                if file_ext == '.pdf':
                    pdf_folder = os.path.join(outdir, f"{_Path(fp).stem}_pages")
                    os.makedirs(pdf_folder, exist_ok=True)
                    pages_converted = self.convert_pdf_to_png_all_pages(fp, pdf_folder)
                    if pages_converted == 0:
                        return {"success": False,
                                "error": self.translate_text('PDF to PNG conversion failed'),
                                "pages_converted": 0}
                    return {"success": True, "error": "",
                            "pages_converted": pages_converted}
                elif file_ext in ['.jpg', '.jpeg', '.bmp', '.tiff']:
                    out = os.path.join(outdir, f"{_Path(fp).stem}.png")
                    self.convert_to_png(fp, out)
                elif file_ext == '.png':
                    out = os.path.join(outdir, f"{_Path(fp).stem}.png")
                    _shutil.copy2(fp, out)
                else:
                    return {"success": False,
                            "error": self.translate_text('format not supported for PNG'),
                            "pages_converted": 0}
            else:
                return {"success": False, "error": "Unknown target format",
                        "pages_converted": 0}

            return {"success": True, "error": "", "pages_converted": pages_converted}

        tasks = [
            {"index": i, "total": len(_files),
             "input_path": fp,
             "output_path": "",
             "output_dir": _outdir,
             "target_format": _target}
            for i, fp in enumerate(_files)
        ]

        _total_pages = [0]

        def _on_file_done(result):
            _total_pages[0] += result.get("pages_converted", 0)
            if result.get("success"):
                from datetime import datetime as _dt
                if 0 <= _dt.now().hour < 6:
                    self.achievement_system.increment_stat("night_conversions", 1)

        def _on_finished(summary):
            self.show_progress(False)
            self._set_ui_enabled(True)
            sc         = summary["success_count"]
            total      = summary["total"]
            total_time = summary["total_time"]
            failed     = summary["failed"]
            total_pages_converted = _total_pages[0]

            self.achievement_system.update_stat("recent_batch_files", sc)
            self.achievement_system.update_stat("recent_batch_time", total_time)
            if sc >= 50 and total_time <= 300:
                self.achievement_system.check_speed_conversion(sc, total_time)

            failed_names = [f["name"] + ": " + f["error"] for f in failed]

            if _target == self.translate_text("Images PNG") and total_pages_converted > 0:
                message = (f"{sc} {self.translate_text('PDF file(s) converted to')} "
                           f"{total_pages_converted} {self.translate_text('PNG image(s)')}")
                if failed_names:
                    message += f"\n\n{self.translate_text('Failed')} ({len(failed_names)}):\n" + "\n".join(failed_names[:3])
            elif sc == total:
                message = f"{self.translate_text('All')} {sc} {self.translate_text('files converted to')} {_target}"
            elif sc > 0:
                message = f"{sc}/{total} {self.translate_text('files converted to')} {_target}"
                if failed_names:
                    message += f"\n\n{self.translate_text('Failed')} ({len(failed_names)}):\n" + "\n".join(failed_names[:3])
            else:
                message = (f"{self.translate_text('No files converted. Check supported formats.')}\n\n"
                           f"{self.translate_text('Failed')}:\n" + "\n".join(failed_names[:5]))

            QMessageBox.information(self, self.translate_text("Batch Conversion Result"),
                                    self.translate_text(message))
            if self.config.get("enable_system_notifications", True):
                self.system_notifier.send("batch_conversion")

        self._worker = ConversionWorker(tasks, _run_batch_file)
        self._worker.progress.connect(self.progress_bar.setValue)
        self._worker.file_done.connect(_on_file_done)
        self._worker.finished.connect(_on_finished)
        self._worker.start()

    def convert_pdf_to_png_all_pages(self, pdf_path, output_dir):
        try:
            import fitz
            pdf_document = fitz.open(pdf_path)
            total_pages = len(pdf_document)
            pages_converted = 0
            
            for page_num in range(total_pages):
                page = pdf_document.load_page(page_num)
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                
                output_file = os.path.join(output_dir, f"{Path(pdf_path).stem}_page_{page_num + 1}.png")
                pix.save(output_file)
                pages_converted += 1
            
            pdf_document.close()
            return pages_converted
        
        except Exception as e:
            print(f"Error converting PDF to PNG: {e}")
            return 0

    def convert_to_png(self, input_path, output_path):
        try:
            img = Image.open(input_path)
            img.save(output_path, 'PNG')
            return True
        except Exception as e:
            raise Exception(f"{self.translate_text('PNG conversion failed')}: {str(e)}")

    def convert_single_image_to_pdf(self, image_path, pdf_path):
        import fitz
        pdf_document = fitz.open()
        img = fitz.open(image_path)
        rect = img[0].rect
        
        page = pdf_document.new_page(width=rect.width, height=rect.height)
        page.insert_image(rect, filename=image_path)
        
        img.close()
        pdf_document.save(pdf_path)
        pdf_document.close()

    def process_batch_rename(self, rename_plan):
        """Execute rename plan: list of (old_path, new_name) tuples."""
        success_count = 0
        new_files_list = []

        for old_path, new_name in rename_plan:
            try:
                new_path = os.path.join(Path(old_path).parent, new_name)

                counter = 1
                base_stem, ext = os.path.splitext(new_path)
                while os.path.exists(new_path) and new_path != old_path:
                    new_path = f"{base_stem}_{counter}{ext}"
                    counter += 1

                if new_path != old_path:
                    os.rename(old_path, new_path)
                new_files_list.append(new_path)
                success_count += 1

            except Exception as e:
                print(f"Error renaming {old_path}: {e}")
                new_files_list.append(old_path)

        self.files_list = new_files_list

        self.files_list_widget.clear()
        for file_path in self.files_list:
            icon = self.get_file_icon(file_path)
            display_name = Path(file_path).name
            if isinstance(icon, QIcon):
                item = QListWidgetItem(display_name)
                item.setIcon(icon)
            else:
                item = QListWidgetItem(f"{icon} {display_name}")
            item.setData(Qt.UserRole, file_path)
            item.setData(Qt.UserRole + 1, "file")
            if os.path.isfile(file_path):
                item.setData(Qt.UserRole + 4, self.format_size(os.path.getsize(file_path)))
            item.setToolTip(file_path)
            self.files_list_widget.addItem(item)

        self.update_file_counter()
        message = f"{success_count} {self.translate_text('files renamed')}"
        QMessageBox.information(self, self.translate_text("Succès"),
                                self.translate_text(message))
        if self.config.get("enable_system_notifications", True):
            self.system_notifier.send("batch_rename")

    def get_output_directory(self, filename=None):
        """Get output directory with optional filename"""
        default_dir = self.config.get("default_output_folder")
        
        if filename:
            start_dir = os.path.join(default_dir, filename) if (default_dir and os.path.exists(default_dir)) else filename
            ext = Path(filename).suffix.lower()
            if ext == '.pdf':
                file_filter = self.translate_text("PDF files (*.pdf)")
            elif ext in ('.docx', '.doc'):
                file_filter = "Word Files (*.docx)"
            else:
                file_filter = "All files (*.*)"
            return QFileDialog.getSaveFileName(
                self,
                self.translate_text("Save file"),
                start_dir,
                file_filter
            )[0]
        else:
            if default_dir and os.path.exists(default_dir):
                return default_dir
            else:
                return QFileDialog.getExistingDirectory(
                    self, 
                    self.translate_text("Select destination folder")
                )